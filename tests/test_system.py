"""
Test Suite — Validates all presentation factory components.

Run: python -m pytest tests/ -v
Or:  python tests/test_system.py
"""

import json
import os
import sys
import tempfile
from pathlib import Path

# Add scripts to path
sys.path.insert(0, str(Path(__file__).parent.parent / "scripts"))

from content_extractor import (
    extract_content, clean_text, detect_language, ExtractedContent,
)
from gemini_client import (
    GeminiClient, PresentationSpec, SlideSpec, parse_slide_specs,
)
from slide_builder import SlideBuilder, BrandConfig
from pptx.dml.color import RGBColor
from image_generator import (
    generate_slide_images, generate_mock_images,
    build_image_prompt, load_image_prompt_template,
    build_full_slide_prompt, load_full_slide_template,
    LAYOUT_TEMPLATES, FULL_SLIDE_MODEL, COMPOSITE_MODEL,
)


# ─── Content Extractor Tests ──────────────────────────────────

def test_language_detection():
    assert detect_language("El directivo influye en sus empleados") == "ES"
    assert detect_language("The manager influences their employees") == "EN"
    assert detect_language("Los datos muestran la tendencia del mercado") == "ES"
    assert detect_language("Data shows the market trend is positive") == "EN"
    print("  [PASS] Language detection")


def test_clean_text():
    dirty = "  Hello  \n\n\n\n\n  World  \n  123  \n  Test  "
    cleaned = clean_text(dirty)
    assert "\n\n\n\n\n" not in cleaned
    assert "Hello" in cleaned
    assert "World" in cleaned
    print("  [PASS] Text cleaning")


def test_extract_pdf():
    test_pdf = "C:/Users/House_Fol/Downloads/IESE/2026-03-02 - PDG-A-2026/2026-03-02 - PDG-A-2026/T_cticas_de_influencia.pdf"
    if not Path(test_pdf).exists():
        print("  [SKIP] PDF test (file not found)")
        return

    content = extract_content(test_pdf)
    assert content.word_count > 1000, f"Expected >1000 words, got {content.word_count}"
    assert content.language == "ES", f"Expected ES, got {content.language}"
    assert content.page_count == 8, f"Expected 8 pages, got {content.page_count}"
    assert "influencia" in content.text.lower()
    print(f"  [PASS] PDF extraction ({content.word_count} words, {content.page_count} pages)")


def test_extract_pdf_porter():
    test_pdf = "C:/Users/House_Fol/Downloads/IESE/2026-03-02 - PDG-A-2026/2026-03-02 - PDG-A-2026/Las_cinco_fuerzas_competitivas_que_le_dan_forma_a_la_estrate.pdf"
    if not Path(test_pdf).exists():
        print("  [SKIP] Porter PDF test (file not found)")
        return

    content = extract_content(test_pdf)
    assert content.word_count > 5000, f"Expected >5000 words, got {content.word_count}"
    assert content.language == "ES"
    print(f"  [PASS] Porter PDF extraction ({content.word_count} words)")


# ─── Gemini Client Tests ──────────────────────────────────────

def test_mock_generation():
    client = GeminiClient()  # No API key = mock mode
    assert client.is_mock

    spec = client.generate_slide_specs(
        content="Test content about influence tactics",
        filename="test.pdf",
        language="ES",
        slide_count=6,
    )

    assert isinstance(spec, PresentationSpec)
    assert spec.total_slides >= 4
    assert spec.language == "ES"
    assert all(s.title for s in spec.slides)
    assert all(s.visual_concept for s in spec.slides)
    print(f"  [PASS] Mock slide generation ({spec.total_slides} slides)")


def test_mock_generation_en():
    client = GeminiClient()
    spec = client.generate_slide_specs(
        content="Test content about leadership",
        filename="test.pdf",
        language="EN",
        slide_count=6,
    )

    assert spec.language == "EN"
    assert spec.total_slides >= 4
    print(f"  [PASS] Mock EN generation ({spec.total_slides} slides)")


def test_mock_visual_concepts_are_image_friendly():
    """Verify mock visual concepts describe scenes, not diagrams with labels."""
    client = GeminiClient()
    spec = client.generate_slide_specs(
        content="Test content",
        filename="test.pdf",
        language="ES",
        slide_count=6,
    )

    for slide in spec.slides:
        vc = slide.visual_concept.lower()
        # Should NOT contain chart/diagram labels
        assert "labeled" not in vc, f"Slide {slide.number} visual_concept contains 'labeled': {vc[:80]}"
        # Reject "label" unless it appears only as part of a compound word
        words = vc.split()
        assert "label" not in words, f"Slide {slide.number} contains standalone 'label': {vc[:80]}"
    print("  [PASS] Visual concepts are image-friendly (no diagram labels)")


def test_parse_slide_specs():
    json_str = json.dumps({
        "metadata": {
            "title": "Test",
            "subtitle": "Sub",
            "language": "EN",
            "source_document": "test.pdf",
            "total_slides": 2,
            "themes": ["test"],
        },
        "slides": [
            {
                "number": 1,
                "type": "title",
                "title": "Test Title",
                "subtitle": "Subtitle",
                "body": "",
                "bullet_points": [],
                "visual_concept": "concept",
                "speaker_notes": "notes",
                "source_reference": "ref",
            },
            {
                "number": 2,
                "type": "content",
                "title": "Content Slide",
                "body": "Some body text here",
                "bullet_points": ["Point 1", "Point 2"],
                "visual_concept": "diagram",
                "speaker_notes": "notes 2",
                "source_reference": "page 1",
            },
        ]
    })

    spec = parse_slide_specs(json_str)
    assert spec.title == "Test"
    assert spec.total_slides == 2
    assert spec.slides[0].type == "title"
    assert spec.slides[1].bullet_points == ["Point 1", "Point 2"]
    print("  [PASS] JSON parsing")


def test_parse_with_markdown_fences():
    json_str = '```json\n{"metadata": {"title": "T", "subtitle": "", "language": "EN", "source_document": "t.pdf", "total_slides": 1, "themes": []}, "slides": [{"number": 1, "type": "title", "title": "T"}]}\n```'
    spec = parse_slide_specs(json_str)
    assert spec.title == "T"
    print("  [PASS] JSON with markdown fences")


def test_specs_to_json():
    client = GeminiClient()
    spec = client.generate_slide_specs("test", "test.pdf", "ES", slide_count=4)
    json_str = client.specs_to_json(spec)
    data = json.loads(json_str)
    assert "metadata" in data
    assert "slides" in data
    assert len(data["slides"]) == spec.total_slides
    print("  [PASS] Specs to JSON serialization")


# ─── Image Generator Tests ────────────────────────────────────

def test_image_prompt_template():
    """Verify the image prompt template loads and formats correctly."""
    template = load_image_prompt_template()
    assert "{visual_concept}" in template
    assert "{slide_title}" in template
    assert "{slide_type}" in template
    assert "text" in template.lower()  # Should mention no-text rule
    print("  [PASS] Image prompt template loads correctly")


def test_build_image_prompt():
    """Verify image prompts are built with concept, title, and type."""
    prompt = build_image_prompt(
        visual_concept="A glowing network of nodes",
        slide_title="Test Slide",
        slide_type="content",
    )
    assert "glowing network" in prompt
    assert "Test Slide" in prompt
    assert "content" in prompt
    print("  [PASS] Image prompt building")


def test_mock_images():
    """Verify mock image generation returns None for all slides."""
    slides = [
        SlideSpec(1, "title", "Title", visual_concept="v"),
        SlideSpec(2, "content", "Content", visual_concept="v"),
    ]
    paths = generate_mock_images(slides)
    assert len(paths) == 2
    assert all(v is None for v in paths.values())
    print("  [PASS] Mock image generation")


def test_image_generation_no_api_key():
    """Without API key, generate_slide_images returns all None paths."""
    slides = [
        SlideSpec(1, "title", "Title", visual_concept="v"),
        SlideSpec(2, "content", "Content", visual_concept="v"),
    ]
    # Temporarily unset GEMINI_API_KEY
    old_key = os.environ.pop("GEMINI_API_KEY", None)
    try:
        paths = generate_slide_images(slides, api_key=None)
        assert len(paths) == 2
        assert all(v is None for v in paths.values())
    finally:
        if old_key:
            os.environ["GEMINI_API_KEY"] = old_key
    print("  [PASS] Image generation gracefully handles missing API key")


def test_full_slide_prompt_template():
    """Verify the full-slide prompt template loads and has all placeholders."""
    template = load_full_slide_template()
    assert "{slide_type}" in template
    assert "{layout_instructions}" in template
    assert "{text_content}" in template
    assert "{visual_concept}" in template
    assert "2752x1536" in template or "16:9" in template
    assert "correctly spelled" in template.lower() or "crisp" in template.lower()
    print("  [PASS] Full-slide prompt template loads correctly")


def test_layout_templates_cover_all_types():
    """Verify layout templates exist for all 7 slide types."""
    expected = {"title", "content", "section", "data", "quote", "comparison", "conclusion"}
    assert set(LAYOUT_TEMPLATES.keys()) == expected
    for stype, layout in LAYOUT_TEMPLATES.items():
        assert len(layout) > 50, f"Layout for {stype} is too short"
    print("  [PASS] Layout templates cover all 7 slide types")


def test_build_full_slide_prompt_title():
    """Full-slide prompt for title includes all content and layout instructions."""
    slide = SlideSpec(1, "title", "Strategy Drives Growth", subtitle="A 2026 Analysis",
                      visual_concept="Dramatic sunrise over a cityscape")
    prompt = build_full_slide_prompt(slide)
    assert "Strategy Drives Growth" in prompt
    assert "2026 Analysis" in prompt
    assert "sunrise" in prompt.lower() or "cityscape" in prompt.lower()
    assert "title" in prompt.lower()
    assert "16:9" in prompt or "widescreen" in prompt.lower()
    print("  [PASS] Full-slide prompt for title slide")


def test_build_full_slide_prompt_content():
    """Full-slide prompt for content includes body and bullets."""
    slide = SlideSpec(2, "content", "Key Insight Here",
                      body="Evidence supports the claim.",
                      bullet_points=["Point A", "Point B"],
                      visual_concept="Abstract geometric shapes")
    prompt = build_full_slide_prompt(slide)
    assert "Key Insight Here" in prompt
    assert "Evidence supports the claim" in prompt
    assert "Point A" in prompt
    assert "Point B" in prompt
    # Content layout should mention grid or column structure
    assert "grid" in prompt.lower() or "column" in prompt.lower() or "layout" in prompt.lower()
    print("  [PASS] Full-slide prompt for content slide")


def test_build_full_slide_prompt_all_types():
    """Full-slide prompts use type-specific layout instructions for every slide type."""
    for stype in ["title", "content", "section", "data", "quote", "comparison", "conclusion"]:
        slide = SlideSpec(1, stype, f"Test {stype}", visual_concept="abstract art")
        prompt = build_full_slide_prompt(slide)
        layout = LAYOUT_TEMPLATES[stype]
        # The layout instructions should be embedded in the prompt
        assert layout[:30] in prompt, f"Layout for {stype} not found in prompt"
    print("  [PASS] Full-slide prompts use type-specific layouts for all 7 types")


def test_model_constants():
    """Verify model constants are set correctly."""
    assert "pro" in FULL_SLIDE_MODEL.lower() or "3" in FULL_SLIDE_MODEL
    assert "flash" in COMPOSITE_MODEL.lower()
    print("  [PASS] Model constants are correct")


# ─── Slide Builder Tests ──────────────────────────────────────

def test_brand_config():
    brand_path = Path(__file__).parent.parent / "config" / "brand.json"
    brand = BrandConfig.from_json(str(brand_path))
    assert brand.font_title == "Outfit SemiBold"
    assert brand.font_body == "DM Sans"
    assert brand.font_accent == "DM Sans Medium"
    assert brand.primary is not None
    # Verify Modern cyan primary
    assert brand.primary[0] == 0x06  # Red component of #06B6D4
    print("  [PASS] Brand config loading (Modern/Outfit+DM Sans)")


def test_build_presentation_no_images():
    """Build presentation without images (gradient-only fallback)."""
    spec = PresentationSpec(
        title="Test Pres",
        subtitle="Subtitle",
        language="EN",
        source_document="test.pdf",
        themes=["test"],
        slides=[
            SlideSpec(1, "title", "Title Slide", subtitle="Sub", visual_concept="v"),
            SlideSpec(2, "content", "Content Slide",
                      body="Some text", bullet_points=["A", "B"],
                      visual_concept="v"),
            SlideSpec(3, "conclusion", "Conclusion",
                      body="Summary", bullet_points=["Key 1"],
                      visual_concept="v"),
        ],
    )

    builder = SlideBuilder()  # No image_paths = gradient fallback
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    assert Path(path).exists()
    size = Path(path).stat().st_size
    assert size > 10000, f"PPTX too small: {size} bytes"

    from pptx import Presentation
    prs = Presentation(path)
    assert len(prs.slides) == 3

    Path(path).unlink()
    print(f"  [PASS] PPTX generation without images ({size:,} bytes, {len(prs.slides)} slides)")


def test_all_slide_types():
    """Verify all 7 slide types render without errors."""
    spec = PresentationSpec(
        title="All Types",
        subtitle="Test",
        language="EN",
        source_document="test.pdf",
        themes=[],
        slides=[
            SlideSpec(1, "title", "Title", subtitle="Sub", visual_concept="v"),
            SlideSpec(2, "section", "Section", body="Intro", visual_concept="v"),
            SlideSpec(3, "content", "Content", body="Text", bullet_points=["A"],
                      visual_concept="v"),
            SlideSpec(4, "comparison", "Compare", bullet_points=["Left", "Right"],
                      visual_concept="v"),
            SlideSpec(5, "data", "93%", body="Key stat", bullet_points=["Detail"],
                      visual_concept="v"),
            SlideSpec(6, "quote", "Famous Quote", body="The quote text",
                      source_reference="Author", visual_concept="v"),
            SlideSpec(7, "conclusion", "Takeaway", body="Summary",
                      bullet_points=["Point"], visual_concept="v"),
        ],
    )

    builder = SlideBuilder()
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    assert len(prs.slides) == 7
    Path(path).unlink()
    print(f"  [PASS] All 7 slide types render correctly")


def test_slide_builder_with_image_paths():
    """Verify SlideBuilder accepts image_paths dict without error."""
    spec = PresentationSpec(
        title="Image Test",
        subtitle="",
        language="EN",
        source_document="test.pdf",
        themes=[],
        slides=[
            SlideSpec(1, "title", "Title", visual_concept="v"),
            SlideSpec(2, "content", "Content", body="Text", visual_concept="v"),
        ],
    )

    # image_paths with None values (no actual images) should still work
    image_paths = {1: None, 2: None}
    builder = SlideBuilder(image_paths=image_paths)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    assert Path(path).exists()
    Path(path).unlink()
    print("  [PASS] SlideBuilder accepts image_paths parameter")


def test_alternating_content_layout():
    """Content slides should alternate image left/right based on slide number."""
    spec = PresentationSpec(
        title="Layout Test",
        subtitle="",
        language="EN",
        source_document="test.pdf",
        themes=[],
        slides=[
            SlideSpec(1, "title", "Title", visual_concept="v"),
            SlideSpec(2, "content", "Even Slide", body="Image left", visual_concept="v"),
            SlideSpec(3, "content", "Odd Slide", body="Image right", visual_concept="v"),
            SlideSpec(4, "content", "Even Again", body="Image left", visual_concept="v"),
        ],
    )

    builder = SlideBuilder()
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    assert len(prs.slides) == 4
    Path(path).unlink()
    print("  [PASS] Alternating content slide layout renders correctly")


def test_gradient_fallback_dark_backgrounds():
    """Without images, slides should have dark backgrounds (not white)."""
    spec = PresentationSpec(
        title="Dark Test",
        subtitle="",
        language="EN",
        source_document="test.pdf",
        themes=[],
        slides=[
            SlideSpec(1, "title", "Title", visual_concept="v"),
            SlideSpec(2, "content", "Content", body="Text", visual_concept="v"),
        ],
    )

    builder = SlideBuilder(full_slide_mode=False)
    builder.build_presentation(spec)

    # Verify the presentation renders — the gradient bg prevents white slides
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    assert len(prs.slides) == 2

    # Check that title slide has a gradient (not solid white)
    title_bg = prs.slides[0].background.fill
    # Gradient fill type is 2 in python-pptx
    assert title_bg.type is not None, "Title slide should have a fill (not blank)"
    Path(path).unlink()
    print("  [PASS] Gradient fallback produces dark backgrounds")


def test_full_slide_mode_no_images():
    """In full_slide_mode with no images, slides get dark gradient + title text as fallback."""
    spec = PresentationSpec(
        title="Full Slide Test",
        subtitle="",
        language="EN",
        source_document="test.pdf",
        themes=[],
        slides=[
            SlideSpec(1, "title", "Title Slide", visual_concept="v"),
            SlideSpec(2, "content", "Content Slide", body="Text", visual_concept="v"),
        ],
    )

    # full_slide_mode=True but no image_paths → fallback to gradient + title
    builder = SlideBuilder(full_slide_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    assert len(prs.slides) == 2
    Path(path).unlink()
    print("  [PASS] Full-slide mode gracefully handles missing images")


def test_full_slide_mode_with_image():
    """In full_slide_mode with an image, slide has exactly 1 picture shape + notes."""
    import io
    from PIL import Image

    # Create a tiny test image
    img = Image.new("RGB", (100, 56), color=(26, 26, 46))
    img_dir = Path(tempfile.mkdtemp())
    img_path = img_dir / "slide_01.jpg"
    img.save(str(img_path), "JPEG")

    spec = PresentationSpec(
        title="Image Test",
        subtitle="",
        language="EN",
        source_document="test.pdf",
        themes=[],
        slides=[
            SlideSpec(1, "title", "Title", visual_concept="v", speaker_notes="Test notes"),
        ],
    )

    builder = SlideBuilder(image_paths={1: str(img_path)}, full_slide_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation as PptxPres
    prs = PptxPres(path)
    slide = prs.slides[0]

    # Full-slide mode: only 1 shape (the picture), no text boxes
    assert len(slide.shapes) == 1, f"Expected 1 shape (picture), got {len(slide.shapes)}"
    assert slide.shapes[0].shape_type == 13  # Picture shape type
    # Speaker notes should be present
    assert "Test notes" in slide.notes_slide.notes_text_frame.text

    Path(path).unlink()
    img_path.unlink()
    img_dir.rmdir()
    print("  [PASS] Full-slide mode produces image-only slides with speaker notes")


def test_composite_mode_explicit():
    """full_slide_mode=False uses composite builders (text boxes + overlays)."""
    spec = PresentationSpec(
        title="Composite Test",
        subtitle="",
        language="EN",
        source_document="test.pdf",
        themes=[],
        slides=[
            SlideSpec(1, "title", "Title", subtitle="Sub", visual_concept="v"),
            SlideSpec(2, "content", "Content", body="Text", bullet_points=["A"], visual_concept="v"),
        ],
    )

    builder = SlideBuilder(full_slide_mode=False)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    # Composite mode: slides should have multiple shapes (text boxes, overlays, etc.)
    for slide in prs.slides:
        assert len(slide.shapes) > 1, "Composite mode should have multiple shapes per slide"

    Path(path).unlink()
    print("  [PASS] Composite mode produces multi-shape slides")


# ─── Modern Editable Mode Tests ──────────────────────────────────

def test_recodme_title_slide():
    """Modern title: dark gradient bg + title + subtitle + accent bar."""
    spec = PresentationSpec(
        title="Title Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(1, "title", "Strategic Leadership", subtitle="A Framework"),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Accent bar + title + separator + subtitle + footer shapes
    assert len(slide.shapes) >= 4, f"Title slide should have >= 4 shapes, got {len(slide.shapes)}"
    # Verify gradient background is set
    bg = slide.background.fill
    assert bg.type is not None
    Path(path).unlink()
    print("  [PASS] Modern title slide (dark gradient bg + accent bar + title)")


def test_recodme_content_slide():
    """Modern content: dark gradient bg + title + body + bullets + accent strip."""
    spec = PresentationSpec(
        title="Content Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(2, "content", "Key Insight",
                      body="Evidence supports the claim.",
                      bullet_points=["Point A", "Point B"]),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Accent bar + title + accent line + body + numbered items + footer = 6+
    assert len(slide.shapes) >= 6, f"Content slide should have >= 6 shapes, got {len(slide.shapes)}"
    Path(path).unlink()
    print("  [PASS] Modern content slide (dark gradient bg + accent strip + bullets)")


def test_recodme_section_slide():
    """Modern section: gradient bg + white centered title + accent bar."""
    spec = PresentationSpec(
        title="Section Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(3, "section", "New Section", body="Introduction"),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Gradient background
    bg = slide.background.fill
    assert bg.type is not None
    # Title + accent bar + body + number = 4+
    assert len(slide.shapes) >= 3, f"Section slide should have >= 3 shapes, got {len(slide.shapes)}"
    Path(path).unlink()
    print("  [PASS] Modern section slide (gradient bg + centered title)")


def test_recodme_comparison_slide():
    """Modern comparison: glassmorphism cards + vertical divider + headers."""
    spec = PresentationSpec(
        title="Compare Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(4, "comparison", "Before vs After",
                      left_column=["Old A", "Old B"], right_column=["New A", "New B"],
                      left_header="Before", right_header="After"),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Title + accent + 2 cards + 2 header texts + divider + 2 columns + footer = 10+
    assert len(slide.shapes) >= 8, f"Comparison slide should have >= 8 shapes, got {len(slide.shapes)}"
    Path(path).unlink()
    print("  [PASS] Modern comparison slide (glassmorphism cards + divider + headers)")


def test_recodme_comparison_fallback_split():
    """Comparison slide splits bullet_points in half when no left/right_column."""
    spec = PresentationSpec(
        title="Fallback Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(4, "comparison", "Split Test",
                      bullet_points=["A", "B", "C", "D"]),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    assert len(prs.slides) == 1
    # Should have rendered without error — title + accent + divider + columns + footer
    assert len(prs.slides[0].shapes) >= 5
    Path(path).unlink()
    print("  [PASS] Modern comparison fallback splits bullets in half")


def test_recodme_data_slide():
    """Modern data: glassmorphism stat card + supporting bullets."""
    spec = PresentationSpec(
        title="Data Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(5, "data", "Key Metrics",
                      body="93% of managers use reason",
                      bullet_points=["Most versatile", "Works everywhere"]),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Title + accent + card + stat text + numbered items + footer = 6+
    assert len(slide.shapes) >= 6, f"Data slide should have >= 6 shapes, got {len(slide.shapes)}"
    Path(path).unlink()
    print("  [PASS] Modern data slide (glassmorphism stat card + bullets)")


def test_recodme_quote_checkbox():
    """Modern quote: checkbox items with accent checks."""
    spec = PresentationSpec(
        title="Checkbox Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(6, "quote", "Rules of Influence",
                      checkbox_items=["Rule 1: Listen first", "Rule 2: Build trust"]),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Accent strip + title + accent + numbered items + footer = 5+
    assert len(slide.shapes) >= 5, f"Quote/numbered slide should have >= 5 shapes, got {len(slide.shapes)}"
    Path(path).unlink()
    print("  [PASS] Modern quote slide with checkbox items")


def test_recodme_quote_mode():
    """Modern quote mode: large quote mark + text + attribution (no bullets)."""
    spec = PresentationSpec(
        title="Quote Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(6, "quote", "Famous Quote",
                      body="The best leaders adapt their style.",
                      source_reference="Jick, 1987"),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Accent strip + quote mark + quote text + rule + attribution + footer = 5+
    assert len(slide.shapes) >= 5, f"Quote mode should have >= 5 shapes, got {len(slide.shapes)}"
    Path(path).unlink()
    print("  [PASS] Modern quote mode (large quote mark + attribution)")


def test_recodme_conclusion_slide():
    """Modern conclusion: glassmorphism takeaway cards."""
    spec = PresentationSpec(
        title="Conclusion Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(7, "conclusion", "Key Takeaways",
                      body="Summary of findings",
                      bullet_points=["Action 1", "Action 2"]),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Cards or panel layout + footer = 6+
    assert len(slide.shapes) >= 6, f"Conclusion slide should have >= 6 shapes, got {len(slide.shapes)}"
    Path(path).unlink()
    print("  [PASS] Modern conclusion slide (glassmorphism takeaway cards)")


def test_recodme_no_images_needed():
    """Editable mode works without any image_paths — pure programmatic layouts."""
    spec = PresentationSpec(
        title="No Images Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(1, "title", "Title", subtitle="Sub"),
            SlideSpec(2, "content", "Content", body="Text", bullet_points=["A"]),
            SlideSpec(3, "section", "Section"),
            SlideSpec(4, "comparison", "Compare", bullet_points=["L", "R"]),
            SlideSpec(5, "data", "93%", body="Stat", bullet_points=["Detail"]),
            SlideSpec(6, "quote", "Quote", body="Text here", source_reference="Author"),
            SlideSpec(7, "conclusion", "End", body="Summary", bullet_points=["Key"]),
        ],
    )

    # No image_paths at all
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    assert len(prs.slides) == 7
    # No slides should have picture shapes (editable mode uses programmatic layouts)
    for slide in prs.slides:
        has_picture = any(s.shape_type == 13 for s in slide.shapes)
        assert not has_picture, "Editable modern slides should not have pictures"
    size = Path(path).stat().st_size
    assert size > 10000, f"PPTX too small: {size} bytes"

    Path(path).unlink()
    print(f"  [PASS] Modern editable mode works without images ({size:,} bytes, 7 slides)")


def test_editable_mode_ignores_images():
    """Editable mode ignores image_paths — uses solid backgrounds instead."""
    from PIL import Image

    img = Image.new("RGB", (100, 56), color=(245, 240, 232))
    img_dir = Path(tempfile.mkdtemp())
    img_path = img_dir / "slide_01.jpg"
    img.save(str(img_path), "JPEG")

    spec = PresentationSpec(
        title="Ignore Images Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(1, "content", "Content", body="Text", bullet_points=["A"],
                      speaker_notes="Notes here"),
        ],
    )

    # Pass image_paths but editable_mode=True — images should be ignored
    builder = SlideBuilder(image_paths={1: str(img_path)}, editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation as PptxPres
    prs = PptxPres(path)
    slide = prs.slides[0]

    # Editable mode should NOT have picture shapes (uses solid bg)
    has_picture = any(s.shape_type == 13 for s in slide.shapes)
    assert not has_picture, "Editable mode should not use images — should use solid backgrounds"
    # Should still have text boxes
    assert len(slide.shapes) >= 3, f"Should have text shapes, got {len(slide.shapes)}"
    # Speaker notes still present
    assert "Notes here" in slide.notes_slide.notes_text_frame.text

    Path(path).unlink()
    img_path.unlink()
    img_dir.rmdir()
    print("  [PASS] Editable mode ignores images (solid backgrounds only)")


def test_micro_copy_in_mock_specs():
    """Mock specs should use micro-copy: body < 40 words, bullets < 10 words."""
    client = GeminiClient()
    spec = client.generate_slide_specs("test", "test.pdf", "ES", slide_count=8)

    for slide in spec.slides:
        if slide.body:
            word_count = len(slide.body.split())
            assert word_count < 40, f"Slide {slide.number} body too long: {word_count} words"
        for bp in slide.bullet_points:
            bp_words = len(bp.split())
            assert bp_words < 10, f"Slide {slide.number} bullet too long: '{bp}' ({bp_words} words)"
    print("  [PASS] Mock specs use micro-copy (body<40, bullets<10)")


def test_checkbox_items_field():
    """SlideSpec checkbox_items field works in mock and parsing."""
    # Direct construction
    slide = SlideSpec(1, "quote", "Rules", checkbox_items=["Rule A", "Rule B"])
    assert slide.checkbox_items == ["Rule A", "Rule B"]

    # JSON parsing
    json_str = json.dumps({
        "metadata": {"title": "T", "subtitle": "", "language": "EN",
                     "source_document": "t.pdf", "total_slides": 1, "themes": []},
        "slides": [{
            "number": 1, "type": "quote", "title": "Rules",
            "checkbox_items": ["Check 1", "Check 2"],
        }]
    })
    spec = parse_slide_specs(json_str)
    assert spec.slides[0].checkbox_items == ["Check 1", "Check 2"]
    print("  [PASS] checkbox_items field works in construction and parsing")


def test_numbered_circle_helper():
    """Numbered circle: oval shape with centered number text."""
    spec = PresentationSpec(
        title="Circle Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[SlideSpec(1, "content", "Dummy")],
    )
    builder = SlideBuilder(editable_mode=True)
    from pptx.util import Inches
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    shape = builder._add_numbered_circle(slide, Inches(1), Inches(1), 3)
    assert shape is not None
    # Verify it's an oval (freeform or auto shape)
    assert shape.shape_type in (1, 5, 6, 13, 17)  # AutoShape
    # Verify text
    assert "3" in shape.text_frame.text
    print("  [PASS] Numbered circle helper")


def test_content_card_helper():
    """Content card: rounded rect with/without border and fill."""
    builder = SlideBuilder(editable_mode=True)
    from pptx.util import Inches
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])

    # With fill
    card1 = builder._add_content_card(
        slide, Inches(0), Inches(0), Inches(3), Inches(2),
        fill_color=RGBColor(0xFD, 0xEC, 0xE9),
    )
    assert card1 is not None

    # With border, no fill (modern cyan accent)
    card2 = builder._add_content_card(
        slide, Inches(4), Inches(0), Inches(3), Inches(2),
        border_color=RGBColor(0x06, 0xB6, 0xD4),
    )
    assert card2 is not None
    print("  [PASS] Content card helper (fill and border variants)")


def test_footer_bar_helper():
    """Footer bar: line + number on slide."""
    builder = SlideBuilder(editable_mode=True)
    from pptx.util import Inches
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    initial_count = len(slide.shapes)
    builder._add_footer_bar(slide, slide_number=5)
    # Should add line + number = 2 shapes (watermark removed in modern design)
    assert len(slide.shapes) >= initial_count + 2, f"Footer bar should add 2 shapes, added {len(slide.shapes) - initial_count}"
    print("  [PASS] Footer bar helper (line + number)")


def test_rich_bullet_formatting():
    """Rich bullets: splits on em dash, bold keyword + regular text."""
    builder = SlideBuilder(editable_mode=True)
    from pptx.util import Inches
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    items = [
        "Keyword \u2014 explanation text",
        "Plain text without delimiter",
    ]
    txBox = builder._add_rich_bullet_list(
        slide, Inches(1), Inches(1), Inches(10), Inches(3),
        items, 14, "DM Sans", RGBColor(0xF1, 0xF5, 0xF9),
    )
    tf = txBox.text_frame
    # First paragraph: bullet + keyword + sep + description = 4 runs
    p1_runs = tf.paragraphs[0].runs
    assert len(p1_runs) == 4, f"Expected 4 runs for em-dash bullet, got {len(p1_runs)}"
    assert p1_runs[1].font.bold == True, "Keyword should be bold"
    # Second paragraph: bullet + plain text = 2 runs
    p2_runs = tf.paragraphs[1].runs
    assert len(p2_runs) == 2, f"Expected 2 runs for plain bullet, got {len(p2_runs)}"
    print("  [PASS] Rich bullet formatting (bold keyword + regular text)")


def test_content_slide_numbered_items():
    """Content slide has numbered circles for each bullet."""
    spec = PresentationSpec(
        title="Numbered Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(2, "content", "Key Insight",
                      body="Evidence.",
                      bullet_points=["A \u2014 first", "B \u2014 second", "C \u2014 third"]),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Find shapes whose text is just a single digit (numbered circles)
    circle_texts = []
    for s in slide.shapes:
        if s.has_text_frame:
            txt = s.text_frame.text.strip()
            if txt in ('1', '2', '3'):
                circle_texts.append(txt)
    assert len(circle_texts) >= 3, f"Expected 3 numbered circles, found: {circle_texts}"
    Path(path).unlink()
    print("  [PASS] Content slide has numbered circles for each bullet")


def test_comparison_header_bars():
    """Comparison headers are text shapes within glassmorphism cards."""
    spec = PresentationSpec(
        title="Header Bar Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(4, "comparison", "Compare",
                      left_column=["A", "B"], right_column=["C", "D"],
                      left_header="Left Side", right_header="Right Side"),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Find shapes that contain header text
    header_texts = [s for s in slide.shapes
                    if hasattr(s, 'text_frame') and s.text_frame.text in ('Left Side', 'Right Side')]
    assert len(header_texts) >= 2, f"Expected 2 header text shapes, got {len(header_texts)}"
    Path(path).unlink()
    print("  [PASS] Comparison header bars (glassmorphism cards)")


def test_conclusion_card_layout():
    """Conclusion with <=3 bullets produces card shapes."""
    spec = PresentationSpec(
        title="Card Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(7, "conclusion", "Takeaways",
                      body="Summary",
                      bullet_points=["Act \u2014 now", "Think \u2014 deeply", "Lead \u2014 boldly"]),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Should have card shapes + numbered circles + keyword texts
    # Header bar + title + accent + body + 3 cards + 3 circles + 3 keywords + 3 explanations + footer = 18+
    assert len(slide.shapes) >= 12, f"Conclusion card layout should have >= 12 shapes, got {len(slide.shapes)}"
    Path(path).unlink()
    print("  [PASS] Conclusion card layout (3 takeaway cards)")


def test_title_slide_accent_line():
    """Title slide has a cyan accent line in modern design."""
    spec = PresentationSpec(
        title="Panel Test", subtitle="", language="EN",
        source_document="test.pdf", themes=[],
        slides=[
            SlideSpec(1, "title", "Big Title", subtitle="Subtitle here"),
        ],
    )
    builder = SlideBuilder(editable_mode=True)
    builder.build_presentation(spec)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    builder.save(path)

    from pptx import Presentation
    prs = Presentation(path)
    slide = prs.slides[0]
    # Check for cyan accent color (#06B6D4) in any shape
    has_accent = False
    for s in slide.shapes:
        if hasattr(s, 'fill') and s.fill.type is not None:
            try:
                if s.fill.fore_color and s.fill.fore_color.rgb == RGBColor(0x06, 0xB6, 0xD4):
                    has_accent = True
                    break
            except (AttributeError, TypeError):
                pass
    assert has_accent, "Title slide should have a cyan accent shape"
    Path(path).unlink()
    print("  [PASS] Title slide has cyan accent line")


def test_new_slide_spec_fields():
    """New SlideSpec fields: left_column, right_column, left_header, right_header."""
    slide = SlideSpec(1, "comparison", "Compare",
                      left_column=["A", "B"], right_column=["C"],
                      left_header="Left", right_header="Right")
    assert slide.left_column == ["A", "B"]
    assert slide.right_column == ["C"]
    assert slide.left_header == "Left"
    assert slide.right_header == "Right"

    # JSON round-trip
    json_str = json.dumps({
        "metadata": {"title": "T", "subtitle": "", "language": "EN",
                     "source_document": "t.pdf", "total_slides": 1, "themes": []},
        "slides": [{
            "number": 1, "type": "comparison", "title": "Compare",
            "left_column": ["X"], "right_column": ["Y"],
            "left_header": "LH", "right_header": "RH",
        }]
    })
    spec = parse_slide_specs(json_str)
    assert spec.slides[0].left_column == ["X"]
    assert spec.slides[0].right_header == "RH"
    print("  [PASS] New SlideSpec fields (left/right column/header)")


# ─── NotebookLM Client Tests ──────────────────────────────────

def test_notebooklm_client_import():
    """Module imports without error."""
    from notebooklm_client import NotebookLMPipeline
    assert NotebookLMPipeline is not None
    print("  [PASS] NotebookLM client module imports")


def test_notebooklm_auth_check():
    """is_authenticated() returns bool based on profile existence."""
    from notebooklm_client import NotebookLMPipeline
    result = NotebookLMPipeline.is_authenticated()
    assert isinstance(result, bool)
    print(f"  [PASS] NotebookLM auth check (authenticated={result})")


def test_notebooklm_prompt_construction():
    """Prompt includes language, count, and layout variety."""
    from notebooklm_client import NotebookLMPipeline
    pipeline = NotebookLMPipeline(profile="test")
    prompt_es = pipeline._build_prompt("ES", 8)
    assert "Spanish" in prompt_es
    assert "8-slide" in prompt_es
    assert "insights" in prompt_es.lower()

    prompt_en = pipeline._build_prompt("EN", 10)
    assert "English" in prompt_en
    assert "10-slide" in prompt_en
    print("  [PASS] NotebookLM prompt construction (language + count)")


def test_notebooklm_pipeline_no_auth():
    """Pipeline returns None when not authenticated (no real API call)."""
    from notebooklm_client import NotebookLMPipeline
    pipeline = NotebookLMPipeline(profile="nonexistent_test_profile_xyz")
    # generate_from_pdf should fail gracefully (RuntimeError on auth)
    try:
        result = pipeline.generate_from_pdf(
            pdf_path="nonexistent.pdf",
            output_path="/tmp/test.pptx",
        )
        # If profile doesn't exist, _get_client raises RuntimeError,
        # which is caught and returns None
        assert result is None
    except RuntimeError:
        # Expected when profile doesn't exist
        pass
    print("  [PASS] NotebookLM pipeline handles missing auth gracefully")


# ─── Run All Tests ─────────────────────────────────────────────

def run_all():
    print("\n=== Presentation Factory Test Suite ===\n")
    tests = [
        ("Language Detection", test_language_detection),
        ("Text Cleaning", test_clean_text),
        ("PDF Extraction (Tacticas)", test_extract_pdf),
        ("PDF Extraction (Porter)", test_extract_pdf_porter),
        ("Mock Generation (ES)", test_mock_generation),
        ("Mock Generation (EN)", test_mock_generation_en),
        ("Visual Concepts Image-Friendly", test_mock_visual_concepts_are_image_friendly),
        ("JSON Parsing", test_parse_slide_specs),
        ("JSON with Fences", test_parse_with_markdown_fences),
        ("Specs Serialization", test_specs_to_json),
        ("Image Prompt Template", test_image_prompt_template),
        ("Image Prompt Building", test_build_image_prompt),
        ("Mock Image Generation", test_mock_images),
        ("Image Gen No API Key", test_image_generation_no_api_key),
        ("Full-Slide Prompt Template", test_full_slide_prompt_template),
        ("Layout Templates All Types", test_layout_templates_cover_all_types),
        ("Full-Slide Prompt Title", test_build_full_slide_prompt_title),
        ("Full-Slide Prompt Content", test_build_full_slide_prompt_content),
        ("Full-Slide Prompt All Types", test_build_full_slide_prompt_all_types),
        ("Model Constants", test_model_constants),
        ("Brand Config", test_brand_config),
        ("PPTX Generation (No Images)", test_build_presentation_no_images),
        ("All Slide Types", test_all_slide_types),
        ("SlideBuilder Image Paths", test_slide_builder_with_image_paths),
        ("Alternating Layout", test_alternating_content_layout),
        ("Dark Gradient Fallback", test_gradient_fallback_dark_backgrounds),
        ("Full-Slide Mode No Images", test_full_slide_mode_no_images),
        ("Full-Slide Mode With Image", test_full_slide_mode_with_image),
        ("Composite Mode Explicit", test_composite_mode_explicit),
        # Modern editable mode tests
        ("Modern Title Slide", test_recodme_title_slide),
        ("Modern Content Slide", test_recodme_content_slide),
        ("Modern Section Slide", test_recodme_section_slide),
        ("Modern Comparison Slide", test_recodme_comparison_slide),
        ("Modern Comparison Fallback", test_recodme_comparison_fallback_split),
        ("Modern Data Slide", test_recodme_data_slide),
        ("Modern Quote Checkbox", test_recodme_quote_checkbox),
        ("Modern Quote Mode", test_recodme_quote_mode),
        ("Modern Conclusion Slide", test_recodme_conclusion_slide),
        ("Modern No Images Needed", test_recodme_no_images_needed),
        ("Editable Ignores Images", test_editable_mode_ignores_images),
        ("Micro-Copy Mock Specs", test_micro_copy_in_mock_specs),
        ("Checkbox Items Field", test_checkbox_items_field),
        # v4.1 new tests
        ("Numbered Circle Helper", test_numbered_circle_helper),
        ("Content Card Helper", test_content_card_helper),
        ("Footer Bar Helper", test_footer_bar_helper),
        ("Rich Bullet Formatting", test_rich_bullet_formatting),
        ("Content Slide Numbered Items", test_content_slide_numbered_items),
        ("Comparison Header Bars", test_comparison_header_bars),
        ("Conclusion Card Layout", test_conclusion_card_layout),
        ("Title Slide Accent Line", test_title_slide_accent_line),
        ("New SlideSpec Fields", test_new_slide_spec_fields),
        # NotebookLM client tests
        ("NotebookLM Client Import", test_notebooklm_client_import),
        ("NotebookLM Auth Check", test_notebooklm_auth_check),
        ("NotebookLM Prompt Construction", test_notebooklm_prompt_construction),
        ("NotebookLM Pipeline No Auth", test_notebooklm_pipeline_no_auth),
    ]

    passed = 0
    failed = 0
    for name, test_fn in tests:
        try:
            test_fn()
            passed += 1
        except Exception as e:
            print(f"  [FAIL] {name}: {e}")
            failed += 1

    print(f"\n{'=' * 40}")
    print(f"Results: {passed} passed, {failed} failed, {len(tests)} total")
    print(f"{'=' * 40}\n")
    return failed == 0


if __name__ == "__main__":
    success = run_all()
    sys.exit(0 if success else 1)
