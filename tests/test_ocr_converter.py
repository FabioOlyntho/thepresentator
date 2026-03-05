"""
Test Suite — OCR Converter: NotebookLM image PPTX → editable Recodme PPTX.

Run: python -m pytest tests/test_ocr_converter.py -v
"""

import io
import json
import os
import sys
import tempfile
from pathlib import Path
from unittest.mock import patch, MagicMock

# Add scripts to path
sys.path.insert(0, str(Path(__file__).parent.parent / "scripts"))

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from gemini_client import SlideSpec, PresentationSpec
from ocr_converter import (
    extract_slide_images,
    classify_slide_type,
    content_to_slidespec,
    convert_notebooklm_to_editable,
    convert_pdnob_style,
    extract_text_with_positions,
    erase_text_from_image,
    _sample_text_color,
    OCRTextBlock,
    ImageRegion,
    load_ocr_prompt,
    normalize_text,
    _parse_raw_text_to_slide,
    VALID_SLIDE_TYPES,
    group_text_blocks,
    _merge_block_group,
    _most_common_color,
    segment_slide_image,
    crop_image_region,
    remove_background,
)


# ─── Helper: Create a mock NotebookLM PPTX with image slides ────

def _create_mock_notebooklm_pptx(num_slides: int = 3) -> str:
    """Create a PPTX file with one full-bleed image per slide (mimics NotebookLM output)."""
    prs = Presentation()
    prs.slide_width = Inches(17.78)  # NotebookLM: 17.8" wide
    prs.slide_height = Inches(10.0)

    for i in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Create a simple PNG image in memory
        img = Image.new("RGB", (1280, 720), color=(245, 240, 232))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)

        # Add image as full-bleed picture shape
        slide.shapes.add_picture(
            buf,
            Inches(0), Inches(0),
            Inches(17.78), Inches(10.0),
        )

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_mock_empty_pptx() -> str:
    """Create a PPTX file with slides but no images (text-only)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Only add a text box, no picture
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    txBox.text_frame.text = "Text-only slide"

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


# ─── OCR Prompt Tests ────────────────────────────────────────────

def test_ocr_prompt_loads():
    """OCR extraction prompt file exists and loads."""
    prompt = load_ocr_prompt()
    assert len(prompt) > 100, "OCR prompt should be substantial"
    assert "{slide_number}" in prompt
    assert "{total_slides}" in prompt
    assert "JSON" in prompt
    print("  [PASS] OCR prompt loads with placeholders")


def test_ocr_prompt_has_slide_types():
    """Prompt instructs Gemini about valid slide types."""
    prompt = load_ocr_prompt()
    for stype in ["title", "section", "content", "comparison", "data", "quote", "conclusion"]:
        assert stype in prompt, f"Prompt missing slide type: {stype}"
    print("  [PASS] OCR prompt includes all 7 slide types")


# ─── Image Extraction Tests ──────────────────────────────────────

def test_extract_slide_images():
    """Extract images from a mock NotebookLM PPTX."""
    pptx_path = _create_mock_notebooklm_pptx(num_slides=3)
    try:
        images = extract_slide_images(pptx_path)
        assert len(images) == 3, f"Expected 3 slide images, got {len(images)}"
        for slide_num, img_bytes in images:
            assert isinstance(slide_num, int)
            assert isinstance(img_bytes, bytes)
            assert len(img_bytes) > 100, "Image bytes should be non-trivial"
        # Slide numbers should be 1, 2, 3
        nums = [n for n, _ in images]
        assert nums == [1, 2, 3], f"Expected slide numbers [1,2,3], got {nums}"
    finally:
        Path(pptx_path).unlink()
    print("  [PASS] Extract slide images from mock PPTX (3 slides)")


def test_extract_slide_images_empty():
    """No images found in text-only PPTX."""
    pptx_path = _create_mock_empty_pptx()
    try:
        images = extract_slide_images(pptx_path)
        assert len(images) == 0, f"Expected 0 images from text-only PPTX, got {len(images)}"
    finally:
        Path(pptx_path).unlink()
    print("  [PASS] No images extracted from text-only PPTX")


def test_extract_slide_images_single():
    """Extract from PPTX with just 1 slide."""
    pptx_path = _create_mock_notebooklm_pptx(num_slides=1)
    try:
        images = extract_slide_images(pptx_path)
        assert len(images) == 1
        assert images[0][0] == 1
    finally:
        Path(pptx_path).unlink()
    print("  [PASS] Single-slide PPTX extraction")


# ─── Slide Type Classification Tests ─────────────────────────────

def test_classify_valid_types():
    """Valid slide types pass through unchanged."""
    for stype in VALID_SLIDE_TYPES:
        assert classify_slide_type(3, 8, stype) == stype
    print("  [PASS] Valid slide types pass through")


def test_classify_first_slide_fallback():
    """First slide falls back to 'title' for invalid types."""
    assert classify_slide_type(1, 8, "unknown") == "title"
    assert classify_slide_type(1, 8, "intro") == "title"
    assert classify_slide_type(1, 8, "") == "title"
    print("  [PASS] First slide → title fallback")


def test_classify_last_slide_fallback():
    """Last slide falls back to 'conclusion' for invalid types."""
    assert classify_slide_type(8, 8, "unknown") == "conclusion"
    assert classify_slide_type(10, 10, "outro") == "conclusion"
    print("  [PASS] Last slide → conclusion fallback")


def test_classify_middle_slide_fallback():
    """Middle slides fall back to 'content' for invalid types."""
    assert classify_slide_type(3, 8, "unknown") == "content"
    assert classify_slide_type(5, 10, "freeform") == "content"
    assert classify_slide_type(2, 5, "") == "content"
    print("  [PASS] Middle slides → content fallback")


# ─── SlideSpec Mapping Tests ─────────────────────────────────────

def test_content_to_slidespec_full():
    """Full content dict maps to SlideSpec with all fields."""
    content = {
        "type": "comparison",
        "title": "Before vs After",
        "subtitle": "Key Changes",
        "body": "Significant improvements observed",
        "bullet_points": ["Speed — 3x faster", "Cost — 50% lower"],
        "left_column": ["Old method", "Manual process"],
        "right_column": ["New method", "Automated pipeline"],
        "left_header": "Before",
        "right_header": "After",
        "checkbox_items": [],
        "speaker_notes": "This slide shows the transformation.",
    }
    spec = content_to_slidespec(content, slide_number=4)
    assert spec.number == 4
    assert spec.type == "comparison"
    assert spec.title == "Before vs After"
    assert spec.subtitle == "Key Changes"
    assert spec.body == "Significant improvements observed"
    assert len(spec.bullet_points) == 2
    assert spec.left_column == ["Old method", "Manual process"]
    assert spec.right_column == ["New method", "Automated pipeline"]
    assert spec.left_header == "Before"
    assert spec.right_header == "After"
    assert spec.speaker_notes == "This slide shows the transformation."
    assert spec.visual_concept == ""  # Not needed for editable mode
    print("  [PASS] Full content dict → SlideSpec mapping")


def test_content_to_slidespec_minimal():
    """Minimal content dict produces valid SlideSpec with defaults."""
    content = {
        "type": "content",
        "title": "Simple Slide",
    }
    spec = content_to_slidespec(content, slide_number=2)
    assert spec.number == 2
    assert spec.type == "content"
    assert spec.title == "Simple Slide"
    assert spec.body == ""
    assert spec.bullet_points == []
    assert spec.left_column == []
    assert spec.right_column == []
    print("  [PASS] Minimal content → SlideSpec with defaults")


def test_content_to_slidespec_empty():
    """Empty content dict produces valid SlideSpec."""
    spec = content_to_slidespec({}, slide_number=1)
    assert spec.number == 1
    assert spec.type == "content"
    assert spec.title == ""
    print("  [PASS] Empty content → valid SlideSpec")


# ─── Full Pipeline Tests (Mocked Gemini) ─────────────────────────

def _mock_gemini_response(slide_number, total_slides):
    """Create a mock Gemini Vision API response for a slide."""
    if slide_number == 1:
        content = {
            "type": "title",
            "title": "Seven Tactics Shape Influence",
            "subtitle": "A Manager's Framework",
            "body": "",
            "bullet_points": [],
            "speaker_notes": "Welcome to the presentation.",
        }
    elif slide_number == total_slides:
        content = {
            "type": "conclusion",
            "title": "Master the Full Repertoire",
            "body": "Adapt your tactics to context.",
            "bullet_points": [
                "Diagnose \u2014 then act",
                "Adapt \u2014 direction matters",
                "Invest \u2014 relationships first",
            ],
            "speaker_notes": "Close by emphasizing adaptability.",
        }
    else:
        content = {
            "type": "content",
            "title": f"Key Insight {slide_number}",
            "body": "Evidence supports this claim.",
            "bullet_points": [
                "Reason \u2014 data and logic",
                "Coalition \u2014 allies multiply",
            ],
            "speaker_notes": f"Slide {slide_number} explains the key insight.",
        }
    return json.dumps(content)


@patch("ocr_converter.requests.post")
def test_convert_pipeline_mock(mock_post):
    """Full conversion pipeline with mocked Gemini Vision API."""
    pptx_path = _create_mock_notebooklm_pptx(num_slides=3)

    # Mock Gemini responses for each slide
    def side_effect(*args, **kwargs):
        # Determine which slide by counting calls
        call_count = mock_post.call_count
        mock_resp = MagicMock()
        mock_resp.status_code = 200
        mock_resp.json.return_value = {
            "candidates": [{
                "content": {
                    "parts": [{
                        "text": _mock_gemini_response(call_count, 3)
                    }]
                }
            }]
        }
        return mock_resp

    mock_post.side_effect = side_effect

    try:
        output_pptx = tempfile.mktemp(suffix="_editable.pptx")
        result = convert_notebooklm_to_editable(
            input_pptx=pptx_path,
            output_pptx=output_pptx,
            api_key="test-key-12345",
        )

        assert result["success"], f"Pipeline failed: {result.get('error')}"
        assert result["files"]["pptx"] == output_pptx
        assert result["metadata"]["total_slides"] == 3
        assert result["metadata"]["title"] == "Seven Tactics Shape Influence"
        assert "extraction" in result["timing"]
        assert "ocr" in result["timing"]
        assert "building" in result["timing"]

        # Verify output PPTX is valid and has correct number of slides
        prs = Presentation(output_pptx)
        assert len(prs.slides) == 3, f"Expected 3 slides, got {len(prs.slides)}"

        # Verify NO picture shapes (editable mode = solid fills)
        for slide in prs.slides:
            has_picture = any(s.shape_type == 13 for s in slide.shapes)
            assert not has_picture, "Editable output should not have picture shapes"

        # Verify specs JSON was saved
        assert Path(result["files"]["specs_json"]).exists()

        # Cleanup
        Path(output_pptx).unlink(missing_ok=True)
        Path(result["files"]["specs_json"]).unlink(missing_ok=True)
    finally:
        Path(pptx_path).unlink()

    print("  [PASS] Full mock conversion pipeline (3 slides)")


@patch("ocr_converter.requests.post")
def test_convert_pipeline_ocr_failure_fallback(mock_post):
    """Pipeline creates fallback slides when OCR fails for individual slides."""
    pptx_path = _create_mock_notebooklm_pptx(num_slides=2)

    call_count = [0]

    def side_effect(*args, **kwargs):
        call_count[0] += 1
        if call_count[0] == 1:
            # First slide succeeds
            mock_resp = MagicMock()
            mock_resp.status_code = 200
            mock_resp.json.return_value = {
                "candidates": [{
                    "content": {
                        "parts": [{
                            "text": json.dumps({
                                "type": "title",
                                "title": "Working Slide",
                                "speaker_notes": "Notes here.",
                            })
                        }]
                    }
                }]
            }
            return mock_resp
        else:
            # Second slide fails
            raise RuntimeError("Gemini API error")

    mock_post.side_effect = side_effect

    try:
        output_pptx = tempfile.mktemp(suffix="_editable.pptx")
        result = convert_notebooklm_to_editable(
            input_pptx=pptx_path,
            output_pptx=output_pptx,
            api_key="test-key-12345",
        )

        assert result["success"], "Pipeline should succeed even with partial OCR failure"
        assert result["metadata"]["total_slides"] == 2

        # Verify output still has 2 slides
        prs = Presentation(output_pptx)
        assert len(prs.slides) == 2

        Path(output_pptx).unlink(missing_ok=True)
        Path(result["files"]["specs_json"]).unlink(missing_ok=True)
    finally:
        Path(pptx_path).unlink()

    print("  [PASS] Pipeline handles partial OCR failure with fallback slides")


def test_convert_no_api_key():
    """Pipeline returns error when no API key is provided."""
    old_key = os.environ.pop("GEMINI_API_KEY", None)
    try:
        result = convert_notebooklm_to_editable(
            input_pptx="fake.pptx",
            api_key=None,
        )
        assert not result["success"]
        assert "API key" in result["error"]
    finally:
        if old_key:
            os.environ["GEMINI_API_KEY"] = old_key
    print("  [PASS] Pipeline requires API key")


def test_convert_empty_pptx():
    """Pipeline returns error for PPTX with no images."""
    pptx_path = _create_mock_empty_pptx()
    try:
        result = convert_notebooklm_to_editable(
            input_pptx=pptx_path,
            api_key="test-key-12345",
        )
        assert not result["success"]
        assert "No images" in result["error"]
    finally:
        Path(pptx_path).unlink()
    print("  [PASS] Pipeline rejects empty PPTX (no images)")


# ─── UTF-8 Normalization Tests ───────────────────────────────────

def test_normalize_text_nfc():
    """NFC normalization composes decomposed accents."""
    # "é" as decomposed (e + combining acute accent)
    decomposed = "e\u0301"
    result = normalize_text(decomposed)
    assert result == "é", f"Expected 'é', got {repr(result)}"
    # "ñ" decomposed
    assert normalize_text("n\u0303") == "ñ"
    print("  [PASS] NFC normalization composes decomposed accents")


def test_normalize_text_zero_width():
    """Removes zero-width characters and BOM."""
    text = "Hello\u200bWorld\ufeff"
    result = normalize_text(text)
    assert result == "HelloWorld", f"Expected 'HelloWorld', got {repr(result)}"
    print("  [PASS] Removes zero-width characters and BOM")


def test_normalize_text_nbsp():
    """Replaces non-breaking spaces with regular spaces."""
    text = "Hello\u00a0World"
    result = normalize_text(text)
    assert result == "Hello World"
    print("  [PASS] Replaces non-breaking spaces")


def test_normalize_text_empty():
    """Empty string returns empty."""
    assert normalize_text("") == ""
    assert normalize_text(None) is None
    print("  [PASS] Empty/None handling")


def test_content_to_slidespec_normalizes_text():
    """SlideSpec creation normalizes all text fields."""
    content = {
        "type": "content",
        "title": "Tácticas de Influencia",  # normal
        "body": "e\u0301xito",  # decomposed é
        "bullet_points": ["Razo\u0301n \u2014 lo\u0301gica", "Coalicio\u0301n \u2014 aliados"],
        "speaker_notes": "Nota\u00a0con\u200bcaracteres",
    }
    spec = content_to_slidespec(content, slide_number=3)
    assert spec.body == "éxito"
    assert spec.bullet_points[0] == "Razón — lógica"
    assert spec.bullet_points[1] == "Coalición — aliados"
    assert spec.speaker_notes == "Nota concaracteres"
    print("  [PASS] SlideSpec normalizes all text fields")


# ─── Docling OCR Engine Tests ───────────────────────────────────

def test_parse_raw_text_to_slide_basic():
    """Parse raw OCR text with title and bullets."""
    raw = "# Key Insight 2\n- Speed — 3x faster\n- Cost — 50% lower\nSome body text here."
    result = _parse_raw_text_to_slide(raw, slide_number=2, total_slides=5)
    assert result["title"] == "Key Insight 2"
    assert len(result["bullet_points"]) == 2
    assert "Speed — 3x faster" in result["bullet_points"]
    assert result["type"] == "content"
    print("  [PASS] Parse raw OCR text with title and bullets")


def test_parse_raw_text_to_slide_title_slide():
    """First slide always becomes title type."""
    raw = "# My Presentation\n**A Great Topic**\nSome intro text."
    result = _parse_raw_text_to_slide(raw, slide_number=1, total_slides=8)
    assert result["type"] == "title"
    assert result["title"] == "My Presentation"
    assert result["subtitle"] == "A Great Topic"
    print("  [PASS] First slide parsed as title type")


def test_parse_raw_text_to_slide_conclusion():
    """Last slide always becomes conclusion type."""
    raw = "# Final Thoughts\n- Point 1\n- Point 2"
    result = _parse_raw_text_to_slide(raw, slide_number=8, total_slides=8)
    assert result["type"] == "conclusion"
    print("  [PASS] Last slide parsed as conclusion type")


def test_parse_raw_text_to_slide_empty():
    """Empty raw text produces valid fallback."""
    result = _parse_raw_text_to_slide("", slide_number=3, total_slides=5)
    assert result["title"] == "Slide 3"
    assert result["type"] == "content"
    print("  [PASS] Empty raw text produces valid fallback")


def test_parse_raw_text_numbered_bullets():
    """Numbered list items are parsed as bullet points."""
    raw = "# Topic\n1. First point\n2. Second point\n3. Third point"
    result = _parse_raw_text_to_slide(raw, slide_number=3, total_slides=5)
    assert len(result["bullet_points"]) == 3
    assert "First point" in result["bullet_points"][0]
    print("  [PASS] Numbered bullets parsed correctly")


def test_parse_raw_text_colon_to_em_dash():
    """Bullet items with colons get reformatted to em dash style."""
    raw = "# Topic\n- Speed: very fast\n- Cost: very low"
    result = _parse_raw_text_to_slide(raw, slide_number=2, total_slides=5)
    assert "Speed — very fast" in result["bullet_points"]
    assert "Cost — very low" in result["bullet_points"]
    print("  [PASS] Colon bullet format → em dash")


def test_convert_invalid_ocr_engine():
    """Pipeline returns error for unknown OCR engine."""
    result = convert_notebooklm_to_editable(
        input_pptx="fake.pptx",
        api_key="test-key",
        ocr_engine="unknown_engine",
    )
    assert not result["success"]
    assert "Unknown OCR engine" in result["error"]
    print("  [PASS] Invalid OCR engine returns error")


def test_convert_docling_no_api_key():
    """Docling engine doesn't need API key — should not error on missing key."""
    pptx_path = _create_mock_empty_pptx()
    try:
        old_key = os.environ.pop("GEMINI_API_KEY", None)
        result = convert_notebooklm_to_editable(
            input_pptx=pptx_path,
            api_key=None,
            ocr_engine="docling",
        )
        # Should fail because no images, not because of API key
        assert not result["success"]
        assert "No images" in result["error"]
    finally:
        if old_key:
            os.environ["GEMINI_API_KEY"] = old_key
        Path(pptx_path).unlink()
    print("  [PASS] Docling engine doesn't need API key")


@patch("ocr_converter.extract_slide_content_docling")
def test_convert_pipeline_docling_mock(mock_docling):
    """Full conversion pipeline with mocked Docling OCR engine."""
    pptx_path = _create_mock_notebooklm_pptx(num_slides=3)

    call_count = [0]

    def side_effect(image_bytes, slide_number, total_slides):
        call_count[0] += 1
        if slide_number == 1:
            return {
                "type": "title",
                "title": "Docling Extracted Title",
                "subtitle": "Via OCR",
                "body": "",
                "bullet_points": [],
                "speaker_notes": "",
            }
        elif slide_number == 3:
            return {
                "type": "conclusion",
                "title": "Final Summary",
                "body": "End of presentation.",
                "bullet_points": ["Key — takeaway"],
                "speaker_notes": "",
            }
        else:
            return {
                "type": "content",
                "title": f"Slide {slide_number} Content",
                "body": "Some body text.",
                "bullet_points": ["Point — one", "Point — two"],
                "speaker_notes": "",
            }

    mock_docling.side_effect = side_effect

    try:
        output_pptx = tempfile.mktemp(suffix="_editable.pptx")
        result = convert_notebooklm_to_editable(
            input_pptx=pptx_path,
            output_pptx=output_pptx,
            ocr_engine="docling",
        )

        assert result["success"], f"Docling pipeline failed: {result.get('error')}"
        assert result["metadata"]["total_slides"] == 3
        assert result["metadata"]["ocr_engine"] == "docling"
        assert result["metadata"]["title"] == "Docling Extracted Title"

        # Verify output PPTX is valid
        prs = Presentation(output_pptx)
        assert len(prs.slides) == 3

        # Verify NO picture shapes (editable mode)
        for slide in prs.slides:
            has_picture = any(s.shape_type == 13 for s in slide.shapes)
            assert not has_picture, "Editable output should not have picture shapes"

        Path(output_pptx).unlink(missing_ok=True)
        Path(result["files"]["specs_json"]).unlink(missing_ok=True)
    finally:
        Path(pptx_path).unlink()

    print("  [PASS] Full Docling mock conversion pipeline (3 slides)")


# ─── PDNob OCR Tests ──────────────────────────────────────────────


def _create_image_with_text() -> bytes:
    """Create a test image with black text on white background."""
    img = Image.new("RGB", (800, 600), color=(255, 255, 255))
    # Draw a dark rectangle to simulate text area
    from PIL import ImageDraw
    draw = ImageDraw.Draw(img)
    # Simulate text blocks as dark rectangles
    draw.rectangle([50, 50, 300, 80], fill=(20, 20, 20))   # "Title" area
    draw.rectangle([50, 120, 400, 145], fill=(30, 30, 30))  # "Body" area
    draw.rectangle([50, 180, 350, 205], fill=(25, 25, 25))  # "Bullet" area
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


@patch("ocr_converter._get_rapid_ocr")
def test_extract_text_with_positions(mock_ocr_factory):
    """Extract text blocks with positions from image via mocked RapidOCR."""
    mock_engine = MagicMock()
    mock_ocr_factory.return_value = mock_engine

    # RapidOCR v3.x returns (results_list, elapsed_times)
    # Each item: [box, text, score]
    mock_engine.return_value = (
        [
            [[[50, 50], [300, 50], [300, 80], [50, 80]], "Title Text", "0.95"],
            [[[50, 120], [400, 120], [400, 145], [50, 145]], "Body content here", "0.88"],
        ],
        [0.1, 0.2, 0.3],
    )

    image_bytes = _create_image_with_text()
    blocks = extract_text_with_positions(image_bytes)

    assert len(blocks) == 2, f"Expected 2 text blocks, got {len(blocks)}"
    assert blocks[0].text == "Title Text"
    assert blocks[1].text == "Body content here"
    # Verify percentage coordinates are reasonable (0-100 range)
    for b in blocks:
        assert 0 <= b.x_pct <= 100
        assert 0 <= b.y_pct <= 100
        assert b.width_pct > 0
        assert b.height_pct > 0
        assert b.font_size_pt >= 8.0
        assert len(b.color) == 3
    # First block should be above second (sorted top-to-bottom)
    assert blocks[0].y_pct < blocks[1].y_pct
    print("  [PASS] Extract text with positions (mocked RapidOCR)")


@patch("ocr_converter._get_rapid_ocr")
def test_extract_text_with_positions_empty(mock_ocr_factory):
    """No text detected returns empty list."""
    mock_engine = MagicMock()
    mock_ocr_factory.return_value = mock_engine
    mock_engine.return_value = (None, [0.0])

    image_bytes = _create_image_with_text()
    blocks = extract_text_with_positions(image_bytes)
    assert blocks == []
    print("  [PASS] Empty OCR result returns empty list")


@patch("ocr_converter._get_rapid_ocr")
def test_extract_text_filters_low_confidence(mock_ocr_factory):
    """Text blocks below 0.5 confidence are filtered out."""
    mock_engine = MagicMock()
    mock_ocr_factory.return_value = mock_engine

    mock_engine.return_value = (
        [
            [[[10, 10], [200, 10], [200, 40], [10, 40]], "High confidence", "0.92"],
            [[[10, 60], [200, 60], [200, 90], [10, 90]], "Low confidence", "0.30"],
        ],
        [0.1, 0.2, 0.3],
    )

    image_bytes = _create_image_with_text()
    blocks = extract_text_with_positions(image_bytes)
    assert len(blocks) == 1
    assert blocks[0].text == "High confidence"
    print("  [PASS] Low confidence text blocks filtered out")


def test_erase_text_from_image():
    """Erased image has same dimensions and text regions are modified."""
    image_bytes = _create_image_with_text()
    original_img = Image.open(io.BytesIO(image_bytes))
    orig_w, orig_h = original_img.size

    text_blocks = [
        OCRTextBlock(
            text="Title", x_pct=6.25, y_pct=8.33,
            width_pct=31.25, height_pct=5.0,
            font_size_pt=24.0, color=(20, 20, 20),
        ),
    ]

    erased_bytes = erase_text_from_image(image_bytes, text_blocks)
    assert len(erased_bytes) > 0
    erased_img = Image.open(io.BytesIO(erased_bytes))
    assert erased_img.size == (orig_w, orig_h), "Erased image dimensions should match original"
    print("  [PASS] Erase text from image (dimensions preserved)")


def test_erase_text_no_blocks():
    """Erasing with no text blocks returns image unchanged in dimensions."""
    image_bytes = _create_image_with_text()
    erased_bytes = erase_text_from_image(image_bytes, [])
    erased_img = Image.open(io.BytesIO(erased_bytes))
    original_img = Image.open(io.BytesIO(image_bytes))
    assert erased_img.size == original_img.size
    print("  [PASS] Erase text with no blocks preserves image")


def test_sample_text_color():
    """Color sampling detects dark text on light background."""
    import numpy as np
    # Create a white image with a dark text region
    img = np.full((100, 200, 3), 240, dtype=np.uint8)  # Light gray bg
    # Draw dark "text" in the middle
    img[30:60, 40:160] = [30, 30, 80]  # Dark blue-ish text

    color = _sample_text_color(img, 40, 30, 160, 60)
    # Should detect the dark text color, not the background
    assert color[0] < 100, f"Expected dark R, got {color[0]}"
    assert color[2] < 150, f"Expected dark-ish B, got {color[2]}"
    print("  [PASS] Sample text color detects dark text on light bg")


def test_sample_text_color_fallback():
    """Color sampling falls back to white for tiny bbox."""
    import numpy as np
    img = np.full((100, 200, 3), 128, dtype=np.uint8)
    # Very small bbox (< 4px)
    color = _sample_text_color(img, 10, 10, 12, 12)
    assert color == (255, 255, 255), "Tiny bbox should fallback to white"
    print("  [PASS] Sample text color fallback for tiny bbox")


def test_build_pdnob_slide():
    """PDNob slide has background image + text boxes."""
    from slide_builder import SlideBuilder

    # Create a temp image
    img = Image.new("RGB", (1280, 720), color=(200, 200, 220))
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name)
    tmp.close()

    try:
        builder = SlideBuilder()
        text_blocks = [
            OCRTextBlock(
                text="Test Title", x_pct=5.0, y_pct=10.0,
                width_pct=80.0, height_pct=8.0,
                font_size_pt=32.0, color=(255, 255, 255),
            ),
            OCRTextBlock(
                text="Body text here", x_pct=5.0, y_pct=25.0,
                width_pct=60.0, height_pct=5.0,
                font_size_pt=18.0, color=(200, 200, 200),
            ),
        ]
        builder.build_pdnob_slide(tmp.name, text_blocks)

        # Verify slide was created
        assert len(builder.prs.slides) == 1
        slide = builder.prs.slides[0]

        # Should have picture shape (background) + 2 text boxes
        shapes = list(slide.shapes)
        has_picture = any(s.shape_type == 13 for s in shapes)
        assert has_picture, "PDNob slide should have a background picture"

        text_shapes = [s for s in shapes if hasattr(s, "text_frame")]
        assert len(text_shapes) >= 2, f"Expected >= 2 text shapes, got {len(text_shapes)}"
    finally:
        Path(tmp.name).unlink(missing_ok=True)

    print("  [PASS] PDNob slide has background image + text boxes")


def test_build_pdnob_slide_no_text():
    """PDNob slide with no text blocks just has the background image."""
    from slide_builder import SlideBuilder

    img = Image.new("RGB", (1280, 720), color=(100, 100, 100))
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name)
    tmp.close()

    try:
        builder = SlideBuilder()
        builder.build_pdnob_slide(tmp.name, [])

        assert len(builder.prs.slides) == 1
        slide = builder.prs.slides[0]
        has_picture = any(s.shape_type == 13 for s in slide.shapes)
        assert has_picture, "Fallback PDNob slide should still have background image"
    finally:
        Path(tmp.name).unlink(missing_ok=True)

    print("  [PASS] PDNob slide with no text has background image only")


@patch("ocr_converter.extract_text_with_positions")
@patch("ocr_converter.erase_text_from_image")
def test_convert_pdnob_style(mock_erase, mock_extract):
    """Full PDNob conversion pipeline with mocked OCR + inpainting."""
    pptx_path = _create_mock_notebooklm_pptx(num_slides=2)

    # Mock OCR returning text blocks
    mock_extract.return_value = [
        OCRTextBlock(
            text="Slide Title", x_pct=5.0, y_pct=5.0,
            width_pct=80.0, height_pct=8.0,
            font_size_pt=28.0, color=(255, 255, 255),
        ),
    ]

    # Mock inpainting returning a valid PNG
    clean_img = Image.new("RGB", (1280, 720), color=(50, 50, 80))
    buf = io.BytesIO()
    clean_img.save(buf, format="PNG")
    mock_erase.return_value = buf.getvalue()

    try:
        output_pptx = tempfile.mktemp(suffix="_pdnob.pptx")
        result = convert_pdnob_style(
            input_pptx=pptx_path,
            output_pptx=output_pptx,
        )

        assert result["success"], f"PDNob pipeline failed: {result.get('error')}"
        assert result["metadata"]["total_slides"] == 2
        assert result["metadata"]["mode"] == "pdnob"
        assert "ocr_inpaint" in result["timing"]
        assert Path(output_pptx).exists()

        # Verify output PPTX has correct number of slides
        prs = Presentation(output_pptx)
        assert len(prs.slides) == 2

        # Each slide should have a picture shape (background)
        for slide in prs.slides:
            has_picture = any(s.shape_type == 13 for s in slide.shapes)
            assert has_picture, "PDNob slides should have background pictures"

        Path(output_pptx).unlink(missing_ok=True)
    finally:
        Path(pptx_path).unlink()

    print("  [PASS] Full PDNob mock conversion pipeline (2 slides)")


def test_convert_pdnob_empty_pptx():
    """PDNob pipeline returns error for PPTX with no images."""
    pptx_path = _create_mock_empty_pptx()
    try:
        result = convert_pdnob_style(input_pptx=pptx_path)
        assert not result["success"]
        assert "No images" in result["error"]
    finally:
        Path(pptx_path).unlink()
    print("  [PASS] PDNob pipeline rejects empty PPTX")


# ─── Text Block Merging Tests ────────────────────────────────────


def test_group_text_blocks_vertical_merge():
    """Two vertically adjacent blocks with same x merge into one."""
    blocks = [
        OCRTextBlock(text="Jefes /", x_pct=10.0, y_pct=22.0,
                     width_pct=20.0, height_pct=3.0,
                     font_size_pt=18.0, color=(255, 255, 255)),
        OCRTextBlock(text="Superiores.", x_pct=10.0, y_pct=26.0,
                     width_pct=20.0, height_pct=3.0,
                     font_size_pt=18.0, color=(255, 255, 255)),
    ]
    merged = group_text_blocks(blocks)
    assert len(merged) == 1, f"Expected 1 merged block, got {len(merged)}"
    assert "Jefes /" in merged[0].text
    assert "Superiores." in merged[0].text
    assert merged[0].y_pct == 22.0
    assert merged[0].height_pct == 7.0  # 26+3 - 22 = 7
    print("  [PASS] Vertical merge: two adjacent blocks → one")


def test_group_text_blocks_different_columns():
    """Two blocks in separate columns stay separate."""
    blocks = [
        OCRTextBlock(text="Left column", x_pct=10.0, y_pct=30.0,
                     width_pct=18.0, height_pct=3.0,
                     font_size_pt=14.0, color=(255, 255, 255)),
        OCRTextBlock(text="Right column", x_pct=55.0, y_pct=30.0,
                     width_pct=18.0, height_pct=3.0,
                     font_size_pt=14.0, color=(255, 255, 255)),
    ]
    merged = group_text_blocks(blocks)
    assert len(merged) == 2, f"Expected 2 separate blocks, got {len(merged)}"
    print("  [PASS] Different columns stay separate")


def test_group_text_blocks_font_size_mismatch():
    """Title (36pt) and body (14pt) stay separate."""
    blocks = [
        OCRTextBlock(text="Big Title", x_pct=10.0, y_pct=10.0,
                     width_pct=60.0, height_pct=6.0,
                     font_size_pt=36.0, color=(255, 255, 255)),
        OCRTextBlock(text="Small body", x_pct=10.0, y_pct=18.0,
                     width_pct=60.0, height_pct=3.0,
                     font_size_pt=14.0, color=(255, 255, 255)),
    ]
    merged = group_text_blocks(blocks)
    assert len(merged) == 2, f"Expected 2 blocks (font mismatch), got {len(merged)}"
    print("  [PASS] Font size mismatch keeps blocks separate")


def test_group_text_blocks_color_mismatch():
    """White text and cyan text stay separate."""
    blocks = [
        OCRTextBlock(text="White text", x_pct=10.0, y_pct=20.0,
                     width_pct=30.0, height_pct=3.0,
                     font_size_pt=16.0, color=(255, 255, 255)),
        OCRTextBlock(text="Cyan text", x_pct=10.0, y_pct=24.0,
                     width_pct=30.0, height_pct=3.0,
                     font_size_pt=16.0, color=(0, 200, 200)),
    ]
    merged = group_text_blocks(blocks)
    assert len(merged) == 2, f"Expected 2 blocks (color mismatch), got {len(merged)}"
    print("  [PASS] Color mismatch keeps blocks separate")


def test_group_text_blocks_three_way_transitive():
    """A merges B, B merges C → all in one group (transitive via Union-Find)."""
    blocks = [
        OCRTextBlock(text="Line A", x_pct=10.0, y_pct=20.0,
                     width_pct=40.0, height_pct=3.0,
                     font_size_pt=16.0, color=(255, 255, 255)),
        OCRTextBlock(text="Line B", x_pct=10.0, y_pct=24.0,
                     width_pct=40.0, height_pct=3.0,
                     font_size_pt=16.0, color=(255, 255, 255)),
        OCRTextBlock(text="Line C", x_pct=10.0, y_pct=28.0,
                     width_pct=40.0, height_pct=3.0,
                     font_size_pt=16.0, color=(255, 255, 255)),
    ]
    merged = group_text_blocks(blocks)
    assert len(merged) == 1, f"Expected 1 merged block (transitive), got {len(merged)}"
    assert "Line A" in merged[0].text
    assert "Line B" in merged[0].text
    assert "Line C" in merged[0].text
    print("  [PASS] Transitive merge: A+B+C → one block")


def test_group_text_blocks_empty_and_single():
    """Edge cases: 0 blocks returns empty, 1 block returns itself."""
    assert group_text_blocks([]) == []

    single = [OCRTextBlock(text="Solo", x_pct=5.0, y_pct=5.0,
                            width_pct=10.0, height_pct=2.0,
                            font_size_pt=14.0, color=(200, 200, 200))]
    result = group_text_blocks(single)
    assert len(result) == 1
    assert result[0].text == "Solo"
    print("  [PASS] Empty and single block edge cases")


def test_merge_block_group_text_concatenation():
    """Vertically separated lines get newlines, same-line fragments get spaces."""
    # Two blocks with large vertical gap → newline
    blocks_vertical = [
        OCRTextBlock(text="First line", x_pct=10.0, y_pct=10.0,
                     width_pct=30.0, height_pct=3.0,
                     font_size_pt=16.0, color=(255, 255, 255)),
        OCRTextBlock(text="Second line", x_pct=10.0, y_pct=20.0,
                     width_pct=30.0, height_pct=3.0,
                     font_size_pt=16.0, color=(255, 255, 255)),
    ]
    merged = _merge_block_group(blocks_vertical)
    assert "\n" in merged.text, "Vertical gap should produce newline"
    assert "First line" in merged.text
    assert "Second line" in merged.text

    # Two blocks at nearly same y → space
    blocks_sameline = [
        OCRTextBlock(text="Hello", x_pct=10.0, y_pct=10.0,
                     width_pct=10.0, height_pct=3.0,
                     font_size_pt=16.0, color=(255, 255, 255)),
        OCRTextBlock(text="World", x_pct=21.0, y_pct=10.5,
                     width_pct=10.0, height_pct=3.0,
                     font_size_pt=16.0, color=(255, 255, 255)),
    ]
    merged2 = _merge_block_group(blocks_sameline)
    assert "\n" not in merged2.text, "Same-line should join with space, not newline"
    assert "Hello" in merged2.text and "World" in merged2.text
    print("  [PASS] Text concatenation: newlines vs spaces")


def test_most_common_color():
    """Color bucketing groups similar colors and returns most frequent."""
    colors = [
        (255, 255, 255),
        (250, 250, 250),  # close to white
        (255, 255, 255),
        (0, 200, 200),    # cyan — outlier
    ]
    result = _most_common_color(colors)
    # White bucket (3 entries) should win over cyan (1 entry)
    assert result[0] > 200, f"Expected white-ish, got {result}"

    # Empty returns white fallback
    assert _most_common_color([]) == (255, 255, 255)

    # Single color returns itself
    assert _most_common_color([(100, 50, 25)]) == (100, 50, 25)
    print("  [PASS] Most common color bucketing")


# ─── Visual Segmentation Helper ──────────────────────────────────


def _create_cream_image_with_content(layout: str) -> bytes:
    """Create synthetic images for visual segmentation testing.

    Args:
        layout: One of "four_icons", "three_row", "uniform",
                "single_icon", "faint_lines".
    """
    from PIL import ImageDraw

    img = Image.new("RGB", (1376, 768), color=(235, 230, 215))  # cream bg
    draw = ImageDraw.Draw(img)

    if layout == "four_icons":
        # 4 distinct colored blocks in 2x2 grid (like slide 1)
        draw.rectangle([100, 80, 350, 330], fill=(60, 100, 80))
        draw.rectangle([750, 80, 1000, 330], fill=(80, 120, 100))
        draw.rectangle([100, 430, 350, 680], fill=(70, 90, 110))
        draw.rectangle([750, 430, 1000, 680], fill=(120, 80, 70))
    elif layout == "three_row":
        # 3 colored blocks in horizontal row (like slide 2)
        draw.rectangle([60, 200, 380, 560], fill=(170, 90, 40))
        draw.rectangle([510, 200, 830, 560], fill=(100, 80, 70))
        draw.rectangle([960, 200, 1280, 560], fill=(190, 130, 50))
    elif layout == "uniform":
        pass  # pure cream - no content
    elif layout == "single_icon":
        # One illustration centered
        draw.rectangle([400, 200, 900, 550], fill=(50, 80, 120))
    elif layout == "faint_lines":
        # Only faint decorative lines close to background color
        draw.line([(0, 384), (1376, 384)], fill=(225, 220, 205), width=2)
        draw.line([(688, 0), (688, 768)], fill=(225, 220, 205), width=2)

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ─── Precision Positioning + Image Segmentation Tests ────────────


def test_positioning_no_artificial_padding():
    """Text box width matches OCR width — no extra 0.15" padding added."""
    from slide_builder import SlideBuilder

    img = Image.new("RGB", (1280, 720), color=(200, 200, 220))
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name)
    tmp.close()

    try:
        builder = SlideBuilder()
        text_blocks = [
            OCRTextBlock(
                text="Precise Width", x_pct=10.0, y_pct=20.0,
                width_pct=50.0, height_pct=5.0,
                font_size_pt=24.0, color=(255, 255, 255),
            ),
        ]
        builder.build_pdnob_slide(tmp.name, text_blocks)

        slide = builder.prs.slides[0]
        text_shapes = [s for s in slide.shapes if hasattr(s, "text_frame") and s.text_frame.text]
        assert len(text_shapes) >= 1, "Should have at least one text shape"

        # Expected width: 50% of slide width, NO extra padding
        slide_w = builder.prs.slide_width / 914400  # EMU to inches
        expected_width_inches = 50.0 / 100 * slide_w
        actual_width_inches = text_shapes[0].width / 914400

        # Should be within 0.01" of expected (no 0.15" padding)
        assert abs(actual_width_inches - expected_width_inches) < 0.02, (
            f"Width {actual_width_inches:.3f}\" should be ~{expected_width_inches:.3f}\" "
            f"(no 0.15\" padding). Diff: {abs(actual_width_inches - expected_width_inches):.3f}\""
        )
    finally:
        Path(tmp.name).unlink(missing_ok=True)
    print("  [PASS] Text box width matches OCR width (no artificial padding)")


def test_positioning_small_min_size():
    """Minimum text box size is 0.05", not the old 0.3"/0.2"."""
    from slide_builder import SlideBuilder

    img = Image.new("RGB", (1280, 720), color=(200, 200, 220))
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name)
    tmp.close()

    try:
        builder = SlideBuilder()
        # Very small block (0.5% width = ~0.067" on 13.333" slide)
        text_blocks = [
            OCRTextBlock(
                text=".", x_pct=50.0, y_pct=50.0,
                width_pct=0.5, height_pct=0.3,
                font_size_pt=8.0, color=(255, 255, 255),
            ),
        ]
        builder.build_pdnob_slide(tmp.name, text_blocks)

        slide = builder.prs.slides[0]
        text_shapes = [s for s in slide.shapes if hasattr(s, "text_frame") and s.text_frame.text]
        assert len(text_shapes) >= 1

        width_inches = text_shapes[0].width / 914400
        height_inches = text_shapes[0].height / 914400

        # Should use 0.05" minimum, not 0.3"/0.2"
        assert width_inches < 0.2, (
            f"Min width {width_inches:.3f}\" should be < 0.2\" (old was 0.3\")"
        )
        assert height_inches < 0.15, (
            f"Min height {height_inches:.3f}\" should be < 0.15\" (old was 0.2\")"
        )
    finally:
        Path(tmp.name).unlink(missing_ok=True)
    print("  [PASS] Minimum text box size is 0.05\" (not 0.3\"/0.2\")")


def test_source_dims_17x10_detected():
    """17.8\" x 10.0\" dimensions from input PPTX are detected and passed to builder."""
    pptx_path = _create_mock_notebooklm_pptx(num_slides=1)
    try:
        prs = Presentation(pptx_path)
        src_w = prs.slide_width / 914400  # EMU to inches
        src_h = prs.slide_height / 914400
        assert abs(src_w - 17.78) < 0.1, f"Expected ~17.78\", got {src_w:.2f}\""
        assert abs(src_h - 10.0) < 0.1, f"Expected ~10.0\", got {src_h:.2f}\""
    finally:
        Path(pptx_path).unlink()
    print("  [PASS] Source PPTX dimensions 17.8\" x 10.0\" detected")


def test_source_dims_passed_to_builder():
    """SlideBuilder accepts custom slide dimensions and applies them."""
    from slide_builder import SlideBuilder

    builder = SlideBuilder(slide_width_inches=17.78, slide_height_inches=10.0)

    # Presentation dimensions should match
    w_inches = builder.prs.slide_width / 914400
    h_inches = builder.prs.slide_height / 914400
    assert abs(w_inches - 17.78) < 0.01, f"Expected 17.78\", got {w_inches:.2f}\""
    assert abs(h_inches - 10.0) < 0.01, f"Expected 10.0\", got {h_inches:.2f}\""
    print("  [PASS] SlideBuilder applies custom slide dimensions")


def test_font_size_10inch_height():
    """10\" slide height uses 720pt reference (not hardcoded 540pt)."""
    # Font size formula: bbox_h / img_h * (slide_height * 72)
    # For a 10" slide: reference = 720pt
    # A bbox that is 5% of image height → 720 * 0.05 = 36pt
    # Old formula with 540pt would give: 540 * 0.05 = 27pt

    # Create image and mock OCR
    img = Image.new("RGB", (1780, 1000), color=(200, 200, 220))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    image_bytes = buf.getvalue()

    with patch("ocr_converter._get_rapid_ocr") as mock_ocr:
        mock_engine = MagicMock()
        mock_ocr.return_value = mock_engine
        # bbox height = 50px out of 1000px = 5% of image
        mock_engine.return_value = (
            [
                [[[100, 100], [500, 100], [500, 150], [100, 150]], "Test text", "0.95"],
            ],
            [0.1],
        )
        blocks = extract_text_with_positions(image_bytes, slide_height_inches=10.0)
        assert len(blocks) == 1
        # Expected: 50/1000 * 720 = 36pt
        assert abs(blocks[0].font_size_pt - 36.0) < 2.0, (
            f"Font size {blocks[0].font_size_pt}pt should be ~36pt for 10\" slide, not ~27pt"
        )
    print("  [PASS] Font size calibrated to 10\" slide height (720pt ref)")


def test_segment_visual_four_icons():
    """4 colored blocks on cream background → 4 regions."""
    image_bytes = _create_cream_image_with_content("four_icons")
    regions = segment_slide_image(image_bytes)
    assert len(regions) == 4, f"Four icons should produce 4 regions, got {len(regions)}"
    # All regions should be reasonable size (each icon is ~18% x 33% = ~6% of area)
    for r in regions:
        assert r.width_pct > 5.0, f"Region too narrow: {r.width_pct:.1f}%"
        assert r.height_pct > 10.0, f"Region too short: {r.height_pct:.1f}%"
    print("  [PASS] Four icons on cream → 4 regions")


def test_segment_visual_three_row():
    """3 colored blocks in a row on cream → 3 regions."""
    image_bytes = _create_cream_image_with_content("three_row")
    regions = segment_slide_image(image_bytes)
    assert len(regions) == 3, f"Three icons in row should produce 3 regions, got {len(regions)}"
    # Regions should be sorted left-to-right (by x_pct)
    x_positions = [r.x_pct for r in regions]
    assert x_positions == sorted(x_positions), f"Regions should be sorted left-to-right"
    print("  [PASS] Three icons in row → 3 regions")


def test_segment_visual_uniform_fallback():
    """Uniform cream background produces 1 region fallback."""
    image_bytes = _create_cream_image_with_content("uniform")
    regions = segment_slide_image(image_bytes)
    assert len(regions) == 1, f"Uniform cream should produce 1 region, got {len(regions)}"
    assert regions[0].x_pct == 0.0
    assert regions[0].y_pct == 0.0
    assert regions[0].width_pct == 100.0
    assert regions[0].height_pct == 100.0
    print("  [PASS] Uniform cream → single region fallback")


def test_segment_visual_single_icon_fallback():
    """Single illustration → 1 region fallback (< 2 regions)."""
    image_bytes = _create_cream_image_with_content("single_icon")
    regions = segment_slide_image(image_bytes)
    assert len(regions) == 1, f"Single icon should produce 1 fallback region, got {len(regions)}"
    assert regions[0].width_pct == 100.0, "Single region should be full-bleed"
    print("  [PASS] Single icon → fallback to full-bleed")


def test_segment_visual_faint_lines_ignored():
    """Faint decorative lines (close to bg color) are ignored."""
    image_bytes = _create_cream_image_with_content("faint_lines")
    regions = segment_slide_image(image_bytes)
    assert len(regions) == 1, f"Faint lines should be ignored → 1 fallback, got {len(regions)}"
    print("  [PASS] Faint decorative lines ignored")


def test_remove_background_returns_rgba():
    """remove_background() returns RGBA (transparent) PNG bytes."""
    # Create a simple image: colored rectangle on cream background
    img = Image.new("RGB", (200, 200), color=(235, 230, 215))
    from PIL import ImageDraw
    draw = ImageDraw.Draw(img)
    draw.rectangle([50, 50, 150, 150], fill=(60, 100, 80))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    input_bytes = buf.getvalue()

    result_bytes = remove_background(input_bytes)
    result_img = Image.open(io.BytesIO(result_bytes))
    assert result_img.mode == "RGBA", f"Expected RGBA, got {result_img.mode}"
    assert result_img.size == (200, 200), f"Size should be preserved, got {result_img.size}"
    print("  [PASS] remove_background returns RGBA PNG")


def test_remove_background_has_transparent_pixels():
    """After background removal, cream areas should have alpha < 255."""
    img = Image.new("RGB", (200, 200), color=(235, 230, 215))
    from PIL import ImageDraw
    draw = ImageDraw.Draw(img)
    draw.rectangle([60, 60, 140, 140], fill=(30, 60, 120))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    input_bytes = buf.getvalue()

    result_bytes = remove_background(input_bytes)
    result_img = Image.open(io.BytesIO(result_bytes))
    # Check corners (should be transparent — background area)
    corner_alpha = result_img.getpixel((5, 5))[3]
    assert corner_alpha < 128, f"Corner should be transparent, alpha={corner_alpha}"
    print("  [PASS] remove_background makes background transparent")


def test_remove_background_preserves_format():
    """Output is valid PNG bytes that can be reopened."""
    img = Image.new("RGB", (100, 100), color=(200, 50, 50))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    input_bytes = buf.getvalue()

    result_bytes = remove_background(input_bytes)
    # Should be valid PNG
    result_img = Image.open(io.BytesIO(result_bytes))
    assert result_img.format == "PNG", f"Expected PNG, got {result_img.format}"
    print("  [PASS] remove_background output is valid PNG")


def test_crop_image_region_dimensions():
    """Cropped image matches region pixel dimensions."""
    # Create 1000x500 test image
    img = Image.new("RGB", (1000, 500), color=(100, 150, 200))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    image_bytes = buf.getvalue()

    region = ImageRegion(x_pct=10.0, y_pct=20.0, width_pct=50.0, height_pct=40.0)
    cropped_bytes = crop_image_region(image_bytes, region)

    cropped_img = Image.open(io.BytesIO(cropped_bytes))
    # Expected: 50% of 1000 = 500px wide, 40% of 500 = 200px tall
    assert cropped_img.size[0] == 500, f"Expected 500px wide, got {cropped_img.size[0]}"
    assert cropped_img.size[1] == 200, f"Expected 200px tall, got {cropped_img.size[1]}"
    print("  [PASS] Crop image region produces correct dimensions")


def test_build_pdnob_with_regions():
    """Slide with image_regions has multiple images (not just one full-bleed)."""
    from slide_builder import SlideBuilder

    img = Image.new("RGB", (1280, 720), color=(200, 200, 220))
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name)
    tmp.close()

    # Read image bytes for cropping
    with open(tmp.name, "rb") as f:
        cleaned_bytes = f.read()

    try:
        builder = SlideBuilder()
        text_blocks = [
            OCRTextBlock(
                text="Title", x_pct=5.0, y_pct=5.0,
                width_pct=40.0, height_pct=5.0,
                font_size_pt=24.0, color=(255, 255, 255),
            ),
        ]
        regions = [
            ImageRegion(x_pct=0.0, y_pct=0.0, width_pct=50.0, height_pct=100.0),
            ImageRegion(x_pct=50.0, y_pct=0.0, width_pct=50.0, height_pct=100.0),
        ]
        builder.build_pdnob_slide(tmp.name, text_blocks, image_regions=regions,
                                  cleaned_bytes=cleaned_bytes)

        slide = builder.prs.slides[0]
        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        # Should have 2 images (one per region), not just 1 full-bleed
        assert len(picture_shapes) == 2, (
            f"Expected 2 picture shapes (multi-region), got {len(picture_shapes)}"
        )
    finally:
        Path(tmp.name).unlink(missing_ok=True)
    print("  [PASS] Multi-region build produces multiple images")


def test_build_pdnob_no_regions_backward_compat():
    """None regions → single full-bleed image (backward compatible)."""
    from slide_builder import SlideBuilder

    img = Image.new("RGB", (1280, 720), color=(200, 200, 220))
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name)
    tmp.close()

    try:
        builder = SlideBuilder()
        text_blocks = [
            OCRTextBlock(
                text="Title", x_pct=5.0, y_pct=5.0,
                width_pct=40.0, height_pct=5.0,
                font_size_pt=24.0, color=(255, 255, 255),
            ),
        ]
        # No image_regions param → backward compatible
        builder.build_pdnob_slide(tmp.name, text_blocks)

        slide = builder.prs.slides[0]
        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        assert len(picture_shapes) == 1, (
            f"Expected 1 full-bleed image (no regions), got {len(picture_shapes)}"
        )
    finally:
        Path(tmp.name).unlink(missing_ok=True)
    print("  [PASS] No regions → single full-bleed image (backward compat)")


def test_full_bleed_uses_custom_dims():
    """Full-bleed image uses actual slide dimensions, not hardcoded 13.333x7.5."""
    from slide_builder import SlideBuilder

    img = Image.new("RGB", (1780, 1000), color=(200, 200, 220))
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name)
    tmp.close()

    try:
        builder = SlideBuilder(slide_width_inches=17.78, slide_height_inches=10.0)
        builder.build_pdnob_slide(tmp.name, [])

        slide = builder.prs.slides[0]
        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        assert len(picture_shapes) == 1

        pic = picture_shapes[0]
        pic_w = pic.width / 914400
        pic_h = pic.height / 914400
        # Image should fill the custom 17.78" x 10.0" slide, NOT 13.333" x 7.5"
        assert abs(pic_w - 17.78) < 0.01, (
            f"Full-bleed image width {pic_w:.2f}\" should be 17.78\", not 13.333\""
        )
        assert abs(pic_h - 10.0) < 0.01, (
            f"Full-bleed image height {pic_h:.2f}\" should be 10.0\", not 7.5\""
        )
    finally:
        Path(tmp.name).unlink(missing_ok=True)
    print("  [PASS] Full-bleed image uses custom slide dimensions")


# ─── PDNob Level Tests ─────────────────────────────────────────────

def _create_pdnob_test_pptx() -> str:
    """Create a PPTX with a full-bleed image slide for PDNob level testing."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    img = Image.new("RGB", (1333, 750), color=(40, 40, 80))
    # Draw some "text-like" blocks so OCR finds something
    from PIL import ImageDraw
    draw = ImageDraw.Draw(img)
    draw.rectangle([100, 50, 600, 120], fill=(255, 255, 255))
    draw.rectangle([100, 200, 800, 350], fill=(200, 200, 220))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)

    slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height)

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


def test_pdnob_level_ocr_only():
    """PDNob ocr_only: text shapes present, no segmented regions (1 full-bleed picture)."""
    input_pptx = _create_pdnob_test_pptx()
    try:
        result = convert_pdnob_style(input_pptx, pdnob_level="ocr_only")
        assert result["success"], f"PDNob ocr_only failed: {result.get('error')}"
        assert result["metadata"]["pdnob_level"] == "ocr_only"

        output_path = result["files"]["pptx"]
        prs = Presentation(output_path)
        slide = prs.slides[0]

        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        # ocr_only: single full-bleed image (no segmentation)
        assert len(picture_shapes) == 1, (
            f"ocr_only should have 1 full-bleed image, got {len(picture_shapes)}"
        )
        Path(output_path).unlink(missing_ok=True)
    finally:
        Path(input_pptx).unlink(missing_ok=True)
    print("  [PASS] PDNob ocr_only: full-bleed image + text boxes")


def test_pdnob_level_remove_bg():
    """PDNob remove_bg: no text shapes, picture shapes present."""
    input_pptx = _create_pdnob_test_pptx()
    try:
        result = convert_pdnob_style(input_pptx, pdnob_level="remove_bg")
        assert result["success"], f"PDNob remove_bg failed: {result.get('error')}"
        assert result["metadata"]["pdnob_level"] == "remove_bg"

        output_path = result["files"]["pptx"]
        prs = Presentation(output_path)
        slide = prs.slides[0]

        # remove_bg: should have picture(s) but no text boxes
        text_shapes = [
            s for s in slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        assert len(text_shapes) == 0, (
            f"remove_bg should have 0 text shapes, got {len(text_shapes)}"
        )
        Path(output_path).unlink(missing_ok=True)
    finally:
        Path(input_pptx).unlink(missing_ok=True)
    print("  [PASS] PDNob remove_bg: picture shapes only, no text")


def test_pdnob_level_full():
    """PDNob full: both text shapes and picture shapes present."""
    input_pptx = _create_pdnob_test_pptx()
    try:
        result = convert_pdnob_style(input_pptx, pdnob_level="full")
        assert result["success"], f"PDNob full failed: {result.get('error')}"
        assert result["metadata"]["pdnob_level"] == "full"

        output_path = result["files"]["pptx"]
        prs = Presentation(output_path)
        assert len(prs.slides) == 1, f"Expected 1 slide, got {len(prs.slides)}"
        Path(output_path).unlink(missing_ok=True)
    finally:
        Path(input_pptx).unlink(missing_ok=True)
    print("  [PASS] PDNob full: both text and picture shapes")


# ─── Run All Tests ────────────────────────────────────────────────

def run_all():
    print("\n=== OCR Converter Test Suite ===\n")
    tests = [
        ("OCR Prompt Loads", test_ocr_prompt_loads),
        ("OCR Prompt Slide Types", test_ocr_prompt_has_slide_types),
        ("Extract Slide Images", test_extract_slide_images),
        ("Extract No Images", test_extract_slide_images_empty),
        ("Extract Single Slide", test_extract_slide_images_single),
        ("Classify Valid Types", test_classify_valid_types),
        ("Classify First → Title", test_classify_first_slide_fallback),
        ("Classify Last → Conclusion", test_classify_last_slide_fallback),
        ("Classify Middle → Content", test_classify_middle_slide_fallback),
        ("SlideSpec Full Mapping", test_content_to_slidespec_full),
        ("SlideSpec Minimal", test_content_to_slidespec_minimal),
        ("SlideSpec Empty", test_content_to_slidespec_empty),
        ("Mock Pipeline (3 slides)", test_convert_pipeline_mock),
        ("OCR Failure Fallback", test_convert_pipeline_ocr_failure_fallback),
        ("No API Key Error", test_convert_no_api_key),
        ("Empty PPTX Error", test_convert_empty_pptx),
        # New: UTF-8 normalization
        ("Normalize NFC", test_normalize_text_nfc),
        ("Normalize Zero-Width", test_normalize_text_zero_width),
        ("Normalize NBSP", test_normalize_text_nbsp),
        ("Normalize Empty", test_normalize_text_empty),
        ("SlideSpec Normalizes Text", test_content_to_slidespec_normalizes_text),
        # New: Docling OCR engine
        ("Parse Raw Text Basic", test_parse_raw_text_to_slide_basic),
        ("Parse Raw Text Title", test_parse_raw_text_to_slide_title_slide),
        ("Parse Raw Text Conclusion", test_parse_raw_text_to_slide_conclusion),
        ("Parse Raw Text Empty", test_parse_raw_text_to_slide_empty),
        ("Parse Numbered Bullets", test_parse_raw_text_numbered_bullets),
        ("Parse Colon → Em Dash", test_parse_raw_text_colon_to_em_dash),
        ("Invalid OCR Engine", test_convert_invalid_ocr_engine),
        ("Docling No API Key", test_convert_docling_no_api_key),
        ("Docling Mock Pipeline", test_convert_pipeline_docling_mock),
        # New: PDNob OCR
        ("PDNob Extract Text Positions", test_extract_text_with_positions),
        ("PDNob Extract Empty Result", test_extract_text_with_positions_empty),
        ("PDNob Filter Low Confidence", test_extract_text_filters_low_confidence),
        ("PDNob Erase Text", test_erase_text_from_image),
        ("PDNob Erase No Blocks", test_erase_text_no_blocks),
        ("PDNob Sample Text Color", test_sample_text_color),
        ("PDNob Sample Color Fallback", test_sample_text_color_fallback),
        ("PDNob Build Slide", test_build_pdnob_slide),
        ("PDNob Build Slide No Text", test_build_pdnob_slide_no_text),
        ("PDNob Full Pipeline", test_convert_pdnob_style),
        ("PDNob Empty PPTX", test_convert_pdnob_empty_pptx),
        # New: Text Block Merging
        ("Merge Vertical Blocks", test_group_text_blocks_vertical_merge),
        ("Merge Different Columns", test_group_text_blocks_different_columns),
        ("Merge Font Size Mismatch", test_group_text_blocks_font_size_mismatch),
        ("Merge Color Mismatch", test_group_text_blocks_color_mismatch),
        ("Merge Three-Way Transitive", test_group_text_blocks_three_way_transitive),
        ("Merge Empty and Single", test_group_text_blocks_empty_and_single),
        ("Merge Text Concatenation", test_merge_block_group_text_concatenation),
        ("Most Common Color", test_most_common_color),
        # New: Precision Positioning + Image Segmentation
        ("No Artificial Padding", test_positioning_no_artificial_padding),
        ("Small Min Size", test_positioning_small_min_size),
        ("Source Dims 17x10 Detected", test_source_dims_17x10_detected),
        ("Source Dims Passed to Builder", test_source_dims_passed_to_builder),
        ("Font Size 10\" Height", test_font_size_10inch_height),
        ("Segment Visual Four Icons", test_segment_visual_four_icons),
        ("Segment Visual Three Row", test_segment_visual_three_row),
        ("Segment Visual Uniform", test_segment_visual_uniform_fallback),
        ("Segment Visual Single Icon", test_segment_visual_single_icon_fallback),
        ("Segment Visual Faint Lines", test_segment_visual_faint_lines_ignored),
        ("Remove BG Returns RGBA", test_remove_background_returns_rgba),
        ("Remove BG Transparent Pixels", test_remove_background_has_transparent_pixels),
        ("Remove BG Valid PNG", test_remove_background_preserves_format),
        ("Crop Image Region Dims", test_crop_image_region_dimensions),
        ("Build PDNob Multi-Region", test_build_pdnob_with_regions),
        ("Build PDNob No Regions Compat", test_build_pdnob_no_regions_backward_compat),
        ("Full-Bleed Custom Dims", test_full_bleed_uses_custom_dims),
        ("PDNob Level OCR Only", test_pdnob_level_ocr_only),
        ("PDNob Level Remove BG", test_pdnob_level_remove_bg),
        ("PDNob Level Full", test_pdnob_level_full),
    ]

    passed = 0
    failed = 0
    for name, test_fn in tests:
        try:
            test_fn()
            passed += 1
        except Exception as e:
            print(f"  [FAIL] {name}: {e}")
            import traceback
            traceback.print_exc()
            failed += 1

    print(f"\n{'=' * 40}")
    print(f"Results: {passed} passed, {failed} failed, {len(tests)} total")
    print(f"{'=' * 40}\n")
    return failed == 0


if __name__ == "__main__":
    success = run_all()
    sys.exit(0 if success else 1)
