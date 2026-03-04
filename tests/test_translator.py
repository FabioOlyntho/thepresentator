"""
Test Suite — Translator + Multi-Format Content Extraction.

Covers:
- TXT, Markdown, PPTX content extraction (Session 2 additions)
- PPTX translation via Gemini (mocked)
- SlideSpec translation via Gemini (mocked)
- CLI integration

Run: python -m pytest tests/test_translator.py -v
"""

import io
import json
import os
import sys
import tempfile
from copy import deepcopy
from pathlib import Path
from unittest.mock import patch, MagicMock

# Add scripts to path
sys.path.insert(0, str(Path(__file__).parent.parent / "scripts"))

from pptx import Presentation
from pptx.util import Inches

from content_extractor import (
    extract_content,
    extract_txt,
    extract_markdown,
    extract_pptx_text,
    detect_language,
    clean_text,
)
from gemini_client import SlideSpec, PresentationSpec
from translator import (
    translate_pptx,
    translate_specs,
    _call_gemini_translate,
    LANGUAGE_LABELS,
)


# ─── Helper: Create temp files ──────────────────────────────────

def _create_temp_txt(content: str = "This is a test document.\nWith multiple lines.\nFor extraction testing.") -> str:
    tmp = tempfile.NamedTemporaryFile(suffix=".txt", mode="w", encoding="utf-8", delete=False)
    tmp.write(content)
    tmp.close()
    return tmp.name


def _create_temp_markdown() -> str:
    md_content = """# Main Heading

## Section One

This is a paragraph with **bold** and *italic* text.

- First bullet point
- Second bullet point
- Third bullet point

### Subsection

A [link to example](https://example.com) and an image:
![Alt text](image.png)

> This is a blockquote

```python
code_block = "should be stripped"
```

---

## Section Two

Final paragraph with __emphasis__.
"""
    tmp = tempfile.NamedTemporaryFile(suffix=".md", mode="w", encoding="utf-8", delete=False)
    tmp.write(md_content)
    tmp.close()
    return tmp.name


def _create_temp_pptx(num_slides: int = 3) -> str:
    """Create a PPTX with text content for extraction testing."""
    prs = Presentation()
    for i in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.text = f"Slide {i + 1} Title"
        p = tf.add_paragraph()
        p.text = f"Content for slide {i + 1} goes here."
        p2 = tf.add_paragraph()
        p2.text = f"Additional text on slide {i + 1}."

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_translatable_pptx() -> str:
    """Create a PPTX with Spanish text for translation testing."""
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.text = "Siete Tácticas de Influencia"
    p = tf.add_paragraph()
    p.text = "Marco práctico para directivos"

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    tf2 = txBox2.text_frame
    tf2.text = "Razón — datos y lógica"
    p2 = tf2.add_paragraph()
    p2.text = "Coalición — los aliados multiplican"

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


# ─── TXT Extraction Tests ───────────────────────────────────────

def test_extract_txt_basic():
    """Extract text from a plain TXT file."""
    path = _create_temp_txt("Hello world.\nSecond line.\nThird line.")
    try:
        text, pages, meta = extract_txt(path)
        assert "Hello world." in text
        assert "Second line." in text
        assert pages >= 1
    finally:
        Path(path).unlink()
    print("  [PASS] TXT basic extraction")


def test_extract_txt_via_extract_content():
    """extract_content() dispatches .txt files correctly."""
    path = _create_temp_txt("Test content for extraction.\nLine two.")
    try:
        result = extract_content(path)
        assert result.file_type == "TXT"
        assert "Test content" in result.text
        assert result.word_count > 0
        assert result.page_count >= 1
    finally:
        Path(path).unlink()
    print("  [PASS] TXT via extract_content()")


def test_extract_txt_language_detection():
    """TXT extraction detects language correctly."""
    es_path = _create_temp_txt("Esta es una presentación sobre las tácticas de influencia para los directivos de la empresa.")
    en_path = _create_temp_txt("This is a presentation about influence tactics for the managers of the company.")
    try:
        es_result = extract_content(es_path)
        en_result = extract_content(en_path)
        assert es_result.language == "ES"
        assert en_result.language == "EN"
    finally:
        Path(es_path).unlink()
        Path(en_path).unlink()
    print("  [PASS] TXT language detection")


# ─── Markdown Extraction Tests ───────────────────────────────────

def test_extract_markdown_basic():
    """Extract and clean Markdown content."""
    path = _create_temp_markdown()
    try:
        text, pages, meta = extract_markdown(path)
        # Headers stripped of # markers
        assert "Main Heading" in text
        assert "Section One" in text
        # Bold markers stripped
        assert "**bold**" not in text
        assert "bold" in text
        # Links converted to text
        assert "link to example" in text
        assert "https://example.com" not in text
        # Image alt text preserved, URL removed
        assert "image.png" not in text
        # Code fences removed
        assert "```" not in text
        # Blockquote marker removed
        assert meta.get("source_format") == "markdown"
    finally:
        Path(path).unlink()
    print("  [PASS] Markdown extraction strips syntax")


def test_extract_markdown_via_extract_content():
    """extract_content() dispatches .md files correctly."""
    path = _create_temp_markdown()
    try:
        result = extract_content(path)
        assert result.file_type == "MD"
        assert result.word_count > 0
        assert "Main Heading" in result.text
    finally:
        Path(path).unlink()
    print("  [PASS] Markdown via extract_content()")


# ─── PPTX Text Extraction Tests ─────────────────────────────────

def test_extract_pptx_text_basic():
    """Extract text from PPTX shapes."""
    path = _create_temp_pptx(num_slides=3)
    try:
        text, slide_count, meta = extract_pptx_text(path)
        assert slide_count == 3
        assert "Slide 1 Title" in text
        assert "Slide 2 Title" in text
        assert "Slide 3 Title" in text
        assert "Content for slide 1" in text
        assert meta.get("source_format") == "pptx"
    finally:
        Path(path).unlink()
    print("  [PASS] PPTX text extraction (3 slides)")


def test_extract_pptx_via_extract_content():
    """extract_content() dispatches .pptx files correctly."""
    path = _create_temp_pptx(num_slides=2)
    try:
        result = extract_content(path)
        assert result.file_type == "PPTX"
        assert result.page_count == 2
        assert "Slide 1 Title" in result.text
    finally:
        Path(path).unlink()
    print("  [PASS] PPTX via extract_content()")


def test_extract_pptx_single_slide():
    """PPTX with 1 slide extracts correctly."""
    path = _create_temp_pptx(num_slides=1)
    try:
        text, count, meta = extract_pptx_text(path)
        assert count == 1
        assert "Slide 1 Title" in text
    finally:
        Path(path).unlink()
    print("  [PASS] Single-slide PPTX extraction")


# ─── Unsupported File Type Test ──────────────────────────────────

def test_unsupported_file_type():
    """extract_content() raises ValueError for unsupported extensions."""
    tmp = tempfile.NamedTemporaryFile(suffix=".xyz", delete=False)
    tmp.close()
    try:
        raised = False
        try:
            extract_content(tmp.name)
        except ValueError as e:
            raised = True
            assert ".xyz" in str(e)
            assert ".pdf" in str(e)
        assert raised, "Expected ValueError for .xyz file"
    finally:
        Path(tmp.name).unlink()
    print("  [PASS] Unsupported file type raises ValueError")


# ─── Translation Tests (Mocked Gemini) ──────────────────────────

def _mock_gemini_translate_response(translated_texts: list[str]):
    """Create a mock Gemini API response for translation."""
    mock_resp = MagicMock()
    mock_resp.status_code = 200
    mock_resp.json.return_value = {
        "candidates": [{
            "content": {
                "parts": [{
                    "text": json.dumps(translated_texts)
                }]
            }
        }]
    }
    return mock_resp


@patch("translator.requests.post")
def test_call_gemini_translate_basic(mock_post):
    """Batch translate strings via Gemini API."""
    texts = ["Hola mundo", "Buenas tardes"]
    expected = ["Hello world", "Good afternoon"]

    mock_post.return_value = _mock_gemini_translate_response(expected)

    result = _call_gemini_translate(
        texts=texts,
        source_lang="ES",
        target_lang="EN",
        api_key="test-key",
    )
    assert result == expected
    assert mock_post.call_count == 1
    print("  [PASS] Gemini translate basic (2 strings)")


@patch("translator.requests.post")
def test_call_gemini_translate_preserves_order(mock_post):
    """Translation preserves order of input texts."""
    texts = ["Uno", "Dos", "Tres", "Cuatro"]
    expected = ["One", "Two", "Three", "Four"]

    mock_post.return_value = _mock_gemini_translate_response(expected)

    result = _call_gemini_translate(
        texts=texts,
        source_lang="ES",
        target_lang="EN",
        api_key="test-key",
    )
    assert result == expected
    print("  [PASS] Translation preserves order")


@patch("translator.requests.post")
def test_call_gemini_translate_count_mismatch(mock_post):
    """Translation handles count mismatch by padding with originals."""
    texts = ["Uno", "Dos", "Tres"]
    # API returns only 2 translations
    mock_post.return_value = _mock_gemini_translate_response(["One", "Two"])

    result = _call_gemini_translate(
        texts=texts,
        source_lang="ES",
        target_lang="EN",
        api_key="test-key",
    )
    assert len(result) == 3
    assert result[0] == "One"
    assert result[1] == "Two"
    assert result[2] == "Tres"  # Padded with original
    print("  [PASS] Translation handles count mismatch")


@patch("translator.requests.post")
def test_translate_pptx_mock(mock_post):
    """Full PPTX translation pipeline with mocked Gemini."""
    pptx_path = _create_translatable_pptx()

    # Mock returns English translations
    mock_post.return_value = _mock_gemini_translate_response([
        "Seven Influence Tactics",
        "A practical framework for managers",
        "Reason — data and logic",
        "Coalition — allies multiply",
    ])

    try:
        output_path = tempfile.mktemp(suffix="_en.pptx")
        result = translate_pptx(
            input_pptx=pptx_path,
            output_pptx=output_path,
            source_lang="ES",
            target_lang="EN",
            api_key="test-key",
        )

        assert result["success"], f"Translation failed: {result.get('error')}"
        assert result["files"]["pptx"] == output_path
        assert result["metadata"]["source_lang"] == "ES"
        assert result["metadata"]["target_lang"] == "EN"
        assert result["metadata"]["total_runs"] == 4
        assert result["metadata"]["slide_count"] == 2

        # Verify output PPTX contains translated text
        prs = Presentation(output_path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                all_text.append(run.text)

        assert "Seven Influence Tactics" in all_text
        assert "Coalition — allies multiply" in all_text

        Path(output_path).unlink(missing_ok=True)
    finally:
        Path(pptx_path).unlink()

    print("  [PASS] Full PPTX translation pipeline (mocked)")


def test_translate_pptx_no_api_key():
    """Translation fails gracefully without API key."""
    old_key = os.environ.pop("GEMINI_API_KEY", None)
    try:
        result = translate_pptx(
            input_pptx="fake.pptx",
            source_lang="ES",
            target_lang="EN",
            api_key=None,
        )
        assert not result["success"]
        assert "API key" in result["error"]
    finally:
        if old_key:
            os.environ["GEMINI_API_KEY"] = old_key
    print("  [PASS] Translation requires API key")


def test_translate_pptx_file_not_found():
    """Translation fails gracefully for missing file."""
    result = translate_pptx(
        input_pptx="nonexistent.pptx",
        source_lang="ES",
        target_lang="EN",
        api_key="test-key",
    )
    assert not result["success"]
    assert "not found" in result["error"]
    print("  [PASS] Translation file not found error")


def test_translate_pptx_no_text():
    """Translation handles PPTX with no text content."""
    # Create PPTX with only a picture, no text
    from PIL import Image
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = Image.new("RGB", (100, 100), color=(255, 0, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(1), Inches(1), Inches(2), Inches(2))

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()

    try:
        result = translate_pptx(
            input_pptx=tmp.name,
            source_lang="ES",
            target_lang="EN",
            api_key="test-key",
        )
        assert not result["success"]
        assert "No text" in result["error"]
    finally:
        Path(tmp.name).unlink()
    print("  [PASS] Translation handles no-text PPTX")


# ─── SlideSpec Translation Tests ─────────────────────────────────

@patch("translator.requests.post")
def test_translate_specs_basic(mock_post):
    """Translate PresentationSpec content."""
    specs = PresentationSpec(
        title="Tácticas de Influencia",
        subtitle="Marco directivo",
        language="ES",
        source_document="test.pdf",
        themes=["influence"],
        slides=[
            SlideSpec(
                number=1,
                type="title",
                title="Siete Tácticas",
                subtitle="Para directivos",
                body="",
                speaker_notes="Bienvenidos a la presentación.",
            ),
            SlideSpec(
                number=2,
                type="content",
                title="Razón",
                body="Los datos impulsan decisiones",
                bullet_points=["Datos — son fundamentales", "Lógica — guía la acción"],
                speaker_notes="Esta diapositiva explica la razón.",
            ),
        ],
    )

    # Collect all non-empty translatable texts
    expected_count = 0
    for t in [specs.title, specs.subtitle]:
        if t.strip():
            expected_count += 1
    for s in specs.slides:
        for t in [s.title, s.subtitle, s.body, s.speaker_notes, s.left_header, s.right_header]:
            if t.strip():
                expected_count += 1
        expected_count += len(s.bullet_points) + len(s.left_column) + len(s.right_column) + len(s.checkbox_items)

    # Create mock responses that match the expected non-empty count
    translations = [
        "Influence Tactics",      # specs.title
        "Management Framework",   # specs.subtitle
        "Seven Tactics",          # slide 1 title
        "For managers",           # slide 1 subtitle
        "Welcome to the presentation.",  # slide 1 notes
        "Reason",                 # slide 2 title
        "Data drives decisions",  # slide 2 body
        "This slide explains reason.",  # slide 2 notes
        "Data — are fundamental", # slide 2 bullet 1
        "Logic — guides action",  # slide 2 bullet 2
    ]

    mock_post.return_value = _mock_gemini_translate_response(translations)

    result = translate_specs(specs, target_lang="EN", api_key="test-key")

    assert result.title == "Influence Tactics"
    assert result.subtitle == "Management Framework"
    assert result.language == "EN"
    assert result.slides[0].title == "Seven Tactics"
    assert result.slides[0].subtitle == "For managers"
    assert result.slides[1].title == "Reason"
    assert result.slides[1].body == "Data drives decisions"
    assert len(result.slides[1].bullet_points) == 2
    assert result.slides[1].bullet_points[0] == "Data — are fundamental"

    # Original should be unchanged
    assert specs.title == "Tácticas de Influencia"
    assert specs.slides[0].title == "Siete Tácticas"

    print("  [PASS] SlideSpec translation (2 slides)")


def test_translate_specs_no_api_key():
    """translate_specs raises ValueError without API key."""
    old_key = os.environ.pop("GEMINI_API_KEY", None)
    specs = PresentationSpec(
        title="Test", subtitle="", language="ES",
        source_document="", themes=[], slides=[],
    )
    try:
        raised = False
        try:
            translate_specs(specs, target_lang="EN", api_key=None)
        except ValueError as e:
            raised = True
            assert "API key" in str(e)
        assert raised, "Expected ValueError"
    finally:
        if old_key:
            os.environ["GEMINI_API_KEY"] = old_key
    print("  [PASS] translate_specs requires API key")


# ─── Language Labels Test ────────────────────────────────────────

def test_language_labels():
    """All expected language codes have labels."""
    for code in ["ES", "EN", "FR", "DE", "PT", "IT"]:
        assert code in LANGUAGE_LABELS
        assert len(LANGUAGE_LABELS[code]) > 2
    print("  [PASS] Language labels for 6 languages")


# ─── Run All Tests ───────────────────────────────────────────────

def run_all():
    print("\n=== Multi-Format + Translator Test Suite ===\n")
    tests = [
        # TXT extraction
        ("TXT Basic Extraction", test_extract_txt_basic),
        ("TXT via extract_content", test_extract_txt_via_extract_content),
        ("TXT Language Detection", test_extract_txt_language_detection),
        # Markdown extraction
        ("Markdown Extraction", test_extract_markdown_basic),
        ("Markdown via extract_content", test_extract_markdown_via_extract_content),
        # PPTX extraction
        ("PPTX Text Extraction", test_extract_pptx_text_basic),
        ("PPTX via extract_content", test_extract_pptx_via_extract_content),
        ("PPTX Single Slide", test_extract_pptx_single_slide),
        # Unsupported
        ("Unsupported File Type", test_unsupported_file_type),
        # Gemini translate
        ("Gemini Translate Basic", test_call_gemini_translate_basic),
        ("Translate Preserves Order", test_call_gemini_translate_preserves_order),
        ("Translate Count Mismatch", test_call_gemini_translate_count_mismatch),
        # PPTX translation
        ("PPTX Translation (mock)", test_translate_pptx_mock),
        ("Translation No API Key", test_translate_pptx_no_api_key),
        ("Translation File Not Found", test_translate_pptx_file_not_found),
        ("Translation No Text", test_translate_pptx_no_text),
        # SlideSpec translation
        ("SlideSpec Translation", test_translate_specs_basic),
        ("SlideSpec No API Key", test_translate_specs_no_api_key),
        # Language labels
        ("Language Labels", test_language_labels),
    ]

    passed = 0
    failed = 0
    for name, test_fn in tests:
        try:
            test_fn()
            passed += 1
        except Exception as e:
            print(f"  [FAIL] {name}: {e}")
            import traceback
            traceback.print_exc()
            failed += 1

    print(f"\n{'=' * 40}")
    print(f"Results: {passed} passed, {failed} failed, {len(tests)} total")
    print(f"{'=' * 40}\n")
    return failed == 0


if __name__ == "__main__":
    success = run_all()
    sys.exit(0 if success else 1)
