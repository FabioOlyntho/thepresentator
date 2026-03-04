"""
Content Extractor — Extracts text from PDF and DOCX files.

Supports:
- PDF extraction via pdfplumber (primary) with PyPDF2 fallback
- Image-based PDF OCR via pypdfium2 + Gemini Vision API
- DOCX extraction via python-docx
- Metadata extraction (page count, word count, language detection)
- Text cleaning and normalization
"""

import base64
import io
import json
import os
import re
import sys
import logging
import unicodedata
from pathlib import Path
from dataclasses import dataclass, field

import pdfplumber
import PyPDF2
import pypdfium2 as pdfium
import requests
from docx import Document
from PIL import Image

logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")
logger = logging.getLogger(__name__)


@dataclass
class ExtractedContent:
    """Result of content extraction."""
    text: str
    filename: str
    file_type: str
    page_count: int
    word_count: int
    char_count: int
    language: str  # "ES" or "EN" (heuristic)
    metadata: dict = field(default_factory=dict)

    def preview(self, chars: int = 500) -> str:
        return self.text[:chars] + ("..." if len(self.text) > chars else "")

    def summary(self) -> str:
        return (
            f"File: {self.filename}\n"
            f"Type: {self.file_type} | Pages: {self.page_count}\n"
            f"Words: {self.word_count:,} | Chars: {self.char_count:,}\n"
            f"Language: {self.language}\n"
        )


def detect_language(text: str) -> str:
    """Simple heuristic language detection based on common words."""
    spanish_markers = {
        "de", "la", "el", "en", "los", "las", "del", "que", "por",
        "una", "con", "para", "como", "más", "pero", "ser", "también",
        "puede", "entre", "sobre", "desde", "este", "esta", "estos",
        "cuando", "donde", "porque", "según", "hacia", "mediante",
    }
    english_markers = {
        "the", "and", "of", "to", "in", "is", "that", "for", "it",
        "with", "as", "was", "are", "by", "be", "this", "which",
        "from", "or", "an", "but", "not", "they", "have", "has",
        "their", "been", "would", "could", "should", "about",
    }

    words = set(re.findall(r"\b[a-záéíóúñü]+\b", text.lower()))
    es_score = len(words & spanish_markers)
    en_score = len(words & english_markers)

    if es_score > en_score * 1.2:
        return "ES"
    elif en_score > es_score * 1.2:
        return "EN"
    return "ES" if es_score >= en_score else "EN"


def clean_text(text: str) -> str:
    """Clean extracted text: normalize Unicode, whitespace, remove artifacts."""
    # Unicode NFC normalization: compose decomposed accents (e.g. e + ´ → é)
    text = unicodedata.normalize("NFC", text)
    # Replace non-breaking spaces and zero-width characters
    text = text.replace("\u00a0", " ")
    text = text.replace("\u200b", "").replace("\u200c", "").replace("\u200d", "")
    text = text.replace("\ufeff", "")  # BOM
    # Normalize line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Remove excessive blank lines (keep max 2)
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    # Remove page number artifacts (standalone numbers on a line)
    text = re.sub(r"\n\s*\d{1,3}\s*\n", "\n", text)
    # Normalize whitespace within lines
    text = re.sub(r"[ \t]+", " ", text)
    # Strip leading/trailing whitespace from each line
    lines = [line.strip() for line in text.split("\n")]
    text = "\n".join(lines)
    return text.strip()


def extract_pdf_pdfplumber(filepath: str) -> tuple[str, int, dict]:
    """Extract text from PDF using pdfplumber (better for complex layouts)."""
    pages_text = []
    metadata = {}

    with pdfplumber.open(filepath) as pdf:
        page_count = len(pdf.pages)
        metadata["pdf_metadata"] = pdf.metadata or {}

        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages_text.append(text)

    full_text = "\n\n".join(pages_text)
    return full_text, page_count, metadata


def extract_pdf_pypdf2(filepath: str) -> tuple[str, int, dict]:
    """Fallback PDF extraction using PyPDF2."""
    pages_text = []
    metadata = {}

    with open(filepath, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        page_count = len(reader.pages)
        if reader.metadata:
            metadata["pdf_metadata"] = {
                k: str(v) for k, v in reader.metadata.items() if v
            }

        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages_text.append(text)

    full_text = "\n\n".join(pages_text)
    return full_text, page_count, metadata


def render_pdf_pages_to_images(filepath: str, dpi: int = 200) -> list[bytes]:
    """Render PDF pages to PNG images using pypdfium2."""
    pdf = pdfium.PdfDocument(filepath)
    images_bytes = []
    scale = dpi / 72  # 72 is the default PDF DPI

    for i in range(len(pdf)):
        page = pdf[i]
        bitmap = page.render(scale=scale)
        pil_image = bitmap.to_pil()
        buf = io.BytesIO()
        pil_image.save(buf, format="PNG", optimize=True)
        images_bytes.append(buf.getvalue())

    pdf.close()
    return images_bytes


def ocr_with_gemini(image_bytes_list: list[bytes], api_key: str) -> str:
    """
    OCR scanned PDF pages using Gemini Vision API.

    Sends page images to Gemini and asks it to extract text.
    """
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"

    all_text = []

    # Process pages in batches of 4 to stay within token limits
    batch_size = 4
    for batch_start in range(0, len(image_bytes_list), batch_size):
        batch = image_bytes_list[batch_start:batch_start + batch_size]
        batch_end = batch_start + len(batch)

        parts = []
        parts.append({
            "text": (
                f"Extract ALL text from these scanned document pages ({batch_start + 1}-{batch_end}). "
                "Preserve the original language, formatting structure (headings, paragraphs, lists), "
                "and any emphasis (bold, italic). Return ONLY the extracted text, no commentary. "
                "Separate pages with '---PAGE BREAK---'."
            )
        })

        for img_bytes in batch:
            b64 = base64.b64encode(img_bytes).decode("utf-8")
            parts.append({
                "inlineData": {
                    "mimeType": "image/png",
                    "data": b64,
                }
            })

        payload = {
            "contents": [{"parts": parts}],
            "generationConfig": {
                "temperature": 0.1,
                "maxOutputTokens": 8192,
            }
        }

        logger.info("OCR via Gemini: pages %d-%d of %d", batch_start + 1, batch_end, len(image_bytes_list))
        resp = requests.post(url, json=payload, timeout=120)
        resp.raise_for_status()
        data = resp.json()

        text = data["candidates"][0]["content"]["parts"][0]["text"]
        all_text.append(text)

    return "\n\n".join(all_text)


def extract_pdf(filepath: str, gemini_api_key: str | None = None) -> tuple[str, int, dict]:
    """Extract text from PDF, trying pdfplumber first, then PyPDF2, then OCR."""
    page_count = 0

    # Try text extraction first
    try:
        text, pages, meta = extract_pdf_pdfplumber(filepath)
        page_count = pages
        if len(text.strip()) > 100:
            logger.info("PDF extracted via pdfplumber (%d pages)", pages)
            return text, pages, meta
        logger.warning("pdfplumber returned minimal text, trying PyPDF2")
    except Exception as e:
        logger.warning("pdfplumber failed: %s — trying PyPDF2", e)

    try:
        text, pages, meta = extract_pdf_pypdf2(filepath)
        page_count = pages
        if len(text.strip()) > 100:
            logger.info("PDF extracted via PyPDF2 (%d pages)", pages)
            return text, pages, meta
        logger.warning("PyPDF2 returned minimal text")
    except Exception as e:
        logger.warning("PyPDF2 failed: %s", e)

    # Text extraction failed — this is likely a scanned/image PDF
    if not gemini_api_key:
        gemini_api_key = os.environ.get("GEMINI_API_KEY")

    if gemini_api_key:
        logger.info("Attempting OCR via Gemini Vision API...")
        try:
            images = render_pdf_pages_to_images(filepath, dpi=200)
            page_count = len(images)
            text = ocr_with_gemini(images, gemini_api_key)
            meta = {"extraction_method": "gemini_ocr", "ocr_pages": page_count}
            logger.info("PDF OCR completed via Gemini (%d pages)", page_count)
            return text, page_count, meta
        except Exception as e:
            logger.error("Gemini OCR failed: %s", e)
            raise
    else:
        logger.error(
            "Image-based PDF detected but no GEMINI_API_KEY available for OCR. "
            "Set GEMINI_API_KEY environment variable or pass it to extract_pdf()."
        )
        raise RuntimeError(
            f"Image-based PDF requires OCR. Set GEMINI_API_KEY env var. "
            f"({page_count} pages detected, 0 text chars extracted)"
        )


def extract_docx(filepath: str) -> tuple[str, int, dict]:
    """Extract text from DOCX file."""
    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n\n".join(paragraphs)

    metadata = {}
    core = doc.core_properties
    if core.title:
        metadata["title"] = core.title
    if core.author:
        metadata["author"] = core.author
    if core.subject:
        metadata["subject"] = core.subject

    # Estimate "pages" from paragraph count (rough heuristic)
    page_estimate = max(1, len(paragraphs) // 8)
    logger.info("DOCX extracted (%d paragraphs, ~%d pages)", len(paragraphs), page_estimate)
    return full_text, page_estimate, metadata


def extract_txt(filepath: str) -> tuple[str, int, dict]:
    """Extract text from a plain text file."""
    text = Path(filepath).read_text(encoding="utf-8")
    # Estimate "pages" from line count (rough: ~40 lines per page)
    line_count = text.count("\n") + 1
    page_estimate = max(1, line_count // 40)
    logger.info("TXT extracted (%d lines, ~%d pages)", line_count, page_estimate)
    return text, page_estimate, {}


def extract_markdown(filepath: str) -> tuple[str, int, dict]:
    """Extract text from a Markdown file, stripping syntax."""
    raw = Path(filepath).read_text(encoding="utf-8")

    # Strip common Markdown syntax while preserving content
    text = raw
    # Remove images: ![alt](url)
    text = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", text)
    # Convert links: [text](url) → text
    text = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)
    # Remove HTML tags
    text = re.sub(r"<[^>]+>", "", text)
    # Remove heading markers but keep text
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    # Remove emphasis markers: **bold**, *italic*, __bold__, _italic_
    text = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", text)
    text = re.sub(r"_{1,2}([^_]+)_{1,2}", r"\1", text)
    # Remove inline code backticks
    text = re.sub(r"`([^`]+)`", r"\1", text)
    # Remove code block fences
    text = re.sub(r"```[a-z]*\n?", "", text)
    # Remove blockquote markers
    text = re.sub(r"^>\s?", "", text, flags=re.MULTILINE)
    # Remove horizontal rules
    text = re.sub(r"^[-*_]{3,}\s*$", "", text, flags=re.MULTILINE)

    line_count = text.count("\n") + 1
    page_estimate = max(1, line_count // 40)
    logger.info("Markdown extracted (%d lines, ~%d pages)", line_count, page_estimate)
    return text, page_estimate, {"source_format": "markdown"}


def extract_pptx_text(filepath: str) -> tuple[str, int, dict]:
    """
    Extract text from all shapes in a PPTX file.

    Useful for converting an existing presentation into a new format.
    """
    from pptx import Presentation as PptxPresentation

    prs = PptxPresentation(filepath)
    texts = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        slide_texts.append(para_text)
        if slide_texts:
            texts.append(f"--- Slide {slide_idx} ---\n" + "\n".join(slide_texts))

    full_text = "\n\n".join(texts)
    slide_count = len(prs.slides)
    logger.info("PPTX text extracted (%d slides)", slide_count)
    return full_text, slide_count, {"source_format": "pptx", "slide_count": slide_count}


def extract_content(filepath: str, gemini_api_key: str | None = None) -> ExtractedContent:
    """
    Main extraction function. Detects file type and extracts content.

    Args:
        filepath: Path to PDF, DOCX, TXT, MD, or PPTX file
        gemini_api_key: Optional Gemini API key for OCR of scanned PDFs

    Returns:
        ExtractedContent with text, metadata, and statistics
    """
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {filepath}")

    ext = path.suffix.lower()
    filename = path.name

    if ext == ".pdf":
        raw_text, page_count, metadata = extract_pdf(str(path), gemini_api_key)
        file_type = "PDF"
    elif ext in (".docx", ".doc"):
        raw_text, page_count, metadata = extract_docx(str(path))
        file_type = "DOCX"
    elif ext == ".txt":
        raw_text, page_count, metadata = extract_txt(str(path))
        file_type = "TXT"
    elif ext == ".md":
        raw_text, page_count, metadata = extract_markdown(str(path))
        file_type = "MD"
    elif ext == ".pptx":
        raw_text, page_count, metadata = extract_pptx_text(str(path))
        file_type = "PPTX"
    else:
        raise ValueError(f"Unsupported file type: {ext} (supported: .pdf, .docx, .txt, .md, .pptx)")

    text = clean_text(raw_text)
    words = text.split()
    language = detect_language(text)

    return ExtractedContent(
        text=text,
        filename=filename,
        file_type=file_type,
        page_count=page_count,
        word_count=len(words),
        char_count=len(text),
        language=language,
        metadata=metadata,
    )


def main():
    """CLI entry point for testing content extraction."""
    if len(sys.argv) < 2:
        print("Usage: python content_extractor.py <file.pdf|file.docx>")
        sys.exit(1)

    filepath = sys.argv[1]
    print(f"\nExtracting content from: {filepath}\n")

    result = extract_content(filepath)
    print(result.summary())
    print("--- Preview (first 800 chars) ---")
    print(result.preview(800))
    print("--- End Preview ---\n")
    return result


if __name__ == "__main__":
    main()
