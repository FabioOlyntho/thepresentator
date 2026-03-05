"""
OCR Converter — Converts NotebookLM image-based PPTX to editable PPTX.

Two modes:
- PDNob (default): Extract text with bounding boxes, inpaint text from image,
  rebuild slide as cleaned background + editable text boxes at original positions.
- Editable: Extract text-only via Gemini/Docling OCR, rebuild as Recodme-branded layouts.

Usage:
    python scripts/ocr_converter.py input.pptx                       # PDNob mode (default)
    python scripts/ocr_converter.py input.pptx --mode editable        # Recodme editable mode
    python scripts/ocr_converter.py input.pptx --mode editable --ocr-engine docling
"""

import argparse
import base64
import io
import json
import logging
import math
import os
import re
import sys
import tempfile
import time
import unicodedata
from dataclasses import dataclass
from pathlib import Path

import cv2
import numpy as np
import requests
from PIL import Image
from pptx import Presentation
from rapidocr_onnxruntime import RapidOCR

logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")
logger = logging.getLogger(__name__)

# Singleton RapidOCR engine (lazy init to avoid model load at import)
_rapid_ocr: RapidOCR | None = None


def _get_rapid_ocr() -> RapidOCR:
    """Get or initialize the RapidOCR engine."""
    global _rapid_ocr
    if _rapid_ocr is None:
        _rapid_ocr = RapidOCR()
    return _rapid_ocr


# ═══════════════════════════════════════════════════════════════
# PDNob-Style OCR — Text + bounding boxes + inpainting
# ═══════════════════════════════════════════════════════════════


@dataclass
class OCRTextBlock:
    """A text block with position info extracted from OCR."""
    text: str
    x_pct: float       # left edge as % of image width
    y_pct: float        # top edge as % of image height
    width_pct: float    # width as % of image width
    height_pct: float   # height as % of image height
    font_size_pt: float  # estimated from bbox height
    color: tuple[int, int, int]  # RGB


@dataclass
class ImageRegion:
    """A rectangular region of a slide image (percentage-based)."""
    x_pct: float       # left edge as % of slide width
    y_pct: float        # top edge as % of slide height
    width_pct: float    # width as % of slide width
    height_pct: float   # height as % of slide height


def extract_text_with_positions(
    image_bytes: bytes,
    slide_height_inches: float = 7.5,
) -> list[OCRTextBlock]:
    """
    Run RapidOCR on image and return text blocks with bounding box positions.

    Each result contains text, percentage-based coordinates (for slide placement),
    estimated font size, and sampled text color.

    Args:
        image_bytes: Raw PNG/JPG bytes of the slide image.
        slide_height_inches: Actual slide height for font size calibration.
            Default 7.5" (standard 16:9). NotebookLM uses 10.0".
    """
    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    img_array = np.array(img)
    img_h, img_w = img_array.shape[:2]

    engine = _get_rapid_ocr()
    results = engine(img_array)

    # RapidOCR v3.x returns (results_list, elapsed_times)
    # where each item in results_list is [box, text, score]
    if results is None or results[0] is None:
        return []

    ocr_items = results[0]
    if not ocr_items:
        return []

    blocks: list[OCRTextBlock] = []
    for item in ocr_items:
        box, text, score = item[0], item[1], float(item[2])
        if score < 0.5:
            continue
        text = normalize_text(text.strip())
        if not text:
            continue

        # box is 4 points: [[x1,y1], [x2,y2], [x3,y3], [x4,y4]]
        box = np.array(box)
        x_min = float(np.min(box[:, 0]))
        y_min = float(np.min(box[:, 1]))
        x_max = float(np.max(box[:, 0]))
        y_max = float(np.max(box[:, 1]))

        bbox_w = x_max - x_min
        bbox_h = y_max - y_min

        # Convert to percentages
        x_pct = x_min / img_w * 100
        y_pct = y_min / img_h * 100
        width_pct = bbox_w / img_w * 100
        height_pct = bbox_h / img_h * 100

        # Estimate font size: slide height * 72 DPI = reference pt
        # E.g. 7.5" → 540pt, 10.0" → 720pt
        ref_pt = slide_height_inches * 72
        font_size_pt = max(8.0, min(72.0, bbox_h / img_h * ref_pt))

        # Sample text color
        color = _sample_text_color(img_array, int(x_min), int(y_min), int(x_max), int(y_max))

        blocks.append(OCRTextBlock(
            text=text,
            x_pct=x_pct,
            y_pct=y_pct,
            width_pct=width_pct,
            height_pct=height_pct,
            font_size_pt=round(font_size_pt, 1),
            color=color,
        ))

    # Sort top-to-bottom, then left-to-right
    blocks.sort(key=lambda b: (b.y_pct, b.x_pct))
    return blocks


def _sample_text_color(
    img_array: np.ndarray,
    x1: int, y1: int, x2: int, y2: int,
) -> tuple[int, int, int]:
    """
    Detect text color by comparing bbox interior with background.

    Samples the 2px border around the bbox as background, then finds
    pixels inside that differ significantly from background.
    """
    h, w = img_array.shape[:2]
    # Clamp coordinates
    x1, y1 = max(0, x1), max(0, y1)
    x2, y2 = min(w, x2), min(h, y2)

    if x2 - x1 < 4 or y2 - y1 < 4:
        return (255, 255, 255)

    # Sample background: 3px border around bbox
    pad = 3
    border_regions = []
    # Top border
    if y1 >= pad:
        border_regions.append(img_array[max(0, y1-pad):y1, x1:x2])
    # Bottom border
    if y2 + pad <= h:
        border_regions.append(img_array[y2:min(h, y2+pad), x1:x2])
    # Left border
    if x1 >= pad:
        border_regions.append(img_array[y1:y2, max(0, x1-pad):x1])
    # Right border
    if x2 + pad <= w:
        border_regions.append(img_array[y1:y2, x2:min(w, x2+pad)])

    if not border_regions:
        return (255, 255, 255)

    bg_pixels = np.concatenate([r.reshape(-1, 3) for r in border_regions if r.size > 0])
    if len(bg_pixels) == 0:
        return (255, 255, 255)

    bg_color = bg_pixels.mean(axis=0)

    # Sample interior pixels
    interior = img_array[y1:y2, x1:x2].reshape(-1, 3).astype(float)
    if len(interior) == 0:
        return (255, 255, 255)

    # Find pixels differing from background (Euclidean distance > 60)
    diffs = np.sqrt(np.sum((interior - bg_color) ** 2, axis=1))
    text_mask = diffs > 60
    text_pixels = interior[text_mask]

    if len(text_pixels) < 3:
        # Not enough distinct pixels — fallback to white
        return (255, 255, 255)

    avg = text_pixels.mean(axis=0)
    return (int(avg[0]), int(avg[1]), int(avg[2]))


def erase_text_from_image(
    image_bytes: bytes,
    text_blocks: list[OCRTextBlock],
) -> bytes:
    """
    Erase text from image using OpenCV inpainting.

    Creates a binary mask from text block bounding boxes, dilates slightly
    to catch edge pixels, then inpaints to fill text regions.
    """
    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    img_array = np.array(img)
    img_h, img_w = img_array.shape[:2]

    # Convert to BGR for OpenCV
    img_bgr = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)

    # Create binary mask
    mask = np.zeros((img_h, img_w), dtype=np.uint8)

    for block in text_blocks:
        x1 = max(0, int(block.x_pct / 100 * img_w) - 3)
        y1 = max(0, int(block.y_pct / 100 * img_h) - 3)
        x2 = min(img_w, int((block.x_pct + block.width_pct) / 100 * img_w) + 3)
        y2 = min(img_h, int((block.y_pct + block.height_pct) / 100 * img_h) + 3)
        cv2.rectangle(mask, (x1, y1), (x2, y2), 255, -1)

    # Dilate to catch edge pixels
    kernel = np.ones((3, 3), np.uint8)
    mask = cv2.dilate(mask, kernel, iterations=1)

    # Inpaint
    result = cv2.inpaint(img_bgr, mask, inpaintRadius=5, flags=cv2.INPAINT_TELEA)

    # Encode as PNG
    _, png_bytes = cv2.imencode(".png", result)
    return png_bytes.tobytes()


def _most_common_color(
    colors: list[tuple[int, int, int]],
    distance_threshold: float = 30.0,
) -> tuple[int, int, int]:
    """
    Find the most common color by bucketing similar colors together.

    Groups colors within Euclidean distance of `distance_threshold` as "same",
    then returns the representative of the largest bucket.
    """
    if not colors:
        return (255, 255, 255)
    if len(colors) == 1:
        return colors[0]

    # Bucket colors: each bucket has a representative and a count
    buckets: list[tuple[tuple[int, int, int], int]] = []
    for c in colors:
        placed = False
        for i, (rep, count) in enumerate(buckets):
            dist = math.sqrt(sum((a - b) ** 2 for a, b in zip(c, rep)))
            if dist < distance_threshold:
                buckets[i] = (rep, count + 1)
                placed = True
                break
        if not placed:
            buckets.append((c, 1))

    # Return the representative of the largest bucket
    buckets.sort(key=lambda x: x[1], reverse=True)
    return buckets[0][0]


def _merge_block_group(blocks: list[OCRTextBlock]) -> OCRTextBlock:
    """
    Merge a group of OCRTextBlocks into a single block.

    - Text: concatenated with newlines for vertically-separated lines,
      or spaces for same-line fragments.
    - Bounding box: union rectangle.
    - Font size: max in group (titles dominate).
    - Color: most common color in group.
    """
    if len(blocks) == 1:
        return blocks[0]

    # Sort by y then x for proper reading order
    blocks = sorted(blocks, key=lambda b: (b.y_pct, b.x_pct))

    # Average font height for gap detection
    avg_height = sum(b.height_pct for b in blocks) / len(blocks)

    # Build text: use space for same-line, newline for different lines
    text_parts: list[str] = []
    for i, b in enumerate(blocks):
        if i == 0:
            text_parts.append(b.text)
            continue
        prev = blocks[i - 1]
        vertical_gap = b.y_pct - (prev.y_pct + prev.height_pct)
        if vertical_gap < avg_height * 0.3:
            # Same line or very close — join with space
            text_parts.append(" " + b.text)
        else:
            text_parts.append("\n" + b.text)

    merged_text = "".join(text_parts).strip()

    # Union bounding box
    x_min = min(b.x_pct for b in blocks)
    y_min = min(b.y_pct for b in blocks)
    x_max = max(b.x_pct + b.width_pct for b in blocks)
    y_max = max(b.y_pct + b.height_pct for b in blocks)

    # Font size: max in group
    font_size = max(b.font_size_pt for b in blocks)

    # Color: most common
    color = _most_common_color([b.color for b in blocks])

    return OCRTextBlock(
        text=merged_text,
        x_pct=x_min,
        y_pct=y_min,
        width_pct=x_max - x_min,
        height_pct=y_max - y_min,
        font_size_pt=font_size,
        color=color,
    )


def group_text_blocks(blocks: list[OCRTextBlock]) -> list[OCRTextBlock]:
    """
    Merge spatially adjacent, visually similar text blocks using Union-Find.

    Reduces overlapping single-line OCR detections into cohesive multi-line
    editable text boxes. Keeps columns, titles, and differently-styled text
    separate.

    Merge criteria (ALL must be true for a pair):
    1. Vertical proximity: gap < 1.8x average font height
    2. Horizontal alignment: left-edge within 4% OR horizontal overlap > 50%
    3. Font size similarity: within 4pt
    4. Color similarity: RGB Euclidean distance < 60
    5. Height sanity: combined block < 35% of slide height
    """
    n = len(blocks)
    if n <= 1:
        return list(blocks)

    # Union-Find
    parent = list(range(n))

    def find(x: int) -> int:
        while parent[x] != x:
            parent[x] = parent[parent[x]]  # path compression
            x = parent[x]
        return x

    def union(a: int, b: int) -> None:
        ra, rb = find(a), find(b)
        if ra != rb:
            parent[ra] = rb

    # Average font height (as % of slide)
    avg_font_h = sum(b.height_pct for b in blocks) / n

    # Sort by y for efficient scanning
    sorted_indices = sorted(range(n), key=lambda i: blocks[i].y_pct)

    # Compare pairs — only scan within vertical proximity window
    vertical_limit = avg_font_h * 3.6  # 2x the 1.8 threshold for pre-filter
    for idx_a in range(len(sorted_indices)):
        i = sorted_indices[idx_a]
        bi = blocks[i]
        bi_bottom = bi.y_pct + bi.height_pct

        for idx_b in range(idx_a + 1, len(sorted_indices)):
            j = sorted_indices[idx_b]
            bj = blocks[j]

            # Pre-filter: if too far vertically, stop scanning
            if bj.y_pct - bi_bottom > vertical_limit:
                break

            # Criterion 1: Vertical proximity — gap < 1.8x avg font height
            vertical_gap = bj.y_pct - bi_bottom
            if vertical_gap > avg_font_h * 1.8:
                continue

            # Criterion 2: Horizontal alignment
            # Option A: left-edge aligned within 4%
            left_aligned = abs(bi.x_pct - bj.x_pct) < 4.0
            # Option B: horizontal overlap > 50% of narrower block
            overlap_left = max(bi.x_pct, bj.x_pct)
            overlap_right = min(bi.x_pct + bi.width_pct, bj.x_pct + bj.width_pct)
            overlap = max(0.0, overlap_right - overlap_left)
            narrower_width = min(bi.width_pct, bj.width_pct)
            horizontal_overlap = (overlap / narrower_width) > 0.5 if narrower_width > 0 else False
            if not (left_aligned or horizontal_overlap):
                continue

            # Criterion 3: Font size similarity (within 4pt)
            if abs(bi.font_size_pt - bj.font_size_pt) > 4.0:
                continue

            # Criterion 4: Color similarity (Euclidean distance < 60)
            color_dist = math.sqrt(sum(
                (a - b) ** 2 for a, b in zip(bi.color, bj.color)
            ))
            if color_dist > 60.0:
                continue

            # Criterion 5: Combined height < 35% of slide
            combined_top = min(bi.y_pct, bj.y_pct)
            combined_bottom = max(bi_bottom, bj.y_pct + bj.height_pct)
            if (combined_bottom - combined_top) > 35.0:
                continue

            # All criteria met — merge
            union(i, j)

    # Build groups from union-find
    groups: dict[int, list[int]] = {}
    for i in range(n):
        root = find(i)
        groups.setdefault(root, []).append(i)

    # Merge each group
    merged: list[OCRTextBlock] = []
    for indices in groups.values():
        group_blocks = [blocks[i] for i in indices]
        mb = _merge_block_group(group_blocks)
        # Filter whitespace-only blocks
        if mb.text.strip():
            merged.append(mb)

    # Sort merged blocks top-to-bottom, left-to-right
    merged.sort(key=lambda b: (b.y_pct, b.x_pct))
    return merged


def segment_slide_image(
    cleaned_image_bytes: bytes,
    bg_tolerance: float = 35.0,
    min_region_pct: float = 1.5,
    merge_gap_pct: float = 3.0,
) -> list[ImageRegion]:
    """
    Segment a cleaned slide image into visual content regions.

    Detects distinct illustrations/icons on a uniform background by
    finding clusters of non-background pixels. Designed for cleaned
    images (after text erasure) where the only remaining content is
    the visual illustrations.

    Algorithm:
    1. Sample background color from image corners
    2. Threshold: pixels far from background = content
    3. Morphological close to connect nearby content within an icon
    4. Find contours → bounding boxes
    5. Filter by minimum area, merge overlapping/nearby boxes
    6. Fallback: < 2 regions → single full-bleed

    Args:
        cleaned_image_bytes: Raw PNG bytes of the text-erased image.
        bg_tolerance: Euclidean color distance threshold for background vs content.
        min_region_pct: Minimum region area as % of total image to keep.
        merge_gap_pct: Maximum gap (% of image diagonal) to merge nearby boxes.

    Returns:
        List of ImageRegion objects (percentage-based coordinates).
    """
    img = Image.open(io.BytesIO(cleaned_image_bytes)).convert("RGB")
    img_array = np.array(img)
    h, w = img_array.shape[:2]

    # Step 1: Sample background from image corners (20x20px each)
    corner_size = min(20, h // 10, w // 10)
    corners = [
        img_array[:corner_size, :corner_size],
        img_array[:corner_size, -corner_size:],
        img_array[-corner_size:, :corner_size],
        img_array[-corner_size:, -corner_size:],
    ]
    bg_pixels = np.concatenate([c.reshape(-1, 3) for c in corners])
    bg_color = bg_pixels.mean(axis=0).astype(float)

    # Step 2: Per-pixel Euclidean distance from background → binary mask
    diff = np.sqrt(np.sum((img_array.astype(float) - bg_color) ** 2, axis=2))
    mask = (diff > bg_tolerance).astype(np.uint8) * 255

    # Step 3: Morphological close to connect nearby content within an icon
    kernel_size = max(3, int(min(h, w) * 0.04))
    if kernel_size % 2 == 0:
        kernel_size += 1
    kernel = np.ones((kernel_size, kernel_size), np.uint8)
    mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, kernel)

    # Remove small noise with morphological open
    small_kernel = np.ones((5, 5), np.uint8)
    mask = cv2.morphologyEx(mask, cv2.MORPH_OPEN, small_kernel)

    # Step 4: Find contours
    contours, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    # Step 5: Get bounding boxes, filter by minimum area
    min_area = h * w * min_region_pct / 100
    boxes: list[tuple[int, int, int, int]] = []
    for cnt in contours:
        x, y, bw, bh = cv2.boundingRect(cnt)
        if bw * bh >= min_area:
            boxes.append((x, y, x + bw, y + bh))

    if not boxes:
        return [ImageRegion(x_pct=0.0, y_pct=0.0, width_pct=100.0, height_pct=100.0)]

    # Step 6: Merge overlapping or nearby boxes
    diag = math.sqrt(h * h + w * w)
    merge_gap_px = int(diag * merge_gap_pct / 100)
    merged = _merge_boxes(boxes, merge_gap_px)

    # Convert to ImageRegion (percentage-based), sorted left-to-right then top-to-bottom
    regions: list[ImageRegion] = []
    for (x1, y1, x2, y2) in merged:
        regions.append(ImageRegion(
            x_pct=x1 / w * 100,
            y_pct=y1 / h * 100,
            width_pct=(x2 - x1) / w * 100,
            height_pct=(y2 - y1) / h * 100,
        ))
    regions.sort(key=lambda r: (r.y_pct, r.x_pct))

    # Fallback: < 2 regions → single full-bleed
    if len(regions) < 2:
        return [ImageRegion(x_pct=0.0, y_pct=0.0, width_pct=100.0, height_pct=100.0)]

    return regions


def _merge_boxes(
    boxes: list[tuple[int, int, int, int]],
    gap: int,
) -> list[tuple[int, int, int, int]]:
    """
    Merge overlapping or nearby bounding boxes.

    Iteratively merges boxes that overlap or are within `gap` pixels
    of each other until no more merges are possible.

    Args:
        boxes: List of (x1, y1, x2, y2) bounding boxes.
        gap: Maximum pixel distance to consider boxes as "nearby".

    Returns:
        Merged list of (x1, y1, x2, y2) bounding boxes.
    """
    if not boxes:
        return []

    merged = [list(b) for b in boxes]
    changed = True
    while changed:
        changed = False
        new_merged: list[list[int]] = []
        used: set[int] = set()
        for i in range(len(merged)):
            if i in used:
                continue
            box = list(merged[i])
            for j in range(i + 1, len(merged)):
                if j in used:
                    continue
                other = merged[j]
                # Check overlap or proximity
                if (box[0] <= other[2] + gap and box[2] >= other[0] - gap and
                        box[1] <= other[3] + gap and box[3] >= other[1] - gap):
                    box[0] = min(box[0], other[0])
                    box[1] = min(box[1], other[1])
                    box[2] = max(box[2], other[2])
                    box[3] = max(box[3], other[3])
                    used.add(j)
                    changed = True
            new_merged.append(box)
        merged = new_merged

    return [tuple(b) for b in merged]


def crop_image_region(image_bytes: bytes, region: ImageRegion) -> bytes:
    """
    Crop a rectangular sub-image from the full slide image.

    Args:
        image_bytes: Raw PNG bytes of the full slide image.
        region: ImageRegion defining the crop area (percentage-based).

    Returns:
        PNG bytes of the cropped sub-image.
    """
    img = Image.open(io.BytesIO(image_bytes))
    w, h = img.size
    left = int(region.x_pct / 100 * w)
    top = int(region.y_pct / 100 * h)
    right = int((region.x_pct + region.width_pct) / 100 * w)
    bottom = int((region.y_pct + region.height_pct) / 100 * h)
    cropped = img.crop((left, top, right, bottom))
    buf = io.BytesIO()
    cropped.save(buf, format="PNG")
    return buf.getvalue()


def remove_background(image_bytes: bytes) -> bytes:
    """
    Remove background from an image using rembg (U2-Net AI model).

    Takes a PNG image (typically a cropped region from a slide) and returns
    a transparent-background RGBA PNG. Works best on illustrations/icons
    on solid-colored backgrounds.

    Args:
        image_bytes: Raw PNG bytes of the input image.

    Returns:
        PNG bytes with transparent background (RGBA).
    """
    from rembg import remove

    result = remove(image_bytes)
    return result


def convert_pdnob_style(
    input_pptx: str,
    output_pptx: str | None = None,
    pdnob_level: str = "full",
) -> dict:
    """
    PDNob-style conversion: extract text + positions, erase text from image,
    rebuild with cleaned image background + editable text boxes.

    Args:
        input_pptx: Path to NotebookLM-generated PPTX.
        output_pptx: Output path (auto-generated if None).
        pdnob_level: "ocr_only" (text boxes only), "remove_bg" (segmented images only),
                     or "full" (both text boxes + segmented images).

    Returns:
        Dict with success status, files, metadata, timing.
    """
    start_time = time.time()
    result: dict = {
        "success": False,
        "files": {},
        "metadata": {},
        "timing": {},
    }

    # Step 1: Extract images + source dimensions
    logger.info("Step 1: Extracting slide images from %s", input_pptx)
    t0 = time.time()
    slide_images = extract_slide_images(input_pptx)
    src_w, src_h = get_source_slide_dims(input_pptx)
    logger.info("Source slide dimensions: %.2f\" x %.2f\"", src_w, src_h)

    if not slide_images:
        return {
            "success": False,
            "error": "No images found in PPTX. Is this a NotebookLM presentation?",
        }

    result["timing"]["extraction"] = round(time.time() - t0, 2)
    total_slides = len(slide_images)
    logger.info("Found %d slide images", total_slides)

    # Step 2: OCR + inpaint each slide
    logger.info("Step 2: PDNob OCR + inpainting")
    t0 = time.time()

    from slide_builder import SlideBuilder

    builder = SlideBuilder(slide_width_inches=src_w, slide_height_inches=src_h)
    temp_files: list[str] = []

    for slide_num, img_bytes in slide_images:
        logger.info("  Slide %d/%d: OCR...", slide_num, total_slides)
        try:
            text_blocks = extract_text_with_positions(img_bytes, slide_height_inches=src_h)
            logger.info("    Found %d text blocks", len(text_blocks))

            # Erase text from image (uses raw blocks for pixel-accurate mask)
            cleaned_bytes = erase_text_from_image(img_bytes, text_blocks)

            # Merge blocks for clean placement (fewer, non-overlapping text boxes)
            merged_blocks = group_text_blocks(text_blocks)
            logger.info("    %d blocks -> %d groups", len(text_blocks), len(merged_blocks))

            # Segmentation only for remove_bg and full
            if pdnob_level in ("remove_bg", "full"):
                regions = segment_slide_image(cleaned_bytes)
                logger.info("    %d image regions", len(regions))
            else:
                regions = []

            # Text boxes only for ocr_only and full
            if pdnob_level in ("ocr_only", "full"):
                blocks_for_slide = merged_blocks
            else:
                blocks_for_slide = []

            logger.info("    %d blocks -> %d groups, %d regions (level=%s)",
                        len(text_blocks), len(merged_blocks), len(regions), pdnob_level)

            # Save cleaned image to temp file
            tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            tmp.write(cleaned_bytes)
            tmp.close()
            temp_files.append(tmp.name)

            # Build slide: image region(s) + text boxes at OCR positions
            if len(regions) > 1:
                builder.build_pdnob_slide(
                    tmp.name, blocks_for_slide,
                    image_regions=regions, cleaned_bytes=cleaned_bytes,
                )
            else:
                builder.build_pdnob_slide(tmp.name, blocks_for_slide)
        except Exception as e:
            logger.error("    Slide %d failed: %s — using original image", slide_num, e)
            # Fallback: use original image without text overlay
            tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            tmp.write(img_bytes)
            tmp.close()
            temp_files.append(tmp.name)
            builder.build_pdnob_slide(tmp.name, [])

    result["timing"]["ocr_inpaint"] = round(time.time() - t0, 2)

    # Step 3: Save PPTX
    logger.info("Step 3: Saving output PPTX")
    t0 = time.time()

    if not output_pptx:
        input_stem = Path(input_pptx).stem
        output_dir = Path(input_pptx).parent
        output_pptx = str(output_dir / f"{input_stem}_pdnob.pptx")

    builder.save(output_pptx)
    result["timing"]["saving"] = round(time.time() - t0, 2)

    # Cleanup temp files
    for f in temp_files:
        try:
            Path(f).unlink(missing_ok=True)
        except Exception:
            pass

    total_time = round(time.time() - start_time, 2)
    result["success"] = True
    result["files"] = {"pptx": output_pptx}
    result["metadata"] = {
        "title": Path(input_pptx).stem,
        "total_slides": total_slides,
        "mode": "pdnob",
        "pdnob_level": pdnob_level,
    }
    result["timing"]["total"] = total_time

    logger.info("=" * 50)
    logger.info("PDNob Conversion Complete")
    logger.info("  Slides: %d", total_slides)
    logger.info("  Time: %.1fs", total_time)
    logger.info("  Output: %s", output_pptx)
    logger.info("=" * 50)

    return result

# Ensure sibling modules are importable
_scripts_dir = str(Path(__file__).parent)
if _scripts_dir not in sys.path:
    sys.path.insert(0, _scripts_dir)

PROMPTS_DIR = Path(__file__).parent.parent / "prompts"

# Valid slide types (must match slide_builder.py's dispatcher)
VALID_SLIDE_TYPES = {"title", "section", "content", "comparison", "data", "quote", "conclusion"}


def load_ocr_prompt() -> str:
    """Load the OCR extraction prompt template."""
    path = PROMPTS_DIR / "ocr_extraction_prompt.txt"
    if not path.exists():
        raise FileNotFoundError(f"OCR prompt not found: {path}")
    return path.read_text(encoding="utf-8")


def normalize_text(text: str) -> str:
    """
    Normalize Unicode text to NFC form and clean common OCR artifacts.

    Fixes issues like:
    - Decomposed accents (e.g. e + ´ → é) via NFC normalization
    - Mojibake from encoding mismatches
    - Non-breaking spaces and other invisible Unicode characters
    """
    if not text:
        return text
    # NFC normalization: compose decomposed characters
    text = unicodedata.normalize("NFC", text)
    # Replace non-breaking spaces with regular spaces
    text = text.replace("\u00a0", " ")
    # Replace zero-width characters
    text = text.replace("\u200b", "")  # zero-width space
    text = text.replace("\u200c", "")  # zero-width non-joiner
    text = text.replace("\u200d", "")  # zero-width joiner
    text = text.replace("\ufeff", "")  # BOM
    return text


def _patch_hf_symlinks():
    """Patch HuggingFace Hub to use file copy instead of symlinks on Windows."""
    if sys.platform != "win32":
        return
    try:
        import huggingface_hub.file_download as hf_dl
        import shutil

        original = hf_dl._create_symlink

        def _copy_fallback(src, dst, new_blob=False):
            try:
                original(src, dst, new_blob)
            except OSError:
                Path(dst).parent.mkdir(parents=True, exist_ok=True)
                if Path(dst).exists():
                    Path(dst).unlink()
                shutil.copy2(src, dst)

        hf_dl._create_symlink = _copy_fallback
    except (ImportError, AttributeError):
        pass


def extract_slide_content_docling(
    image_bytes: bytes,
    slide_number: int,
    total_slides: int,
) -> dict:
    """
    Extract structured content from a slide image using Docling OCR (offline).

    Docling uses local ML models for text extraction — no API key needed.
    It produces raw text that we then parse into our slide JSON schema.

    Args:
        image_bytes: Raw image bytes (PNG/JPEG).
        slide_number: 1-indexed slide number.
        total_slides: Total number of slides.

    Returns:
        Dict with slide content fields matching SlideSpec.
    """
    import tempfile

    _patch_hf_symlinks()

    from docling.document_converter import DocumentConverter

    # Docling works on files, so write image to a temp file
    suffix = ".png" if image_bytes[:4] == b"\x89PNG" else ".jpg"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(image_bytes)
        tmp_path = tmp.name

    try:
        logger.info("OCR slide %d/%d via Docling (offline)...", slide_number, total_slides)
        converter = DocumentConverter()
        result = converter.convert(tmp_path)
        raw_text = result.document.export_to_markdown()
        raw_text = normalize_text(raw_text)
        # Clean Docling markdown artifacts
        raw_text = _clean_docling_markdown(raw_text)
    finally:
        Path(tmp_path).unlink(missing_ok=True)

    # Parse the raw OCR text into our structured slide format
    return _parse_raw_text_to_slide(raw_text, slide_number, total_slides)


def _clean_docling_markdown(text: str) -> str:
    """Remove Docling-specific markdown artifacts from OCR output."""
    # Remove HTML image comments: <!-- image -->
    text = re.sub(r"<!--\s*image\s*-->", "", text)
    # Remove Docling HTML picture tags
    text = re.sub(r"<picture>.*?</picture>", "", text, flags=re.DOTALL)
    # Normalize excessive whitespace from removals
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _parse_raw_text_to_slide(raw_text: str, slide_number: int, total_slides: int) -> dict:
    """
    Parse raw OCR text into structured slide JSON matching our schema.

    Heuristic parsing: first line → title, bullet-like lines → bullet_points,
    remaining → body. Classifies slide type based on position and content.
    """
    lines = [l.strip() for l in raw_text.strip().split("\n") if l.strip()]

    if not lines:
        return {
            "type": classify_slide_type(slide_number, total_slides, ""),
            "title": f"Slide {slide_number}",
            "body": "",
            "bullet_points": [],
            "speaker_notes": "",
        }

    # First non-empty line is the title (strip markdown heading markers)
    title = re.sub(r"^#{1,6}\s*", "", lines[0])
    # Truncate to ~10 words for micro-copy
    title_words = title.split()
    if len(title_words) > 10:
        title = " ".join(title_words[:10])

    remaining = lines[1:]
    bullet_points = []
    body_lines = []
    left_column = []
    right_column = []
    subtitle = ""

    for line in remaining:
        # Detect bullet-like patterns
        if re.match(r"^[-•*▪▸►]\s+", line):
            cleaned = re.sub(r"^[-•*▪▸►]\s+", "", line)
            bullet_points.append(cleaned)
        elif re.match(r"^\d+[.)]\s+", line):
            cleaned = re.sub(r"^\d+[.)]\s+", "", line)
            bullet_points.append(cleaned)
        elif line.startswith("**") and line.endswith("**"):
            # Bold line — could be a subtitle or section header
            if not subtitle:
                subtitle = line.strip("* ")
            else:
                body_lines.append(line.strip("* "))
        else:
            body_lines.append(line)

    body = " ".join(body_lines)
    # Truncate body to ~30 words for micro-copy
    body_words = body.split()
    if len(body_words) > 30:
        body = " ".join(body_words[:30])

    # Format bullets as "Keyword — explanation" if they contain a colon or dash
    formatted_bullets = []
    for bp in bullet_points:
        # If already has em dash, keep it
        if " — " in bp:
            formatted_bullets.append(bp)
        # Convert "Key: explanation" to "Key — explanation"
        elif ": " in bp and bp.index(": ") < 30:
            parts = bp.split(": ", 1)
            formatted_bullets.append(f"{parts[0]} — {parts[1]}")
        else:
            formatted_bullets.append(bp)

    # Determine slide type heuristically
    if slide_number == 1:
        slide_type = "title"
    elif slide_number == total_slides:
        slide_type = "conclusion"
    elif len(formatted_bullets) >= 2 and any("|" in bp for bp in formatted_bullets):
        slide_type = "comparison"
    elif re.search(r"\d{2,}%|\d{1,3}\.\d", title + body):
        slide_type = "data"
    elif len(formatted_bullets) >= 2:
        slide_type = "content"
    else:
        slide_type = "content"

    return {
        "type": slide_type,
        "title": title,
        "subtitle": subtitle,
        "body": body,
        "bullet_points": formatted_bullets,
        "left_column": left_column,
        "right_column": right_column,
        "left_header": "",
        "right_header": "",
        "checkbox_items": [],
        "speaker_notes": "",
    }


def extract_slide_images(
    pptx_path: str,
) -> list[tuple[int, bytes]]:
    """
    Extract the primary image from each slide in a NotebookLM PPTX.

    NotebookLM slides are single full-bleed images (17.8" x 10.0").
    Each slide contains one picture shape covering the entire canvas.

    Args:
        pptx_path: Path to NotebookLM-generated PPTX file.

    Returns:
        List of (slide_number, image_bytes) tuples, 1-indexed.
    """
    prs = Presentation(pptx_path)
    results = []

    for idx, slide in enumerate(prs.slides, start=1):
        image_bytes = None
        largest_area = 0

        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                # Pick the largest image on the slide (should be the full-bleed one)
                area = shape.width * shape.height
                if area > largest_area:
                    largest_area = area
                    image_bytes = shape.image.blob

        if image_bytes:
            results.append((idx, image_bytes))
            logger.info("Extracted image from slide %d (%d bytes)", idx, len(image_bytes))
        else:
            logger.warning("No image found on slide %d", idx)

    return results


def get_source_slide_dims(pptx_path: str) -> tuple[float, float]:
    """
    Read source PPTX slide dimensions.

    Returns:
        (width_inches, height_inches) tuple.
    """
    prs = Presentation(pptx_path)
    return (prs.slide_width / 914400, prs.slide_height / 914400)


def classify_slide_type(slide_number: int, total_slides: int, raw_type: str) -> str:
    """
    Validate and apply heuristic fallbacks for slide type classification.

    Args:
        slide_number: 1-indexed slide number.
        total_slides: Total number of slides.
        raw_type: The type returned by Gemini Vision.

    Returns:
        A valid slide type string.
    """
    if raw_type in VALID_SLIDE_TYPES:
        return raw_type

    # Heuristic fallbacks
    if slide_number == 1:
        return "title"
    if slide_number == total_slides:
        return "conclusion"
    return "content"


def extract_slide_content(
    image_bytes: bytes,
    slide_number: int,
    total_slides: int,
    api_key: str,
    model: str = "gemini-2.5-flash",
) -> dict:
    """
    Send a slide image to Gemini Vision and extract structured content.

    Args:
        image_bytes: Raw image bytes (PNG/JPEG).
        slide_number: 1-indexed slide number.
        total_slides: Total number of slides.
        api_key: Gemini API key.
        model: Gemini model for vision extraction.

    Returns:
        Dict with slide content fields matching SlideSpec.
    """
    prompt_template = load_ocr_prompt()
    prompt = prompt_template.replace("{slide_number}", str(slide_number))
    prompt = prompt.replace("{total_slides}", str(total_slides))

    b64_image = base64.b64encode(image_bytes).decode("utf-8")

    # Detect MIME type from image header
    if image_bytes[:4] == b"\x89PNG":
        mime_type = "image/png"
    elif image_bytes[:2] == b"\xff\xd8":
        mime_type = "image/jpeg"
    else:
        mime_type = "image/png"  # Default

    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"

    payload = {
        "contents": [{
            "parts": [
                {"text": prompt},
                {
                    "inlineData": {
                        "mimeType": mime_type,
                        "data": b64_image,
                    }
                },
            ]
        }],
        "generationConfig": {
            "temperature": 0.2,
            "maxOutputTokens": 4096,
            "responseMimeType": "application/json",
        },
    }

    logger.info("OCR slide %d/%d via Gemini Vision...", slide_number, total_slides)
    resp = requests.post(url, json=payload, timeout=120)
    resp.raise_for_status()
    data = resp.json()

    if "candidates" not in data or not data["candidates"]:
        raise RuntimeError(f"Gemini returned no candidates for slide {slide_number}")

    raw_text = data["candidates"][0]["content"]["parts"][0]["text"]

    # Normalize Unicode before parsing
    raw_text = normalize_text(raw_text)

    # Parse JSON response
    try:
        content = json.loads(raw_text)
    except json.JSONDecodeError:
        # Try stripping markdown fences
        cleaned = raw_text.strip()
        cleaned = re.sub(r"^```(?:json)?\s*\n?", "", cleaned)
        cleaned = re.sub(r"\n?```\s*$", "", cleaned)
        content = json.loads(cleaned)

    # Validate and fix slide type
    raw_type = content.get("type", "content")
    content["type"] = classify_slide_type(slide_number, total_slides, raw_type)

    return content


def content_to_slidespec(content: dict, slide_number: int):
    """
    Convert extracted content dict to a SlideSpec object.

    Applies Unicode normalization to all text fields to fix encoding artifacts.

    Args:
        content: Dict from OCR extraction (Gemini or Docling).
        slide_number: 1-indexed slide number.

    Returns:
        SlideSpec instance.
    """
    from gemini_client import SlideSpec

    # Normalize all text fields
    def _n(val):
        if isinstance(val, str):
            return normalize_text(val)
        if isinstance(val, list):
            return [normalize_text(v) if isinstance(v, str) else v for v in val]
        return val

    return SlideSpec(
        number=slide_number,
        type=content.get("type", "content"),
        title=_n(content.get("title", "")),
        subtitle=_n(content.get("subtitle", "")),
        body=_n(content.get("body", "")),
        bullet_points=_n(content.get("bullet_points", [])),
        visual_concept="",  # Not needed for editable mode
        speaker_notes=_n(content.get("speaker_notes", "")),
        source_reference="",
        left_column=_n(content.get("left_column", [])),
        right_column=_n(content.get("right_column", [])),
        left_header=_n(content.get("left_header", "")),
        right_header=_n(content.get("right_header", "")),
        checkbox_items=_n(content.get("checkbox_items", [])),
    )


def convert_notebooklm_to_editable(
    input_pptx: str,
    output_pptx: str | None = None,
    api_key: str | None = None,
    brand_path: str | None = None,
    model: str = "gemini-2.5-flash",
    ocr_engine: str = "gemini",
) -> dict:
    """
    Full pipeline: NotebookLM image PPTX → OCR → editable Recodme PPTX.

    Args:
        input_pptx: Path to NotebookLM-generated PPTX.
        output_pptx: Output path for editable PPTX (auto-generated if None).
        api_key: Gemini API key (falls back to GEMINI_API_KEY env var).
        brand_path: Path to brand config JSON (uses default Recodme if None).
        model: Gemini model for vision extraction.
        ocr_engine: OCR engine to use: "gemini" (default, API-based) or "docling" (offline).

    Returns:
        Dict with success status, file paths, and metadata.
    """
    start_time = time.time()

    if ocr_engine not in ("gemini", "docling"):
        return {
            "success": False,
            "error": f"Unknown OCR engine: {ocr_engine}. Use 'gemini' or 'docling'.",
        }

    # Gemini engine requires an API key; Docling works offline
    if ocr_engine == "gemini":
        api_key = api_key or os.environ.get("GEMINI_API_KEY")
        if not api_key:
            return {
                "success": False,
                "error": "No API key. Set GEMINI_API_KEY or pass api_key parameter.",
            }

    result = {
        "success": False,
        "files": {},
        "metadata": {},
        "timing": {},
    }

    # Step 1: Extract images from PPTX
    logger.info("Step 1: Extracting slide images from %s", input_pptx)
    t0 = time.time()
    slide_images = extract_slide_images(input_pptx)

    if not slide_images:
        return {
            "success": False,
            "error": "No images found in PPTX. Is this a NotebookLM presentation?",
        }

    result["timing"]["extraction"] = round(time.time() - t0, 2)
    total_slides = len(slide_images)
    logger.info("Found %d slide images", total_slides)

    # Step 2: OCR each slide
    logger.info("Step 2: OCR extraction via %s", ocr_engine.capitalize())
    t0 = time.time()

    from gemini_client import PresentationSpec

    slide_specs = []
    for i, (slide_num, img_bytes) in enumerate(slide_images):
        # Rate limiting for Gemini: pause between API calls (skip first)
        if ocr_engine == "gemini" and i > 0:
            time.sleep(1.0)

        try:
            if ocr_engine == "docling":
                content = extract_slide_content_docling(
                    image_bytes=img_bytes,
                    slide_number=slide_num,
                    total_slides=total_slides,
                )
            else:
                content = extract_slide_content(
                    image_bytes=img_bytes,
                    slide_number=slide_num,
                    total_slides=total_slides,
                    api_key=api_key,
                    model=model,
                )
            spec = content_to_slidespec(content, slide_num)
            slide_specs.append(spec)
            logger.info(
                "  Slide %d [%s]: %s",
                slide_num, spec.type, spec.title[:60],
            )
        except Exception as e:
            logger.error("  Slide %d OCR failed: %s", slide_num, e)
            # Create a fallback content slide
            from gemini_client import SlideSpec
            fallback = SlideSpec(
                number=slide_num,
                type=classify_slide_type(slide_num, total_slides, ""),
                title=f"Slide {slide_num}",
                body="(OCR extraction failed for this slide)",
                speaker_notes="",
            )
            slide_specs.append(fallback)

    result["timing"]["ocr"] = round(time.time() - t0, 2)

    # Build PresentationSpec
    title = slide_specs[0].title if slide_specs else "Untitled"
    subtitle = slide_specs[0].subtitle if slide_specs else ""
    presentation_spec = PresentationSpec(
        title=title,
        subtitle=subtitle,
        language="auto",  # Preserved from original
        source_document=Path(input_pptx).name,
        themes=[],
        slides=slide_specs,
    )

    # Step 3: Build editable PPTX
    logger.info("Step 3: Building editable Recodme PPTX")
    t0 = time.time()

    from slide_builder import SlideBuilder, BrandConfig

    if brand_path:
        brand = BrandConfig.from_json(brand_path)
    else:
        default_brand = Path(__file__).parent.parent / "config" / "brand.json"
        if default_brand.exists():
            brand = BrandConfig.from_json(str(default_brand))
        else:
            brand = BrandConfig.default()

    builder = SlideBuilder(brand=brand, editable_mode=True)
    builder.build_presentation(presentation_spec)

    # Determine output path
    if not output_pptx:
        input_stem = Path(input_pptx).stem
        output_dir = Path(input_pptx).parent
        output_pptx = str(output_dir / f"{input_stem}_editable.pptx")

    builder.save(output_pptx)
    result["timing"]["building"] = round(time.time() - t0, 2)

    # Save specs JSON alongside
    specs_path = Path(output_pptx).with_suffix(".json")
    from gemini_client import GeminiClient
    client = GeminiClient()
    specs_path.write_text(client.specs_to_json(presentation_spec), encoding="utf-8")

    # Results
    total_time = round(time.time() - start_time, 2)
    result["success"] = True
    result["files"] = {
        "pptx": output_pptx,
        "specs_json": str(specs_path),
    }
    result["metadata"] = {
        "title": title,
        "total_slides": total_slides,
        "slide_types": [s.type for s in slide_specs],
        "ocr_engine": ocr_engine,
    }
    result["timing"]["total"] = total_time

    logger.info("=" * 50)
    logger.info("OCR Conversion Complete")
    logger.info("  Title: %s", title)
    logger.info("  Slides: %d", total_slides)
    logger.info("  Time: %.1fs (extract=%.1fs, ocr=%.1fs, build=%.1fs)",
                total_time, result["timing"]["extraction"],
                result["timing"]["ocr"], result["timing"]["building"])
    logger.info("  Output: %s", output_pptx)
    logger.info("=" * 50)

    return result


def main():
    """CLI entry point for OCR conversion."""
    parser = argparse.ArgumentParser(
        description="Convert NotebookLM image PPTX to editable PPTX",
    )
    parser.add_argument("input_pptx", help="Path to NotebookLM-generated PPTX file")
    parser.add_argument("--output", "-o", help="Output path for PPTX")
    parser.add_argument("--mode", choices=["pdnob", "editable"], default="pdnob",
                        help="Conversion mode: 'pdnob' (image+text, default) or 'editable' (Recodme layouts)")
    parser.add_argument("--brand", help="Brand config JSON path (editable mode only)")
    parser.add_argument("--model", default="gemini-2.5-flash",
                        help="Gemini model for Vision OCR (editable mode, default: gemini-2.5-flash)")
    parser.add_argument("--ocr-engine", choices=["gemini", "docling"], default="gemini",
                        help="OCR engine for editable mode: 'gemini' (API) or 'docling' (offline)")
    parser.add_argument("--pdnob-level", choices=["ocr_only", "remove_bg", "full"], default="full",
                        help="PDNob level: 'ocr_only' (text only), 'remove_bg' (images only), 'full' (both)")

    args = parser.parse_args()

    if not Path(args.input_pptx).exists():
        print(f"Error: File not found: {args.input_pptx}")
        sys.exit(1)

    if args.mode == "pdnob":
        result = convert_pdnob_style(
            input_pptx=args.input_pptx,
            output_pptx=args.output,
            pdnob_level=args.pdnob_level,
        )
    else:
        result = convert_notebooklm_to_editable(
            input_pptx=args.input_pptx,
            output_pptx=args.output,
            brand_path=args.brand,
            model=args.model,
            ocr_engine=args.ocr_engine,
        )

    if not result["success"]:
        print(f"\nError: {result.get('error', 'Unknown error')}")
        sys.exit(1)

    print(f"\nOutput: {result['files']['pptx']}")


if __name__ == "__main__":
    main()
