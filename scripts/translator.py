"""
Translator — Translate PPTX presentations between languages via Gemini API.

Two modes:
1. translate_pptx(): Translate an existing PPTX file in-place (preserving layout/formatting)
2. translate_specs(): Translate SlideSpec content before building a new PPTX

Usage:
    python scripts/translator.py input.pptx --from ES --to EN [--output translated.pptx]
"""

import argparse
import json
import logging
import os
import re
import sys
import time
from copy import deepcopy
from pathlib import Path

import requests
from pptx import Presentation

logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")
logger = logging.getLogger(__name__)

# Ensure sibling modules are importable
_scripts_dir = str(Path(__file__).parent)
if _scripts_dir not in sys.path:
    sys.path.insert(0, _scripts_dir)

# Language labels for prompts
LANGUAGE_LABELS = {
    "ES": "Spanish",
    "EN": "English",
    "FR": "French",
    "DE": "German",
    "PT": "Portuguese",
    "IT": "Italian",
}


def _call_gemini_translate(
    texts: list[str],
    source_lang: str,
    target_lang: str,
    api_key: str,
    model: str = "gemini-2.5-flash",
) -> list[str]:
    """
    Batch-translate a list of text strings via Gemini API.

    Sends all strings at once for consistent context-aware translation.

    Args:
        texts: List of text strings to translate.
        source_lang: Source language code (e.g., "ES").
        target_lang: Target language code (e.g., "EN").
        api_key: Gemini API key.
        model: Gemini model name.

    Returns:
        List of translated strings in the same order.
    """
    src_label = LANGUAGE_LABELS.get(source_lang, source_lang)
    tgt_label = LANGUAGE_LABELS.get(target_lang, target_lang)

    # Build numbered text list for the prompt
    numbered = "\n".join(f"[{i}] {t}" for i, t in enumerate(texts))

    system_prompt = (
        f"You are a professional translator specializing in business presentations. "
        f"Translate from {src_label} to {tgt_label}. "
        f"Rules:\n"
        f"1. Preserve formatting: em dashes, bullet structure, capitalization style\n"
        f"2. Keep proper nouns, brand names, and acronyms unchanged\n"
        f"3. Maintain the same tone (formal/informal) as the source\n"
        f"4. Keep numbers, percentages, and currency symbols as-is\n"
        f"5. If text is already in {tgt_label}, return it unchanged\n"
        f"6. Micro-copy style: keep translations concise, not verbose\n"
    )

    user_prompt = (
        f"Translate each numbered text from {src_label} to {tgt_label}. "
        f"Return a JSON array with exactly {len(texts)} strings in the same order.\n\n"
        f"{numbered}"
    )

    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"

    payload = {
        "system_instruction": {
            "parts": [{"text": system_prompt}]
        },
        "contents": [{
            "parts": [{"text": user_prompt}]
        }],
        "generationConfig": {
            "temperature": 0.2,
            "maxOutputTokens": 8192,
            "responseMimeType": "application/json",
        },
    }

    logger.info("Translating %d texts (%s → %s) via Gemini...", len(texts), source_lang, target_lang)
    resp = requests.post(url, json=payload, timeout=120)
    resp.raise_for_status()
    data = resp.json()

    if "candidates" not in data or not data["candidates"]:
        raise RuntimeError("Gemini returned no candidates for translation")

    raw_text = data["candidates"][0]["content"]["parts"][0]["text"]

    # Parse JSON array response
    try:
        translated = json.loads(raw_text)
    except json.JSONDecodeError:
        cleaned = raw_text.strip()
        cleaned = re.sub(r"^```(?:json)?\s*\n?", "", cleaned)
        cleaned = re.sub(r"\n?```\s*$", "", cleaned)
        translated = json.loads(cleaned)

    if not isinstance(translated, list):
        raise ValueError(f"Expected JSON array, got {type(translated).__name__}")

    if len(translated) != len(texts):
        logger.warning(
            "Translation count mismatch: sent %d, received %d. Padding/truncating.",
            len(texts), len(translated),
        )
        # Pad with originals or truncate
        if len(translated) < len(texts):
            translated.extend(texts[len(translated):])
        else:
            translated = translated[:len(texts)]

    return translated


def translate_pptx(
    input_pptx: str,
    output_pptx: str | None = None,
    source_lang: str = "ES",
    target_lang: str = "EN",
    api_key: str | None = None,
    model: str = "gemini-2.5-flash",
) -> dict:
    """
    Translate all text in a PPTX file while preserving layout and formatting.

    Extracts text from all shapes, batch-translates via Gemini, then replaces
    text in-place (preserving font, size, color, position).

    Args:
        input_pptx: Path to input PPTX file.
        output_pptx: Output path (auto-generated if None).
        source_lang: Source language code.
        target_lang: Target language code.
        api_key: Gemini API key (falls back to GEMINI_API_KEY env var).
        model: Gemini model name.

    Returns:
        Dict with success status, file paths, and metadata.
    """
    api_key = api_key or os.environ.get("GEMINI_API_KEY")
    if not api_key:
        return {
            "success": False,
            "error": "No API key. Set GEMINI_API_KEY or pass api_key parameter.",
        }

    if not Path(input_pptx).exists():
        return {"success": False, "error": f"File not found: {input_pptx}"}

    prs = Presentation(input_pptx)

    # Step 1: Collect all text strings with their locations
    # Each entry: (slide_idx, shape_idx, para_idx, run_idx, text)
    text_entries: list[tuple[int, int, int, int, str]] = []
    texts_to_translate: list[str] = []

    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(paragraph.runs):
                    text = run.text.strip()
                    if text:
                        text_entries.append((s_idx, sh_idx, p_idx, r_idx, run.text))
                        texts_to_translate.append(run.text)

    if not texts_to_translate:
        return {
            "success": False,
            "error": "No text found in PPTX to translate.",
        }

    logger.info("Found %d text runs across %d slides", len(texts_to_translate), len(prs.slides))

    # Step 2: Batch translate (split into batches of 50 for API limits)
    batch_size = 50
    all_translated: list[str] = []

    for batch_start in range(0, len(texts_to_translate), batch_size):
        batch = texts_to_translate[batch_start:batch_start + batch_size]

        if batch_start > 0:
            time.sleep(1.0)  # Rate limiting

        translated_batch = _call_gemini_translate(
            texts=batch,
            source_lang=source_lang,
            target_lang=target_lang,
            api_key=api_key,
            model=model,
        )
        all_translated.extend(translated_batch)

    # Step 3: Replace text in PPTX while preserving formatting
    slides_list = list(prs.slides)
    for i, (s_idx, sh_idx, p_idx, r_idx, original_text) in enumerate(text_entries):
        slide = slides_list[s_idx]
        shapes_list = list(slide.shapes)
        shape = shapes_list[sh_idx]
        paragraph = shape.text_frame.paragraphs[p_idx]
        run = paragraph.runs[r_idx]
        # Replace text only — font, size, color, bold, italic are preserved
        run.text = all_translated[i]

    # Step 4: Save output
    if not output_pptx:
        stem = Path(input_pptx).stem
        parent = Path(input_pptx).parent
        output_pptx = str(parent / f"{stem}_{target_lang.lower()}.pptx")

    prs.save(output_pptx)

    changed_count = sum(
        1 for orig, trans in zip(texts_to_translate, all_translated) if orig != trans
    )

    logger.info("Translation complete: %d/%d runs changed", changed_count, len(texts_to_translate))
    logger.info("Output: %s", output_pptx)

    return {
        "success": True,
        "files": {"pptx": output_pptx},
        "metadata": {
            "source_lang": source_lang,
            "target_lang": target_lang,
            "total_runs": len(texts_to_translate),
            "changed_runs": changed_count,
            "slide_count": len(prs.slides),
        },
    }


def translate_specs(
    specs: "PresentationSpec",
    target_lang: str,
    api_key: str | None = None,
    model: str = "gemini-2.5-flash",
) -> "PresentationSpec":
    """
    Translate SlideSpec content before building a new PPTX.

    Creates a deep copy and translates all text fields. Useful for generating
    a presentation in one language then producing a translated version.

    Args:
        specs: PresentationSpec to translate.
        target_lang: Target language code.
        api_key: Gemini API key (falls back to GEMINI_API_KEY env var).
        model: Gemini model name.

    Returns:
        New PresentationSpec with translated text.
    """
    api_key = api_key or os.environ.get("GEMINI_API_KEY")
    if not api_key:
        raise ValueError("No API key. Set GEMINI_API_KEY or pass api_key parameter.")

    source_lang = specs.language if specs.language != "auto" else "ES"
    translated_specs = deepcopy(specs)

    # Collect all translatable strings from all slides
    texts: list[str] = []
    # Track: title, subtitle of presentation
    texts.append(specs.title)
    texts.append(specs.subtitle)

    for slide in specs.slides:
        texts.append(slide.title)
        texts.append(slide.subtitle)
        texts.append(slide.body)
        texts.append(slide.speaker_notes)
        texts.append(slide.left_header)
        texts.append(slide.right_header)
        for bp in slide.bullet_points:
            texts.append(bp)
        for lc in slide.left_column:
            texts.append(lc)
        for rc in slide.right_column:
            texts.append(rc)
        for ci in slide.checkbox_items:
            texts.append(ci)

    # Filter empty strings but track positions
    non_empty_indices = [i for i, t in enumerate(texts) if t.strip()]
    non_empty_texts = [texts[i] for i in non_empty_indices]

    if not non_empty_texts:
        logger.warning("No text to translate in specs")
        return translated_specs

    # Batch translate
    batch_size = 50
    all_translated: list[str] = []

    for batch_start in range(0, len(non_empty_texts), batch_size):
        batch = non_empty_texts[batch_start:batch_start + batch_size]
        if batch_start > 0:
            time.sleep(1.0)

        translated_batch = _call_gemini_translate(
            texts=batch,
            source_lang=source_lang,
            target_lang=target_lang,
            api_key=api_key,
            model=model,
        )
        all_translated.extend(translated_batch)

    # Rebuild translated texts array (put translations back at correct indices)
    translated_full = list(texts)  # Start with originals
    for idx, translated_text in zip(non_empty_indices, all_translated):
        translated_full[idx] = translated_text

    # Apply translations back to specs
    pos = 0
    translated_specs.title = translated_full[pos]; pos += 1
    translated_specs.subtitle = translated_full[pos]; pos += 1
    translated_specs.language = target_lang

    for slide in translated_specs.slides:
        slide.title = translated_full[pos]; pos += 1
        slide.subtitle = translated_full[pos]; pos += 1
        slide.body = translated_full[pos]; pos += 1
        slide.speaker_notes = translated_full[pos]; pos += 1
        slide.left_header = translated_full[pos]; pos += 1
        slide.right_header = translated_full[pos]; pos += 1
        new_bp = []
        for _ in slide.bullet_points:
            new_bp.append(translated_full[pos]); pos += 1
        slide.bullet_points = new_bp
        new_lc = []
        for _ in slide.left_column:
            new_lc.append(translated_full[pos]); pos += 1
        slide.left_column = new_lc
        new_rc = []
        for _ in slide.right_column:
            new_rc.append(translated_full[pos]); pos += 1
        slide.right_column = new_rc
        new_ci = []
        for _ in slide.checkbox_items:
            new_ci.append(translated_full[pos]); pos += 1
        slide.checkbox_items = new_ci

    logger.info("Translated specs: %d texts (%s → %s)", len(non_empty_texts), source_lang, target_lang)
    return translated_specs


def main():
    """CLI entry point for PPTX translation."""
    parser = argparse.ArgumentParser(
        description="Translate a PPTX presentation between languages via Gemini API",
    )
    parser.add_argument("input_pptx", help="Path to PPTX file to translate")
    parser.add_argument("--from", dest="source_lang", default="ES",
                        help="Source language code (default: ES)")
    parser.add_argument("--to", dest="target_lang", required=True,
                        help="Target language code (e.g., EN, FR, DE)")
    parser.add_argument("--output", "-o", help="Output path for translated PPTX")
    parser.add_argument("--model", default="gemini-2.5-flash",
                        help="Gemini model (default: gemini-2.5-flash)")

    args = parser.parse_args()

    if not Path(args.input_pptx).exists():
        print(f"Error: File not found: {args.input_pptx}")
        sys.exit(1)

    result = translate_pptx(
        input_pptx=args.input_pptx,
        output_pptx=args.output,
        source_lang=args.source_lang,
        target_lang=args.target_lang,
        model=args.model,
    )

    if not result["success"]:
        print(f"\nError: {result.get('error', 'Unknown error')}")
        sys.exit(1)

    meta = result["metadata"]
    print(f"\nTranslation complete ({meta['source_lang']} → {meta['target_lang']})")
    print(f"  Slides: {meta['slide_count']}")
    print(f"  Text runs: {meta['total_runs']} ({meta['changed_runs']} changed)")
    print(f"  Output: {result['files']['pptx']}")


if __name__ == "__main__":
    main()
