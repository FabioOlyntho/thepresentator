"""
Slide Builder — Creates professional PPTX presentations.

Three modes:
- Editable (default): Programmatic Recodme layouts — solid fills, accent bars,
  proper typography. No AI images needed; instant generation.
- Full-slide: Each slide is a single AI-rendered image (NotebookLM approach)
- Composite: AI illustrations + python-pptx text layout (legacy)

Generates presentations using python-pptx with:
- Programmatic Recodme brand layouts (editable mode)
- AI-generated full-slide images (full-slide mode)
- Editable text boxes with Poppins typography
- Speaker notes
- 16:9 widescreen format
"""

import json
import logging
import sys
from pathlib import Path
from dataclasses import dataclass
from lxml import etree

# Ensure sibling modules (image_generator, etc.) are importable
_scripts_dir = str(Path(__file__).parent)
if _scripts_dir not in sys.path:
    sys.path.insert(0, _scripts_dir)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap

logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")
logger = logging.getLogger(__name__)


@dataclass
class BrandConfig:
    """Brand configuration for presentation styling."""
    primary: RGBColor
    secondary: RGBColor
    background: RGBColor
    accent: RGBColor
    text_dark: RGBColor
    text_light: RGBColor
    highlight: RGBColor
    font_title: str
    font_body: str
    font_accent: str

    @classmethod
    def from_json(cls, filepath: str) -> "BrandConfig":
        data = json.loads(Path(filepath).read_text(encoding="utf-8"))
        colors = data["colors"]
        fonts = data["fonts"]
        return cls(
            primary=RGBColor.from_string(colors["primary"][1:]),
            secondary=RGBColor.from_string(colors["secondary"][1:]),
            background=RGBColor.from_string(colors["background"][1:]),
            accent=RGBColor.from_string(colors["accent"][1:]),
            text_dark=RGBColor.from_string(colors["text_dark"][1:]),
            text_light=RGBColor.from_string(colors["text_light"][1:]),
            highlight=RGBColor.from_string(colors["highlight"][1:]),
            font_title=fonts["title"],
            font_body=fonts["body"],
            font_accent=fonts.get("accent", fonts["body"]),
        )

    @classmethod
    def default(cls) -> "BrandConfig":
        return cls(
            primary=RGBColor(0xE8, 0x44, 0x22),
            secondary=RGBColor(0xC2, 0xBF, 0xAA),
            background=RGBColor(0xF5, 0xF0, 0xE8),
            accent=RGBColor(0x01, 0x26, 0x2D),
            text_dark=RGBColor(0x31, 0x31, 0x31),
            text_light=RGBColor(0xFF, 0xFF, 0xFF),
            highlight=RGBColor(0xE8, 0x44, 0x22),
            font_title="Poppins SemiBold",
            font_body="Poppins Medium",
            font_accent="Poppins Light",
        )


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _rgb_to_hex(color: RGBColor) -> str:
    """Convert RGBColor to hex string (without #)."""
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


class SlideBuilder:
    """Builds PPTX presentations with AI-generated images and modern design.

    Three modes:
    - editable_mode=True (new default): AI background image + editable text boxes.
      Text is selectable/editable in PowerPoint with Recodme brand fonts.
    - full_slide_mode=True: Each slide is a single full-bleed image with only
      speaker notes as editable text.
    - full_slide_mode=False: Composite mode with AI illustrations + python-pptx
      text layout (legacy approach).
    """

    # 16:9 widescreen dimensions
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(
        self,
        brand: BrandConfig | None = None,
        image_paths: dict[int, str | None] | None = None,
        full_slide_mode: bool = True,
        editable_mode: bool = False,
    ):
        self.brand = brand or BrandConfig.default()
        self.image_paths = image_paths or {}
        self.full_slide_mode = full_slide_mode
        self.editable_mode = editable_mode
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

    def _set_slide_bg(self, slide, color: RGBColor):
        """Set slide background to a solid color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _set_gradient_bg(self, slide, color_start: RGBColor, color_end: RGBColor):
        """Set slide background to a two-stop gradient."""
        bg = slide.background
        fill = bg.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = color_start
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = color_end
        fill.gradient_stops[1].position = 1.0

    def _add_full_bleed_image(self, slide, image_path: str):
        """Add an image that fills the entire slide."""
        slide.shapes.add_picture(
            image_path,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
        )

    def _add_image_region(self, slide, image_path: str, left, top, width, height):
        """Add an image at a specific position and size."""
        slide.shapes.add_picture(image_path, left, top, width, height)

    def _add_dark_overlay(self, slide, left, top, width, height, opacity_pct: int = 60):
        """Add a semi-transparent dark rectangle for text readability over images."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        shape.line.fill.background()

        # Set fill with transparency via XML
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x1A)

        # Apply alpha transparency
        self._set_shape_alpha(shape, opacity_pct)
        return shape

    def _add_colored_overlay(self, slide, left, top, width, height, color: RGBColor, opacity_pct: int = 85):
        """Add a semi-transparent colored rectangle."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        shape.line.fill.background()
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = color
        self._set_shape_alpha(shape, opacity_pct)
        return shape

    def _set_shape_alpha(self, shape, opacity_pct: int):
        """Set shape fill transparency (0=fully transparent, 100=fully opaque)."""
        # python-pptx doesn't expose alpha directly, so we manipulate XML
        alpha_val = opacity_pct * 1000  # PowerPoint uses units of 1/1000th percent
        sp_pr = shape._element.spPr
        solid_fill = sp_pr.find(qn("a:solidFill"))
        if solid_fill is not None:
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is not None:
                # Remove existing alpha
                for child in list(srgb):
                    if "alpha" in child.tag.lower():
                        srgb.remove(child)
                alpha_elem = etree.SubElement(srgb, qn("a:alpha"))
                alpha_elem.set("val", str(alpha_val))

    def _add_gradient_overlay(self, slide, left, top, width, height, direction: str = "bottom"):
        """Add a gradient overlay (transparent top → dark bottom, or left/right variants)."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        shape.line.fill.background()

        # Build gradient fill via XML for proper transparency control
        sp_pr = shape._element.spPr

        # Remove any existing fill
        for child in list(sp_pr):
            if child.tag.endswith("Fill") or "fill" in child.tag.lower():
                sp_pr.remove(child)

        grad_fill = etree.SubElement(sp_pr, qn("a:gradFill"))

        # Gradient stops
        gs_lst = etree.SubElement(grad_fill, qn("a:gsLst"))

        # Stop 1: transparent
        gs1 = etree.SubElement(gs_lst, qn("a:gs"))
        gs1.set("pos", "0")
        srgb1 = etree.SubElement(gs1, qn("a:srgbClr"))
        srgb1.set("val", "0A0A1A")
        alpha1 = etree.SubElement(srgb1, qn("a:alpha"))
        alpha1.set("val", "0")

        # Stop 2: semi-dark at midpoint
        gs2 = etree.SubElement(gs_lst, qn("a:gs"))
        gs2.set("pos", "50000")
        srgb2 = etree.SubElement(gs2, qn("a:srgbClr"))
        srgb2.set("val", "0A0A1A")
        alpha2 = etree.SubElement(srgb2, qn("a:alpha"))
        alpha2.set("val", "40000")

        # Stop 3: dark
        gs3 = etree.SubElement(gs_lst, qn("a:gs"))
        gs3.set("pos", "100000")
        srgb3 = etree.SubElement(gs3, qn("a:srgbClr"))
        srgb3.set("val", "0A0A1A")
        alpha3 = etree.SubElement(srgb3, qn("a:alpha"))
        alpha3.set("val", "75000")

        # Linear gradient direction
        lin = etree.SubElement(grad_fill, qn("a:lin"))
        angle_map = {
            "bottom": "5400000",   # top to bottom
            "top": "16200000",     # bottom to top
            "left": "10800000",    # right to left
            "right": "0",          # left to right
        }
        lin.set("ang", angle_map.get(direction, "5400000"))
        lin.set("scaled", "0")

        return shape

    def _add_text_box(
        self, slide, left, top, width, height,
        text: str, font_size: int, font_name: str,
        color: RGBColor, bold: bool = False,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
        line_spacing: float | None = None,
    ):
        """Add a text box with formatted text."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = color
        run.font.bold = bold
        p.alignment = alignment

        if line_spacing:
            p.line_spacing = Pt(line_spacing)

        # Set vertical alignment via XML
        txBody = tf._txBody
        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            anchor_map = {
                MSO_ANCHOR.TOP: "t",
                MSO_ANCHOR.MIDDLE: "ctr",
                MSO_ANCHOR.BOTTOM: "b",
            }
            bodyPr.set("anchor", anchor_map.get(anchor, "t"))

        return txBox

    def _add_bullet_list(
        self, slide, left, top, width, height,
        items: list[str], font_size: int, font_name: str,
        color: RGBColor, line_spacing: float | None = None,
    ):
        """Add a bullet point list."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.color.rgb = color
            p.space_after = Pt(10)
            p.level = 0

            if line_spacing:
                p.line_spacing = Pt(line_spacing)

            # Add colored bullet character via XML
            pPr = p._p.get_or_add_pPr()
            for child in list(pPr):
                if "buNone" in child.tag or "buChar" in child.tag or "buColor" in child.tag.lower():
                    pPr.remove(child)
            # Bullet color (teal accent)
            buClr = pPr.makeelement(qn("a:buClr"), {})
            srgbClr = buClr.makeelement(qn("a:srgbClr"), {"val": _rgb_to_hex(self.brand.accent)})
            buClr.append(srgbClr)
            pPr.append(buClr)
            buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
            pPr.append(buChar)

        return txBox

    def _add_speaker_notes(self, slide, notes_text: str):
        """Add speaker notes to a slide."""
        if not notes_text:
            return
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    def _add_slide_number(self, slide, number: int, color: RGBColor | None = None):
        """Add a small slide number in the bottom-right corner."""
        self._add_text_box(
            slide, Inches(12.3), Inches(7.0),
            Inches(0.8), Inches(0.3),
            str(number), 11, self.brand.font_body,
            color or RGBColor(0xAA, 0xAA, 0xAA),
            alignment=PP_ALIGN.RIGHT,
        )

    # ═══════════════════════════════════════════════════════════════
    # RECODME HELPERS — Reusable accent shapes for programmatic layouts
    # ═══════════════════════════════════════════════════════════════

    def _add_accent_bar(self, slide, left, top, width, height=Inches(0.04), color=None):
        """Add a thin solid-fill accent bar (no border). Default: brand primary red."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color or self.brand.primary
        shape.line.fill.background()
        return shape

    def _add_divider_line(self, slide, left, top, width, height=Inches(0.03), color=None, opacity_pct=40):
        """Add a thin divider rectangle with optional transparency."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color or self.brand.text_dark
        shape.line.fill.background()
        if opacity_pct < 100:
            self._set_shape_alpha(shape, opacity_pct)
        return shape

    def _add_watermark(self, slide, text="Recodme", color=None):
        """Add a small watermark in the bottom-right corner."""
        self._add_text_box(
            slide, Inches(11.5), Inches(7.0),
            Inches(1.5), Inches(0.3),
            text, 10, self.brand.font_accent,
            color or self.brand.secondary,
            alignment=PP_ALIGN.RIGHT,
        )

    # ═══════════════════════════════════════════════════════════════
    # v4.1 HELPERS — Professional visual elements
    # ═══════════════════════════════════════════════════════════════

    def _add_numbered_circle(self, slide, left, top, number, size=None, bg_color=None, text_color=None):
        """Draw a teal circle with centered white number."""
        sz = size or Inches(0.4)
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, sz, sz)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color or self.brand.accent
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(number)
        run.font.size = Pt(13)
        run.font.name = "Poppins SemiBold"
        run.font.color.rgb = text_color or self.brand.text_light
        run.font.bold = True

        # Vertical center
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")

        return shape

    def _add_content_card(self, slide, left, top, width, height, fill_color=None, border_color=None):
        """Rounded rectangle container for visual grouping."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.fill.background()
        else:
            shape.fill.background()
            if border_color:
                shape.line.color.rgb = border_color
                shape.line.width = Pt(1)
            else:
                shape.line.fill.background()
        return shape

    def _add_footer_bar(self, slide, slide_number=None, dark_bg=False):
        """Consistent footer across all slides (except title)."""
        line_color = RGBColor(0xC2, 0xBF, 0xAA) if not dark_bg else RGBColor(0x4A, 0x5A, 0x5E)
        num_color = RGBColor(0x99, 0x99, 0x99) if not dark_bg else RGBColor(0x88, 0x99, 0x99)
        wm_color = self.brand.secondary if not dark_bg else RGBColor(0x4A, 0x5A, 0x5E)

        # Horizontal line
        line = self._add_divider_line(
            slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.01),
            color=line_color, opacity_pct=30,
        )

        # Watermark
        self._add_text_box(
            slide, Inches(11.0), Inches(7.1),
            Inches(1.5), Inches(0.3),
            "Recodme", 10, self.brand.font_accent,
            wm_color, alignment=PP_ALIGN.RIGHT,
        )

        # Slide number
        if slide_number is not None:
            self._add_text_box(
                slide, Inches(12.3), Inches(7.1),
                Inches(0.8), Inches(0.3),
                str(slide_number), 11, self.brand.font_body,
                num_color, alignment=PP_ALIGN.RIGHT,
            )

    def _add_rich_bullet_list(self, slide, left, top, width, height, items, font_size, font_name, color, bold_color=None):
        """Enhanced bullet list with bold keyword + regular description."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        bc = bold_color or self.brand.accent

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Teal bullet character
            bullet_run = p.add_run()
            bullet_run.text = "\u2022 "
            bullet_run.font.size = Pt(font_size)
            bullet_run.font.name = font_name
            bullet_run.font.color.rgb = self.brand.accent

            # Split on em dash delimiter
            if " \u2014 " in item:
                keyword, description = item.split(" \u2014 ", 1)
                # Bold keyword
                kw_run = p.add_run()
                kw_run.text = keyword
                kw_run.font.size = Pt(font_size)
                kw_run.font.name = font_name
                kw_run.font.color.rgb = bc
                kw_run.font.bold = True
                # Separator
                sep_run = p.add_run()
                sep_run.text = " \u2014 "
                sep_run.font.size = Pt(font_size)
                sep_run.font.name = font_name
                sep_run.font.color.rgb = color
                # Description
                desc_run = p.add_run()
                desc_run.text = description
                desc_run.font.size = Pt(font_size)
                desc_run.font.name = font_name
                desc_run.font.color.rgb = color
            else:
                # No delimiter — render as regular text
                text_run = p.add_run()
                text_run.text = item
                text_run.font.size = Pt(font_size)
                text_run.font.name = font_name
                text_run.font.color.rgb = color

            p.space_after = Pt(14)

        return txBox

    def _get_image(self, slide_number: int) -> str | None:
        """Get the image path for a slide, if available."""
        path = self.image_paths.get(slide_number)
        if path and Path(path).exists():
            return path
        return None

    # ═══════════════════════════════════════════════════════════════
    # SLIDE BUILDERS — Image-dominant layouts
    # ═══════════════════════════════════════════════════════════════

    def build_title_slide(self, spec) -> None:
        """
        Title slide: Full-bleed image background with dark gradient overlay.
        Large bold title in white at bottom-left, subtitle below.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        image_path = self._get_image(spec.number)

        if image_path:
            # Full-bleed image background
            self._add_full_bleed_image(slide, image_path)
            # Gradient overlay: transparent top → dark bottom
            self._add_gradient_overlay(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, self.SLIDE_HEIGHT, "bottom")
        else:
            # Fallback: rich gradient background
            self._set_gradient_bg(slide, self.brand.text_dark, RGBColor(0x0A, 0x0A, 0x1A))
            # Accent bar at top
            self._add_colored_overlay(
                slide, Inches(0), Inches(0), self.SLIDE_WIDTH, Inches(0.06),
                self.brand.primary, opacity_pct=100,
            )

        # Title — large, bold, white, bottom-left
        self._add_text_box(
            slide, Inches(1.0), Inches(4.2),
            Inches(10.5), Inches(2.0),
            spec.title, 44, self.brand.font_title,
            self.brand.text_light, bold=True,
            anchor=MSO_ANCHOR.BOTTOM,
            line_spacing=52,
        )

        # Subtitle
        if spec.subtitle:
            self._add_text_box(
                slide, Inches(1.0), Inches(6.3),
                Inches(10.5), Inches(0.8),
                spec.subtitle, 22, self.brand.font_accent,
                RGBColor(0xCC, 0xCC, 0xDD),
            )

        self._add_speaker_notes(slide, spec.speaker_notes)

    def build_section_slide(self, spec) -> None:
        """
        Section divider: Full-bleed image with centered text overlay.
        Dark semi-transparent overlay with large white section title.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        image_path = self._get_image(spec.number)

        if image_path:
            self._add_full_bleed_image(slide, image_path)
            self._add_dark_overlay(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, self.SLIDE_HEIGHT, opacity_pct=55)
        else:
            self._set_slide_bg(slide, self.brand.primary)

        # Section title — large, centered
        self._add_text_box(
            slide, Inches(1.5), Inches(2.0),
            Inches(10.333), Inches(3.0),
            spec.title, 42, self.brand.font_title,
            self.brand.text_light, bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=50,
        )

        # Optional body text below
        if spec.body:
            self._add_text_box(
                slide, Inches(2.0), Inches(5.0),
                Inches(9.333), Inches(1.5),
                spec.body, 18, self.brand.font_body,
                RGBColor(0xDD, 0xDD, 0xEE),
                alignment=PP_ALIGN.CENTER,
            )

        self._add_speaker_notes(slide, spec.speaker_notes)

    def build_content_slide(self, spec) -> None:
        """
        Content slide: Split layout — 55% image side, 45% content side.
        Alternates image left/right based on slide number for visual variety.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        image_path = self._get_image(spec.number)
        is_image_left = (spec.number % 2 == 0)

        image_width = Inches(7.0)
        content_width = Inches(6.333)

        if image_path:
            if is_image_left:
                # Image on left
                self._add_image_region(slide, image_path, Inches(0), Inches(0), image_width, self.SLIDE_HEIGHT)
                # Dark content panel on right
                self._add_colored_overlay(
                    slide, image_width, Inches(0), content_width, self.SLIDE_HEIGHT,
                    self.brand.text_dark, opacity_pct=95,
                )
                content_left = Inches(7.8)
                content_area_width = Inches(4.8)
                title_color = self.brand.text_light
                body_color = RGBColor(0xCC, 0xCC, 0xDD)
                bullet_color = RGBColor(0xDD, 0xDD, 0xEE)
                num_color = RGBColor(0x88, 0x88, 0x99)
            else:
                # Image on right
                self._add_colored_overlay(
                    slide, Inches(0), Inches(0), content_width, self.SLIDE_HEIGHT,
                    self.brand.text_dark, opacity_pct=95,
                )
                self._add_image_region(slide, image_path, content_width, Inches(0), image_width, self.SLIDE_HEIGHT)
                content_left = Inches(0.8)
                content_area_width = Inches(4.8)
                title_color = self.brand.text_light
                body_color = RGBColor(0xCC, 0xCC, 0xDD)
                bullet_color = RGBColor(0xDD, 0xDD, 0xEE)
                num_color = RGBColor(0x88, 0x88, 0x99)
        else:
            # Fallback: gradient bg, no image
            self._set_gradient_bg(slide, self.brand.text_dark, RGBColor(0x12, 0x12, 0x2E))
            content_left = Inches(1.0)
            content_area_width = Inches(11.0)
            title_color = self.brand.text_light
            body_color = RGBColor(0xCC, 0xCC, 0xDD)
            bullet_color = RGBColor(0xDD, 0xDD, 0xEE)
            num_color = RGBColor(0x88, 0x88, 0x99)

        # Slide number
        self._add_text_box(
            slide, content_left, Inches(0.4),
            Inches(1), Inches(0.3),
            str(spec.number), 11, self.brand.font_body,
            num_color,
        )

        # Title (conclusion-style)
        self._add_text_box(
            slide, content_left, Inches(0.8),
            content_area_width, Inches(1.8),
            spec.title, 28, self.brand.font_title,
            title_color, bold=True,
            line_spacing=34,
        )

        # Accent line under title
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, content_left, Inches(2.7),
            Inches(1.5), Inches(0.04),
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = self.brand.primary
        accent_shape.line.fill.background()

        # Body text
        body_top = Inches(3.0)
        if spec.body:
            self._add_text_box(
                slide, content_left, body_top,
                content_area_width, Inches(2.0),
                spec.body, 16, self.brand.font_body,
                body_color,
                line_spacing=24,
            )
            body_top = Inches(5.0)

        # Bullet points
        if spec.bullet_points:
            self._add_bullet_list(
                slide, content_left, body_top,
                content_area_width, Inches(2.2),
                spec.bullet_points, 16, self.brand.font_body,
                bullet_color,
                line_spacing=24,
            )

        # Source reference
        if spec.source_reference:
            self._add_text_box(
                slide, content_left, Inches(7.0),
                content_area_width, Inches(0.3),
                f"Source: {spec.source_reference}", 9, self.brand.font_body,
                RGBColor(0x66, 0x66, 0x77),
            )

        self._add_speaker_notes(slide, spec.speaker_notes)

    def build_comparison_slide(self, spec) -> None:
        """
        Comparison slide: Image background with two-column content overlay.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        image_path = self._get_image(spec.number)

        if image_path:
            self._add_full_bleed_image(slide, image_path)
            self._add_dark_overlay(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, self.SLIDE_HEIGHT, opacity_pct=70)
        else:
            self._set_gradient_bg(slide, self.brand.text_dark, RGBColor(0x12, 0x12, 0x2E))

        title_color = self.brand.text_light
        body_color = RGBColor(0xDD, 0xDD, 0xEE)

        # Title
        self._add_text_box(
            slide, Inches(1.0), Inches(0.5),
            Inches(11.333), Inches(1.2),
            spec.title, 30, self.brand.font_title,
            title_color, bold=True,
        )

        # Accent line
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8),
            Inches(2.0), Inches(0.04),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.brand.primary
        accent.line.fill.background()

        # Vertical divider
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(2.2),
            Inches(0.04), Inches(4.5),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.brand.primary
        divider.line.fill.background()
        self._set_shape_alpha(divider, 50)

        # Split bullets into left and right columns
        if len(spec.bullet_points) >= 2:
            half = len(spec.bullet_points) // 2
            left_items = spec.bullet_points[:half]
            right_items = spec.bullet_points[half:]
        else:
            left_items = spec.bullet_points
            right_items = []

        self._add_bullet_list(
            slide, Inches(1.0), Inches(2.5),
            Inches(5.0), Inches(4.5),
            left_items, 16, self.brand.font_body, body_color,
        )

        if right_items:
            self._add_bullet_list(
                slide, Inches(7.2), Inches(2.5),
                Inches(5.0), Inches(4.5),
                right_items, 16, self.brand.font_body, body_color,
            )

        self._add_slide_number(slide, spec.number, RGBColor(0x88, 0x88, 0x99))
        self._add_speaker_notes(slide, spec.speaker_notes)

    def build_data_slide(self, spec) -> None:
        """
        Data/stats slide: Image background with prominent stat card overlay.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        image_path = self._get_image(spec.number)

        if image_path:
            self._add_full_bleed_image(slide, image_path)
            self._add_dark_overlay(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, self.SLIDE_HEIGHT, opacity_pct=65)
        else:
            self._set_gradient_bg(slide, self.brand.text_dark, RGBColor(0x12, 0x12, 0x2E))

        # Title
        self._add_text_box(
            slide, Inches(1.0), Inches(0.5),
            Inches(11.333), Inches(1.2),
            spec.title, 30, self.brand.font_title,
            self.brand.text_light, bold=True,
        )

        # Key stat card
        if spec.body:
            stat_card = self._add_colored_overlay(
                slide, Inches(1.5), Inches(2.2),
                Inches(10.333), Inches(2.2),
                self.brand.primary, opacity_pct=85,
            )
            self._add_text_box(
                slide, Inches(2.0), Inches(2.5),
                Inches(9.333), Inches(1.8),
                spec.body, 24, self.brand.font_title,
                self.brand.text_light, bold=True,
                alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )

        # Supporting data points
        if spec.bullet_points:
            self._add_bullet_list(
                slide, Inches(1.5), Inches(4.8),
                Inches(10.333), Inches(2.2),
                spec.bullet_points, 18, self.brand.font_body,
                RGBColor(0xDD, 0xDD, 0xEE),
            )

        self._add_slide_number(slide, spec.number, RGBColor(0x88, 0x88, 0x99))
        self._add_speaker_notes(slide, spec.speaker_notes)

    def build_quote_slide(self, spec) -> None:
        """
        Quote slide: Full-bleed image with dark overlay and prominent quote.
        Large quotation marks, quote text in white, attribution in accent color.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        image_path = self._get_image(spec.number)

        if image_path:
            self._add_full_bleed_image(slide, image_path)
            self._add_dark_overlay(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, self.SLIDE_HEIGHT, opacity_pct=55)
        else:
            self._set_gradient_bg(slide, RGBColor(0x1A, 0x1A, 0x2E), RGBColor(0x0A, 0x0A, 0x1A))

        # Large quotation mark
        self._add_text_box(
            slide, Inches(1.0), Inches(1.0),
            Inches(2.0), Inches(2.5),
            "\u201C", 120, self.brand.font_title,
            self.brand.primary, bold=True,
        )

        # Quote text — large, centered
        self._add_text_box(
            slide, Inches(2.0), Inches(2.5),
            Inches(9.333), Inches(3.0),
            spec.body or spec.title, 26, self.brand.font_accent,
            self.brand.text_light,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=36,
        )

        # Attribution
        if spec.source_reference:
            self._add_text_box(
                slide, Inches(2.0), Inches(5.8),
                Inches(9.333), Inches(0.6),
                f"\u2014 {spec.source_reference}", 16, self.brand.font_body,
                self.brand.highlight,
                alignment=PP_ALIGN.CENTER,
            )

        self._add_speaker_notes(slide, spec.speaker_notes)

    def build_conclusion_slide(self, spec) -> None:
        """
        Conclusion slide: Dark gradient background with image accent strip.
        Key takeaways in white text with good spacing.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        image_path = self._get_image(spec.number)

        if image_path:
            # Image as a strip on the right side
            self._add_colored_overlay(
                slide, Inches(0), Inches(0), Inches(8.5), self.SLIDE_HEIGHT,
                self.brand.text_dark, opacity_pct=100,
            )
            self._add_image_region(
                slide, image_path,
                Inches(8.5), Inches(0), Inches(4.833), self.SLIDE_HEIGHT,
            )
            # Soft gradient fade at image edge
            self._add_gradient_overlay(
                slide, Inches(7.5), Inches(0), Inches(2.0), self.SLIDE_HEIGHT,
                direction="right",
            )
        else:
            self._set_gradient_bg(slide, self.brand.text_dark, RGBColor(0x0A, 0x0A, 0x1A))

        # Title
        self._add_text_box(
            slide, Inches(1.0), Inches(0.8),
            Inches(7.0), Inches(1.5),
            spec.title, 32, self.brand.font_title,
            self.brand.text_light, bold=True,
            line_spacing=40,
        )

        # Accent line
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.4),
            Inches(2.0), Inches(0.05),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.brand.primary
        accent.line.fill.background()

        # Body text
        if spec.body:
            self._add_text_box(
                slide, Inches(1.0), Inches(2.8),
                Inches(7.0), Inches(1.5),
                spec.body, 18, self.brand.font_body,
                RGBColor(0xBB, 0xBB, 0xCC),
                line_spacing=26,
            )

        # Key takeaways
        if spec.bullet_points:
            self._add_bullet_list(
                slide, Inches(1.5), Inches(4.3),
                Inches(6.5), Inches(2.8),
                spec.bullet_points, 18, self.brand.font_body,
                self.brand.text_light,
                line_spacing=28,
            )

        self._add_speaker_notes(slide, spec.speaker_notes)

    # ═══════════════════════════════════════════════════════════════
    # PDNOB MODE — Cleaned image background + editable text boxes
    # ═══════════════════════════════════════════════════════════════

    def build_pdnob_slide(self, cleaned_image_path: str, text_blocks: list) -> None:
        """
        Build a PDNob-style slide: cleaned background image + text boxes at OCR positions.

        Args:
            cleaned_image_path: Path to the text-erased image.
            text_blocks: List of OCRTextBlock objects with position/style info.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Full-bleed cleaned image as background
        self._add_full_bleed_image(slide, cleaned_image_path)

        # Add editable text boxes at original OCR positions
        for block in text_blocks:
            left = Inches(block.x_pct / 100 * 13.333)
            top = Inches(block.y_pct / 100 * 7.5)
            width = Inches(block.width_pct / 100 * 13.333)
            height = Inches(block.height_pct / 100 * 7.5)

            # Ensure minimum size
            if width < Inches(0.3):
                width = Inches(0.3)
            if height < Inches(0.2):
                height = Inches(0.2)

            # Add small padding to width for text wrapping
            width = Inches(min(13.333, block.width_pct / 100 * 13.333 + 0.15))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None

            # Make text box background transparent
            txBox.fill.background()

            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = block.text
            run.font.size = Pt(max(8, min(60, block.font_size_pt)))
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(block.color[0], block.color[1], block.color[2])

            # Remove paragraph spacing for tight fit
            p.space_before = Pt(0)
            p.space_after = Pt(0)

    # ═══════════════════════════════════════════════════════════════
    # EDITABLE MODE — Programmatic Recodme layouts (no AI images)
    # ═══════════════════════════════════════════════════════════════

    def _add_checkbox_list(
        self, slide, left, top, width, height,
        items: list[str], font_size: int, font_name: str,
        color: RGBColor, line_spacing: float | None = None,
    ):
        """Add a checkbox-prefixed list (manifesto/rules style)."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Checkbox character
            check_run = p.add_run()
            check_run.text = "\u2611 "  # Ballot box with check
            check_run.font.size = Pt(font_size + 2)
            check_run.font.name = font_name
            check_run.font.color.rgb = self.brand.primary

            # Item text
            text_run = p.add_run()
            text_run.text = item
            text_run.font.size = Pt(font_size)
            text_run.font.name = font_name
            text_run.font.color.rgb = color
            p.space_after = Pt(14)

            if line_spacing:
                p.line_spacing = Pt(line_spacing)

        return txBox

    # ── Seven Recodme Layout Builders ─────────────────────────────

    def _build_recodme_title(self, spec) -> None:
        """Title slide: Cream bg + teal bottom panel + structured layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.brand.background)

        # Red accent bar at very top
        self._add_accent_bar(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, Inches(0.06))

        # Title — large, left-aligned
        self._add_text_box(
            slide, Inches(1.2), Inches(1.8),
            Inches(10.5), Inches(2.5),
            spec.title, 36, self.brand.font_title,
            self.brand.text_dark, bold=True,
            anchor=MSO_ANCHOR.BOTTOM,
        )

        # Red separator line below title
        self._add_accent_bar(slide, Inches(1.2), Inches(4.5), Inches(2.0), Inches(0.04))

        # Teal bottom panel
        teal_panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(5.5), self.SLIDE_WIDTH, Inches(2.0),
        )
        teal_panel.fill.solid()
        teal_panel.fill.fore_color.rgb = self.brand.accent
        teal_panel.line.fill.background()

        # Small red dash on panel edge
        self._add_accent_bar(slide, Inches(1.2), Inches(5.55), Inches(0.4), Inches(0.03))

        # Subtitle on teal panel — white
        if spec.subtitle:
            self._add_text_box(
                slide, Inches(1.2), Inches(5.8),
                Inches(10.5), Inches(0.8),
                spec.subtitle, 18, self.brand.font_body,
                self.brand.text_light,
            )

        # Watermark on teal panel
        self._add_text_box(
            slide, Inches(11.0), Inches(7.0),
            Inches(1.5), Inches(0.3),
            "Recodme", 10, self.brand.font_accent,
            RGBColor(0x4A, 0x5A, 0x5E),
            alignment=PP_ALIGN.RIGHT,
        )

        self._add_speaker_notes(slide, spec.speaker_notes)

    def _build_recodme_content(self, spec) -> None:
        """Content slide: Numbered items with rich formatting."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.brand.background)

        # Teal top bar
        self._add_accent_bar(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, Inches(0.04), color=self.brand.accent)

        # Title
        self._add_text_box(
            slide, Inches(1.0), Inches(0.4),
            Inches(11.0), Inches(1.2),
            spec.title, 26, self.brand.font_title,
            self.brand.text_dark, bold=True,
        )

        # Red accent under title
        self._add_accent_bar(slide, Inches(1.0), Inches(1.6), Inches(1.5), Inches(0.04))

        # Body text — lighter gray
        bullet_top = Inches(2.0)
        if spec.body:
            self._add_text_box(
                slide, Inches(1.0), Inches(1.9),
                Inches(10.5), Inches(1.0),
                spec.body, 14, self.brand.font_body,
                RGBColor(0x66, 0x66, 0x66),
            )
            bullet_top = Inches(3.0)

        # Numbered bullet items
        if spec.bullet_points:
            items = spec.bullet_points
            if len(items) > 5:
                # Two-column layout for many items
                half = (len(items) + 1) // 2
                left_items = items[:half]
                right_items = items[half:]

                # Vertical divider
                self._add_divider_line(
                    slide, Inches(6.6), bullet_top, Inches(0.02), Inches(3.8),
                    color=self.brand.secondary, opacity_pct=30,
                )

                self._build_numbered_items(slide, Inches(1.0), bullet_top, Inches(5.3), left_items, 0)
                self._build_numbered_items(slide, Inches(7.0), bullet_top, Inches(5.3), right_items, half)
            else:
                self._build_numbered_items(slide, Inches(1.0), bullet_top, Inches(11.0), items, 0)

        self._add_footer_bar(slide, slide_number=spec.number)
        self._add_speaker_notes(slide, spec.speaker_notes)

    def _build_numbered_items(self, slide, left, top, width, items, start_index):
        """Build numbered items with circles, rich text, and separator lines."""
        y = top
        item_height = Inches(0.5)
        spacing = Inches(0.65)

        for i, item in enumerate(items):
            num = start_index + i + 1

            # Numbered circle
            self._add_numbered_circle(slide, left, y + Inches(0.05), num)

            # Rich text next to circle
            text_left = left + Inches(0.6)
            text_width = width - Inches(0.6)

            txBox = slide.shapes.add_textbox(text_left, y, text_width, item_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]

            # Split on em dash delimiter
            if " \u2014 " in item:
                keyword, description = item.split(" \u2014 ", 1)
                kw_run = p.add_run()
                kw_run.text = keyword
                kw_run.font.size = Pt(14)
                kw_run.font.name = self.brand.font_body
                kw_run.font.color.rgb = self.brand.accent
                kw_run.font.bold = True
                sep_run = p.add_run()
                sep_run.text = " \u2014 "
                sep_run.font.size = Pt(14)
                sep_run.font.name = self.brand.font_body
                sep_run.font.color.rgb = self.brand.text_dark
                desc_run = p.add_run()
                desc_run.text = description
                desc_run.font.size = Pt(14)
                desc_run.font.name = self.brand.font_body
                desc_run.font.color.rgb = self.brand.text_dark
            else:
                run = p.add_run()
                run.text = item
                run.font.size = Pt(14)
                run.font.name = self.brand.font_body
                run.font.color.rgb = self.brand.text_dark

            # Separator line (except last item)
            if i < len(items) - 1:
                sep_y = y + item_height + Inches(0.05)
                self._add_divider_line(
                    slide, text_left, sep_y, text_width, Inches(0.01),
                    color=RGBColor(0xE0, 0xDD, 0xD5), opacity_pct=100,
                )

            y += spacing

    def _build_recodme_section(self, spec) -> None:
        """Section divider: Dark teal bg + corner accent blocks."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.brand.accent)  # Dark teal #01262D

        # Top-left red corner block
        corner_tl = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.25), Inches(1.8),
        )
        corner_tl.fill.solid()
        corner_tl.fill.fore_color.rgb = self.brand.primary
        corner_tl.line.fill.background()

        # Bottom-right red corner block
        corner_br = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(13.083), Inches(5.7), Inches(0.25), Inches(1.8),
        )
        corner_br.fill.solid()
        corner_br.fill.fore_color.rgb = self.brand.primary
        corner_br.line.fill.background()

        # Title — large, centered, white
        self._add_text_box(
            slide, Inches(1.5), Inches(2.0),
            Inches(10.333), Inches(2.5),
            spec.title, 38, self.brand.font_title,
            self.brand.text_light, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Centered red accent bar
        self._add_accent_bar(
            slide, Inches(5.5), Inches(4.5), Inches(2.333), Inches(0.05),
        )

        # Body text below bar
        if spec.body:
            self._add_text_box(
                slide, Inches(2.0), Inches(5.0),
                Inches(9.333), Inches(1.0),
                spec.body, 16, self.brand.font_body,
                self.brand.secondary,
                alignment=PP_ALIGN.CENTER,
            )

        # Slide number + watermark
        self._add_text_box(
            slide, Inches(12.3), Inches(0.3),
            Inches(0.8), Inches(0.3),
            str(spec.number), 11, self.brand.font_accent,
            RGBColor(0x4A, 0x5A, 0x5E),
            alignment=PP_ALIGN.RIGHT,
        )
        self._add_text_box(
            slide, Inches(11.0), Inches(7.0),
            Inches(1.5), Inches(0.3),
            "Recodme", 10, self.brand.font_accent,
            RGBColor(0x4A, 0x5A, 0x5E),
            alignment=PP_ALIGN.RIGHT,
        )

        self._add_speaker_notes(slide, spec.speaker_notes)

    def _build_recodme_comparison(self, spec) -> None:
        """Comparison slide: Column headers in teal bars + rich bullet lists."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.brand.background)

        # Teal top bar
        self._add_accent_bar(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, Inches(0.04), color=self.brand.accent)

        # Title
        self._add_text_box(
            slide, Inches(0.8), Inches(0.4),
            Inches(11.5), Inches(1.2),
            spec.title, 26, self.brand.font_title,
            self.brand.text_dark, bold=True,
        )

        # Red accent bar
        self._add_accent_bar(slide, Inches(0.8), Inches(1.6), Inches(2.0), Inches(0.04))

        # Determine column content
        left_items = spec.left_column if spec.left_column else []
        right_items = spec.right_column if spec.right_column else []
        left_hdr = spec.left_header or ""
        right_hdr = spec.right_header or ""

        # Fallback: split bullet_points in half if columns empty
        if not left_items and not right_items and spec.bullet_points:
            half = len(spec.bullet_points) // 2
            left_items = spec.bullet_points[:half] or spec.bullet_points
            right_items = spec.bullet_points[half:]

        # Left header bar — rounded rect with teal fill
        if left_hdr:
            left_bar = self._add_content_card(
                slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.45),
                fill_color=self.brand.accent,
            )
            self._add_text_box(
                slide, Inches(0.8), Inches(2.0),
                Inches(5.5), Inches(0.45),
                left_hdr, 15, self.brand.font_title,
                self.brand.text_light, bold=True,
                alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )

        # Right header bar — rounded rect with teal fill
        if right_hdr:
            right_bar = self._add_content_card(
                slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.45),
                fill_color=self.brand.accent,
            )
            self._add_text_box(
                slide, Inches(7.0), Inches(2.0),
                Inches(5.5), Inches(0.45),
                right_hdr, 15, self.brand.font_title,
                self.brand.text_light, bold=True,
                alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )

        # Vertical divider
        self._add_divider_line(
            slide, Inches(6.5), Inches(2.0), Inches(0.03), Inches(4.5),
            color=self.brand.secondary, opacity_pct=30,
        )

        # Left column — rich bullet list
        if left_items:
            self._add_rich_bullet_list(
                slide, Inches(0.8), Inches(2.7),
                Inches(5.5), Inches(3.8),
                left_items, 14, self.brand.font_body,
                self.brand.text_dark,
            )

        # Right column — rich bullet list
        if right_items:
            self._add_rich_bullet_list(
                slide, Inches(7.0), Inches(2.7),
                Inches(5.5), Inches(3.8),
                right_items, 14, self.brand.font_body,
                self.brand.text_dark,
            )

        self._add_footer_bar(slide, slide_number=spec.number)
        self._add_speaker_notes(slide, spec.speaker_notes)

    def _build_recodme_data(self, spec) -> None:
        """Data/stats slide: Better stat card + numbered supporting points."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.brand.background)

        # Teal top bar
        self._add_accent_bar(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, Inches(0.04), color=self.brand.accent)

        # Title
        self._add_text_box(
            slide, Inches(1.0), Inches(0.4),
            Inches(11.0), Inches(1.0),
            spec.title, 26, self.brand.font_title,
            self.brand.text_dark, bold=True,
        )

        # Red accent bar
        self._add_accent_bar(slide, Inches(1.0), Inches(1.6), Inches(1.5), Inches(0.04))

        # Stat card — rounded rect with light red fill and red border
        if spec.body:
            card = self._add_content_card(
                slide, Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.8),
                fill_color=RGBColor(0xFD, 0xEC, 0xE9),
            )
            card.line.color.rgb = self.brand.primary
            card.line.width = Pt(1.5)

            # Stat text centered in card
            self._add_text_box(
                slide, Inches(1.5), Inches(2.1),
                Inches(10.333), Inches(1.6),
                spec.body, 24, self.brand.font_title,
                self.brand.primary, bold=True,
                alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )

        # Supporting numbered points
        if spec.bullet_points:
            bullet_top = Inches(4.2)
            self._build_numbered_items(slide, Inches(1.0), bullet_top, Inches(11.0), spec.bullet_points, 0)

        self._add_footer_bar(slide, slide_number=spec.number)
        self._add_speaker_notes(slide, spec.speaker_notes)

    def _build_recodme_quote(self, spec) -> None:
        """Quote/list slide: Numbered items mode or quote mode depending on content."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.brand.background)

        has_items = (spec.checkbox_items and len(spec.checkbox_items) > 0) or spec.bullet_points

        if has_items:
            # ── Numbered items mode (replaces checkbox mode) ──
            # Teal top bar
            self._add_accent_bar(slide, Inches(0), Inches(0), self.SLIDE_WIDTH, Inches(0.04), color=self.brand.accent)

            self._add_text_box(
                slide, Inches(1.0), Inches(0.4),
                Inches(11.0), Inches(1.0),
                spec.title, 26, self.brand.font_title,
                self.brand.text_dark, bold=True,
            )
            self._add_accent_bar(slide, Inches(1.0), Inches(1.5), Inches(1.5), Inches(0.04))

            items = spec.checkbox_items if spec.checkbox_items else spec.bullet_points
            self._build_numbered_items(slide, Inches(1.0), Inches(2.0), Inches(11.0), items, 0)
        else:
            # ── Quote mode — with left teal strip ──
            # Left teal strip
            strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(0.35), self.SLIDE_HEIGHT,
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = self.brand.accent
            strip.line.fill.background()

            # Large quote mark
            self._add_text_box(
                slide, Inches(1.2), Inches(1.0),
                Inches(2.0), Inches(2.0),
                "\u201C", 96, self.brand.font_title,
                self.brand.primary, bold=True,
            )

            # Quote text
            self._add_text_box(
                slide, Inches(1.8), Inches(2.8),
                Inches(10.0), Inches(3.0),
                spec.body or spec.title, 22, self.brand.font_accent,
                self.brand.text_dark,
                alignment=PP_ALIGN.LEFT,
                line_spacing=32,
            )

            # Horizontal rule
            self._add_divider_line(
                slide, Inches(1.8), Inches(5.6), Inches(3.0), Inches(0.02),
                color=self.brand.secondary, opacity_pct=100,
            )

            # Attribution
            if spec.source_reference:
                self._add_text_box(
                    slide, Inches(1.8), Inches(5.8),
                    Inches(10.0), Inches(0.5),
                    f"\u2014 {spec.source_reference}", 14, self.brand.font_body,
                    self.brand.secondary,
                )

        self._add_footer_bar(slide, slide_number=spec.number)
        self._add_speaker_notes(slide, spec.speaker_notes)

    def _build_recodme_conclusion(self, spec) -> None:
        """Conclusion slide: Takeaway cards for ≤3 bullets, split panel for >3."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.brand.background)

        bullet_count = len(spec.bullet_points) if spec.bullet_points else 0

        if bullet_count <= 3 and bullet_count > 0:
            # ── Card layout: teal header bar + 3 takeaway cards ──

            # Teal header bar
            header_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), self.SLIDE_WIDTH, Inches(1.6),
            )
            header_bar.fill.solid()
            header_bar.fill.fore_color.rgb = self.brand.accent
            header_bar.line.fill.background()

            # Title on teal header
            self._add_text_box(
                slide, Inches(1.0), Inches(0.2),
                Inches(11.0), Inches(0.8),
                spec.title, 22, self.brand.font_title,
                self.brand.text_light, bold=True,
            )

            # Red accent on teal header
            self._add_accent_bar(slide, Inches(1.0), Inches(1.1), Inches(1.5), Inches(0.04))

            # Body text below header
            body_top = Inches(1.9)
            if spec.body:
                self._add_text_box(
                    slide, Inches(1.0), body_top,
                    Inches(11.0), Inches(0.8),
                    spec.body, 16, self.brand.font_body,
                    self.brand.accent,
                )
                body_top = Inches(2.8)

            # Takeaway cards
            card_width = Inches(3.7)
            card_height = Inches(3.3)
            card_top = Inches(3.2)
            card_positions = [Inches(0.8), Inches(4.8), Inches(8.8)]

            for i, bullet in enumerate(spec.bullet_points[:3]):
                card_x = card_positions[i]

                # Card border (rounded rect, no fill, teal border)
                self._add_content_card(
                    slide, card_x, card_top, card_width, card_height,
                    border_color=self.brand.accent,
                )

                # Numbered circle centered at top of card
                circle_x = card_x + (card_width - Inches(0.4)) // 2
                self._add_numbered_circle(slide, circle_x, card_top + Inches(0.3), i + 1)

                # Split bullet on em dash for keyword vs explanation
                if " \u2014 " in bullet:
                    keyword, explanation = bullet.split(" \u2014 ", 1)
                else:
                    keyword = bullet
                    explanation = ""

                # Bold keyword below circle
                self._add_text_box(
                    slide, card_x + Inches(0.2), card_top + Inches(0.9),
                    card_width - Inches(0.4), Inches(0.6),
                    keyword, 15, self.brand.font_title,
                    self.brand.accent, bold=True,
                    alignment=PP_ALIGN.CENTER,
                )

                # Explanation text below keyword
                if explanation:
                    self._add_text_box(
                        slide, card_x + Inches(0.2), card_top + Inches(1.6),
                        card_width - Inches(0.4), Inches(1.2),
                        explanation, 12, self.brand.font_body,
                        self.brand.text_dark,
                        alignment=PP_ALIGN.CENTER,
                    )
        else:
            # ── Split panel layout for >3 bullets or no bullets ──
            # Left teal panel
            left_panel = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(5.333), self.SLIDE_HEIGHT,
            )
            left_panel.fill.solid()
            left_panel.fill.fore_color.rgb = self.brand.accent
            left_panel.line.fill.background()

            # Red accent bar on teal panel
            self._add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(1.5), Inches(0.05))

            # Title on teal panel — white
            self._add_text_box(
                slide, Inches(0.8), Inches(0.8),
                Inches(4.0), Inches(1.5),
                spec.title, 26, self.brand.font_title,
                self.brand.text_light, bold=True,
            )

            # Takeaway bullets on teal panel — rich formatting
            if spec.bullet_points:
                self._add_rich_bullet_list(
                    slide, Inches(1.0), Inches(2.8),
                    Inches(3.8), Inches(4.0),
                    spec.bullet_points, 14, self.brand.font_body,
                    self.brand.text_light,
                    bold_color=self.brand.primary,
                )

            # Right side — cream background
            # Right accent bar
            self._add_accent_bar(slide, Inches(6.0), Inches(1.2), Inches(2.0))

            # Right body text
            if spec.body:
                self._add_text_box(
                    slide, Inches(6.0), Inches(1.5),
                    Inches(6.5), Inches(2.0),
                    spec.body, 16, self.brand.font_body,
                    self.brand.accent,
                )

        self._add_footer_bar(slide, slide_number=spec.number)
        self._add_speaker_notes(slide, spec.speaker_notes)

    # ── Editable Mode Dispatcher ──────────────────────────────────

    def _build_editable_slide(self, slide_spec) -> None:
        """
        Build an editable slide using programmatic Recodme layouts.

        No AI images needed — solid fills, accent bars, and proper typography
        give 100% control over placement, contrast, and brand fidelity.
        """
        recodme_builders = {
            "title": self._build_recodme_title,
            "section": self._build_recodme_section,
            "content": self._build_recodme_content,
            "comparison": self._build_recodme_comparison,
            "data": self._build_recodme_data,
            "quote": self._build_recodme_quote,
            "conclusion": self._build_recodme_conclusion,
        }
        slide_type = slide_spec.type or "content"
        builder = recodme_builders.get(slide_type, self._build_recodme_content)
        builder(slide_spec)

    def _build_full_slide(self, slide_spec) -> None:
        """
        Build a slide as a single full-bleed image + speaker notes only.

        This is the primary mode (matching NotebookLM approach): the entire slide
        is a pre-rendered image, and speaker notes are the only editable text.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        image_path = self._get_image(slide_spec.number)

        if image_path:
            self._add_full_bleed_image(slide, image_path)
        else:
            # No image available — this shouldn't happen in full-slide mode
            # but if it does, fall back to a dark gradient placeholder
            self._set_gradient_bg(slide, self.brand.text_dark, RGBColor(0x0A, 0x0A, 0x1A))
            # Add title as a centered text box so the slide isn't blank
            self._add_text_box(
                slide, Inches(1.5), Inches(2.5),
                Inches(10.333), Inches(2.5),
                slide_spec.title, 36, self.brand.font_title,
                self.brand.text_light, bold=True,
                alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )

        self._add_speaker_notes(slide, slide_spec.speaker_notes)

    def build_presentation(self, spec) -> "Presentation":
        """
        Build complete PPTX from PresentationSpec.

        In full_slide_mode (default), each slide with an available image becomes
        a single full-bleed image. Slides without images fall back to composite mode.

        Args:
            spec: PresentationSpec with slides

        Returns:
            The Presentation object
        """
        composite_builders = {
            "title": self.build_title_slide,
            "section": self.build_section_slide,
            "content": self.build_content_slide,
            "comparison": self.build_comparison_slide,
            "data": self.build_data_slide,
            "quote": self.build_quote_slide,
            "conclusion": self.build_conclusion_slide,
        }

        for slide_spec in spec.slides:
            slide_type = slide_spec.type
            has_image = bool(self._get_image(slide_spec.number))

            try:
                if self.editable_mode:
                    # Editable path: AI background + text boxes
                    self._build_editable_slide(slide_spec)
                    logger.info(
                        "Built slide %d [%s] EDIT: %s",
                        slide_spec.number, slide_type, slide_spec.title[:60],
                    )
                elif self.full_slide_mode and has_image:
                    # Full-slide path: single AI image
                    self._build_full_slide(slide_spec)
                    logger.info(
                        "Built slide %d [%s] FULL: %s",
                        slide_spec.number, slide_type, slide_spec.title[:60],
                    )
                else:
                    # Composite path (fallback or explicit composite mode)
                    builder = composite_builders.get(slide_type, self.build_content_slide)
                    builder(slide_spec)
                    logger.info(
                        "Built slide %d [%s] %s: %s",
                        slide_spec.number, slide_type,
                        "IMG" if has_image else "GRAD",
                        slide_spec.title[:60],
                    )
            except Exception as e:
                logger.error("Failed to build slide %d: %s", slide_spec.number, e)
                try:
                    self.build_content_slide(slide_spec)
                except Exception:
                    logger.error("Fallback also failed for slide %d", slide_spec.number)

        return self.prs

    def save(self, output_path: str) -> str:
        """Save the presentation to a PPTX file."""
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        logger.info("Presentation saved: %s", path)
        return str(path)


def main():
    """CLI test: build a sample presentation."""
    sys.path.insert(0, str(Path(__file__).parent))
    from gemini_client import SlideSpec, PresentationSpec

    print("\n=== Slide Builder Test ===\n")

    spec = PresentationSpec(
        title="Test Presentation",
        subtitle="Built with Slide Builder",
        language="EN",
        source_document="test.pdf",
        themes=["test", "demo"],
        slides=[
            SlideSpec(1, "title", "Influence Tactics Shape Leadership Outcomes",
                      subtitle="A research-based framework for managers",
                      visual_concept="Network diagram of influence flows"),
            SlideSpec(2, "content", "Seven categories capture all influence tactics",
                      body="Research identified over 1,400 distinct tactics across hundreds of managers.",
                      bullet_points=["Reason", "Coalition", "Friendliness", "Bargaining"],
                      visual_concept="Pie chart of 7 categories"),
            SlideSpec(3, "data", "93% of managers rely on reason as their primary tactic",
                      body="1,400+ tactics -> 7 categories -> 1 dominant approach",
                      bullet_points=["Most versatile", "Works in all directions", "Lowest resistance"]),
            SlideSpec(4, "quote", "The most effective leaders adapt their influence style to the situation",
                      body="The most effective leaders adapt their influence style to the situation, the target, and the objective.",
                      source_reference="Jick, 1987"),
            SlideSpec(5, "conclusion", "Master multiple tactics to maximize leadership impact",
                      body="No single tactic works universally. The key is diagnosis before action.",
                      bullet_points=["Diagnose before acting", "Match tactic to direction", "Build relationships first"],
                      visual_concept="Compass with 7 directions"),
        ],
    )

    brand_path = Path(__file__).parent.parent / "config" / "brand.json"
    brand = BrandConfig.from_json(str(brand_path))

    builder = SlideBuilder(brand)
    builder.build_presentation(spec)

    output = Path(__file__).parent.parent / "output" / "test_presentation.pptx"
    saved = builder.save(str(output))
    print(f"\nPresentation saved: {saved}")
    print(f"Slides: {len(spec.slides)}")
    return saved


if __name__ == "__main__":
    main()
